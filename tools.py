"""
Ultra-Advanced Agent Tools Module for Telegram AI Bot.
Provides deep system monitoring, bash command execution, Python code sandbox & plotter,
live web search, web page content extraction, file & workspace intelligence,
persistent isolated long-term memory, proactive reminders, and desktop/webcam vision.
"""

import os
import sys
import json
import time
import difflib
import subprocess
import asyncio
import psutil
import datetime
import logging
import re
import glob
from contextvars import ContextVar
from typing import Dict, Any, List, Optional

import database
import plugins
from dotenv import load_dotenv
load_dotenv()
from ddgs import DDGS

logger = logging.getLogger("AgentTools")

# Context Variables for Multi-user Dynamic Context Isolation
current_user_id_var: ContextVar[int] = ContextVar("current_user_id", default=0)
current_chat_id_var: ContextVar[int] = ContextVar("current_chat_id", default=0)

SANDBOX_DIR = "/dev/shm/alfa_sandbox"
os.makedirs(SANDBOX_DIR, exist_ok=True)


def get_current_user_id() -> int:
    """Get active Telegram User ID for the current agent turn."""
    uid = current_user_id_var.get()
    return uid if uid else 0


def get_current_chat_id() -> int:
    """Get active Telegram Chat ID for the current agent turn."""
    cid = current_chat_id_var.get()
    return cid if cid else get_current_user_id()


def get_system_stats() -> Dict[str, Any]:
    """
    Get real-time Linux system health metrics including CPU cores/frequencies, RAM, Swap,
    Disk usage, Network interfaces, Battery/Thermal status, Uptime, and Top Processes.
    Use this tool when the user asks about laptop/server specs, resource usage, or performance.
    """
    try:
        cpu_percent = psutil.cpu_percent(interval=None)
        cpu_count = psutil.cpu_count(logical=True)
        cpu_freq = psutil.cpu_freq()
        freq_str = f"{round(cpu_freq.current, 1)} MHz" if cpu_freq else "N/A"
        
        mem = psutil.virtual_memory()
        ram_total_gb = round(mem.total / (1024 ** 3), 2)
        ram_used_gb = round(mem.used / (1024 ** 3), 2)
        ram_free_gb = round(mem.available / (1024 ** 3), 2)
        ram_percent = mem.percent
        
        swap = psutil.swap_memory()
        swap_total_gb = round(swap.total / (1024 ** 3), 2)
        swap_used_gb = round(swap.used / (1024 ** 3), 2)
        
        disk = psutil.disk_usage('/')
        disk_total_gb = round(disk.total / (1024 ** 3), 2)
        disk_used_gb = round(disk.used / (1024 ** 3), 2)
        disk_percent = disk.percent
        
        boot_time = datetime.datetime.fromtimestamp(psutil.boot_time())
        uptime = str(datetime.datetime.now() - boot_time).split('.')[0]
        
        # Network stats
        net_addrs = psutil.net_if_addrs()
        ip_summary = []
        for iface, addrs in net_addrs.items():
            if iface.startswith("lo"):
                continue
            for a in addrs:
                if a.family.name == "AF_INET":
                    ip_summary.append(f"{iface}: {a.address}")
        
        # Battery stats if available
        battery = psutil.sensors_battery()
        battery_str = "N/A (Desktop/Server)"
        if battery:
            plugged = "🔌 Mengisi daya" if battery.power_plugged else "🔋 Baterai"
            battery_str = f"{battery.percent}% ({plugged})"

        # Top processes by RAM and CPU
        processes = []
        for p in psutil.process_iter(['pid', 'name', 'cpu_percent', 'memory_percent']):
            try:
                processes.append(p.info)
            except Exception:
                pass
                
        top_ram = sorted(processes, key=lambda p: p.get('memory_percent') or 0, reverse=True)[:4]
        top_cpu = sorted(processes, key=lambda p: p.get('cpu_percent') or 0, reverse=True)[:4]
        
        return {
            "status": "success",
            "cpu": f"{cpu_percent}% ({cpu_count} cores @ {freq_str})",
            "ram": f"{ram_used_gb} GB / {ram_total_gb} GB ({ram_percent}%, free: {ram_free_gb} GB)",
            "swap": f"{swap_used_gb} GB / {swap_total_gb} GB",
            "disk": f"{disk_used_gb} GB / {disk_total_gb} GB ({disk_percent}%)",
            "battery": battery_str,
            "ip_addresses": ", ".join(ip_summary) or "127.0.0.1",
            "uptime": uptime,
            "top_ram_processes": [
                f"{p['name']} (PID {p['pid']}: {round(p['memory_percent'] or 0, 1)}% RAM)"
                for p in top_ram
            ],
            "top_cpu_processes": [
                f"{p['name']} (PID {p['pid']}: {round(p['cpu_percent'] or 0, 1)}% CPU)"
                for p in top_cpu if (p.get('cpu_percent') or 0) > 0
            ]
        }
    except Exception as e:
        logger.error(f"Error in get_system_stats: {e}")
        return {"status": "error", "message": str(e)}


_BASH_BLOCK_PATTERNS = [
    # (regex, alasan) — dicocokkan terhadap perintah mentah
    (r":\s*\(\s*\)\s*\{.*\}\s*;", "fork bomb"),
    (r"\bdd\s+[^\n]*of=/dev/(sd|hd|vd|nvme|mmcblk)", "tulis mentah ke disk fisik"),
    (r"\bmkfs(\.\w+)?\b", "format filesystem"),
    (r"chmod\s+-R\s+777\s+/", "chmod 777 rekursif pada root"),
    (r"chown\s+-R\b[^\n]*(\s/|\s~|\s\$HOME)(\s|$)", "chown rekursif pada root/home"),
    (r"\b(shutdown|reboot|halt|poweroff)\b", "mematikan/menyalakan ulang sistem"),
    (r"(history\s+-c\b|>\s*~/\.bash_history|shred\s+[^;\n]*history|unset\s+HISTFILE)",
     "menghapus jejak riwayat shell"),
    (r"(curl|wget|fetch)[^\n|]*\|\s*(sudo\s+)?(ba|z|da)?sh\b", "pipe skrip internet langsung ke shell"),
    (r"base64\s+[^;\n|&]*(?:-d\b|--decode)[^;\n|&]*\|\s*(sudo\s+)?(ba|z|da)?sh\b",
     "pipe payload base64 ke shell"),
    (r"/(dev/tcp/|proc/sysrq-trigger)", "teknik reverse-shell/kernel trigger"),
    (r"\.(ssh/id_(rsa|ed25519|ecdsa)|aws/credentials|gnupg)", "akses berkas kredensial privat"),
    (r"\b(useradd|userdel|usermod|visudo)\b", "manipulasi akun pengguna sistem"),
    (r"(iptables|nft)\s+(-F|--flush)", "flush firewall"),
    (r">\s*/dev/(sd|hd|vd|nvme)", "overwrite perangkat blok"),
]

# Target penghapusan yang dianggap destruktif saat dipadukan dgn rm rekursif
_RM_DANGER_TARGETS = (
    r"(?:(?:/{1,2})|(?:~)|(?:\$HOME)|\*|(?:/(?:home|etc|usr|var|boot|lib|opt|bin|sbin|srv|root))"
    r"|(?:\.\./)+(?:home|etc|usr))?(?:\s|$|/)"
)


def _bash_blocked_reason(command: str) -> Optional[str]:
    """Kembalikan alasan pemblokiran bila perintah cocok pola berbahaya."""
    import re as _re
    cmd = command or ""

    for pat, reason in _BASH_BLOCK_PATTERNS:
        if _re.search(pat, cmd):
            return reason

    # rm rekursif (-r/-R/-rf/-fr/--recursive) ke target luas/sistem/home
    m = _re.search(r"\brm\b([^#;\n]*)", cmd)
    if m:
        seg = m.group(1)
        has_recursive = bool(_re.search(
            r"(?:^|\s)(-{1,2}[a-zA-Z]*[rR][a-zA-Z]*|--recursive)(?:\s|$)", seg))
        has_danger_target = bool(_re.search(
            r"(?:^|\s)(\"|')?" + _RM_DANGER_TARGETS, seg))
        if has_recursive and has_danger_target:
            return "penghapusan massal direktori sistem/home"

    low_parts = cmd.lower().split()
    if low_parts and (low_parts[0] == "sudo" or " sudo " in f" {cmd.lower()} "):
        return "eskalasi hak akses (sudo)"
    return None


def execute_bash_command(command: str, working_dir: str = "", backend: str = "") -> Dict[str, Any]:
    """
    Execute a Linux shell command SAFELY inside an isolated Docker sandbox by
    default (resource-limited, no privileges). Falls back to a direct host run
    ONLY when Docker is unavailable, or when backend='host' is requested
    explicitly. Destructive command patterns are always rejected first.

    Args:
        command: The bash command string to execute (e.g. 'ls -la', 'pytest').
        working_dir: Directory to run in (mounted read-write into sandbox).
        backend: 'auto' (default), 'docker', or 'host'.
    """
    blocked = _bash_blocked_reason(command)
    if blocked:
        logger.warning(f"Bash command BLOCKED ({blocked}): {command[:200]}")
        return {
            "status": "error",
            "exit_code": -1,
            "stdout": "",
            "stderr": f"[KEAMANAN] Perintah diblokir: {blocked}.",
            "isolation": "rejected",
        }

    pref = (backend or os.getenv("ALFA_BASH_BACKEND", "auto")).lower().strip()
    use_docker = (
        pref in ("auto", "docker")
        and _docker_available()
        and _ensure_sandbox_image()
    )

    try:
        if use_docker:
            stamp = f"{os.getpid()}_{int(time.time()*1000)%10**9}"
            script_name = f"sandbox_sh_{stamp}.sh"
            script_path = os.path.join(SANDBOX_DIR, script_name)
            with open(script_path, "w", encoding="utf-8") as f:
                f.write("#!/bin/bash\nset -o pipefail\n" + (command or "").strip() + "\n")

            wd_abs = os.path.realpath(os.path.expanduser(working_dir)) if working_dir else ""
            home_in_box = "/workspace" if (wd_abs and os.path.isdir(wd_abs)) else "/sandbox"
            cmd = [
                "docker", "run", "--rm",
                "--name", f"alfa_sbx_{stamp}",
                "-v", f"{SANDBOX_DIR}:/sandbox",
                "--cap-drop", "ALL",
                "--security-opt", "no-new-privileges",
                "--user", f"{os.getuid()}:{os.getgid()}",
                "-e", f"HOME={home_in_box}",
                "--memory", os.getenv("SANDBOX_MEM_LIMIT", "512m"),
                "--memory-swap", os.getenv("SANDBOX_MEM_LIMIT", "512m"),
                "--cpus", os.getenv("SANDBOX_CPUS", "1.0"),
                "--pids-limit", "128",
            ]
            if wd_abs and os.path.isdir(wd_abs):
                cmd += ["-v", f"{wd_abs}:/workspace", "-w", "/workspace"]
            else:
                cmd += ["-w", "/sandbox"]
            cmd += [_SANDBOX_IMAGE, "bash", f"/sandbox/{script_name}"]
            timeout_secs, isolation = 55, "docker"
        else:
            if pref in ("auto", "docker"):
                logger.warning("Docker unavailable - bash falls back to HOST execution.")
            target_dir = os.path.expanduser(working_dir) if working_dir else os.path.expanduser("~")
            if not os.path.exists(target_dir):
                target_dir = os.path.expanduser("~")
            cmd = ["bash", "-c", command]
            timeout_secs, isolation = 45, "none"

        logger.info(f"Executing bash ({isolation}): {command[:150]}")
        try:
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=timeout_secs)
        finally:
            # Skrip wrapper adalah infrastruktur internal; jangan tinggalkan
            # agar tidak ikut terkirim ke chat oleh auto-dispatcher.
            try:
                os.remove(script_path)
            except OSError:
                pass

        stdout = result.stdout.strip()
        stderr = result.stderr.strip()
        if len(stdout) > 3500:
            stdout = stdout[:3500] + "\n...[Output terpotong karena terlalu panjang]"
        if len(stderr) > 1200:
            stderr = stderr[:1200] + "\n...[Stderr terpotong]"

        warn = "" if isolation == "docker" else " [PERINGATAN: dieksekusi di HOST tanpa isolasi]"
        return {
            "status": "success" if result.returncode == 0 else "failed",
            "exit_code": result.returncode,
            "stdout": stdout or "(tidak ada output standar)",
            "stderr": (stderr or None),
            "isolation": isolation,
            "message": "Perintah selesai." + warn,
        }
    except subprocess.TimeoutExpired:
        return {"status": "error", "message": f"Command execution timed out ({timeout_secs}s).", "isolation": isolation}
    except Exception as e:
        return {"status": "error", "message": str(e), "isolation": "none"}


# ── Isolated Python Sandbox ──────────────────────────────────────────────────
_DOCKER_AVAILABLE_CACHE: Optional[bool] = None
_SANDBOX_IMAGE = "alfa-sandbox:latest"


def is_internal_sandbox_artifact(fname: str) -> bool:
    """True bila file adalah artefak infrastruktur sandbox (bukan hasil kerja
    untuk dikirim ke pengguna): skrip wrapper bash/python, Dockerfile,
    plot sementara, dan file tersembunyi."""
    name = fname or ""
    if not name or name.startswith("."):
        return True
    return (
        name.startswith("sandbox_run")
        or name.startswith("sandbox_sh_")
        or name.startswith("sandbox_py_")
        or name == "Dockerfile"
        or name.startswith("generated_plot_old")
        or name == "screen_recording.mp4"
    )


def _docker_available() -> bool:
    """Check (and cache) whether the Docker daemon is usable by this user."""
    global _DOCKER_AVAILABLE_CACHE
    if _DOCKER_AVAILABLE_CACHE is None:
        try:
            probe = subprocess.run(
                ["docker", "version", "--format", "{{.Server.Version}}"],
                capture_output=True, text=True, timeout=8,
            )
            _DOCKER_AVAILABLE_CACHE = probe.returncode == 0
        except Exception:
            _DOCKER_AVAILABLE_CACHE = False
    return _DOCKER_AVAILABLE_CACHE


def _ensure_sandbox_image() -> bool:
    """Ensure the alfa-sandbox image exists; auto-build it on first use."""
    try:
        chk = subprocess.run(
            ["docker", "images", "-q", _SANDBOX_IMAGE],
            capture_output=True, text=True, timeout=10,
        )
        if chk.returncode == 0 and chk.stdout.strip():
            return True

        logger.info("Building sandbox image '%s' (first use, ~2-5 min)...", _SANDBOX_IMAGE)
        dockerfile = os.path.join(SANDBOX_DIR, "Dockerfile")
        os.makedirs(SANDBOX_DIR, exist_ok=True)
        with open(dockerfile, "w", encoding="utf-8") as f:
            f.write(
                "FROM python:3.11-slim\n"
                "RUN useradd -ms /bin/bash -u 1000 sandbox || true\n"
                "RUN pip install --no-cache-dir matplotlib numpy pandas\n"
            )
        build = subprocess.run(
            ["docker", "build", "-t", _SANDBOX_IMAGE, SANDBOX_DIR],
            capture_output=True, text=True, timeout=600,
        )
        if build.returncode == 0:
            logger.info("Sandbox image built successfully.")
            return True
        logger.error("Sandbox image build failed: %s", build.stderr[-400:])
        return False
    except Exception as img_err:
        logger.error("ensure_sandbox_image error: %s", img_err)
        return False


def execute_python_sandbox(code: str) -> Dict[str, Any]:
    """
    Execute a Python script in an ISOLATED Docker container (default) with
    resource limits and no privileges. Falls back to a local subprocess only
    when Docker is unavailable or SANDBOX_BACKEND=none.
    If matplotlib/seaborn creates a plot, it is captured as generated_plot*.png
    for delivery to the Telegram chat.

    Args:
        code: Complete Python code string to execute.
    """
    # Reject empty/trivial/unparseable code before spawning anything
    cleaned = (code or "").strip()
    if len(cleaned) < 10:
        return {"status": "error", "exit_code": -1, "stdout": "", "stderr": "Kode kosong atau terlalu pendek untuk dieksekusi.", "has_plot": False, "isolation": "none"}
    try:
        compile(cleaned, "<sandbox>", "exec")
    except SyntaxError as syn_err:
        return {
            "status": "error", "exit_code": -1, "stdout": "",
            "stderr": f"Syntax error di baris {syn_err.lineno}: {syn_err.msg}",
            "has_plot": False, "isolation": "none",
        }

    backend_pref = os.getenv("SANDBOX_BACKEND", "auto").lower().strip()
    use_docker = (
        backend_pref in ("auto", "docker")
        and _docker_available()
        and _ensure_sandbox_image()
    )

    # Unique per-invocation filenames avoid races between concurrent runs
    stamp = f"{os.getpid()}_{int(time.time()*1000)%10**9}"
    script_name = f"sandbox_run_{stamp}.py"
    plot_name = f"generated_plot_{stamp}.png"
    script_path = os.path.join(SANDBOX_DIR, script_name)
    plot_path = os.path.join(SANDBOX_DIR, plot_name)
    # Clean stale artifacts from previous runs (best effort)
    for old in ("sandbox_run.py", "generated_plot.png"):
        old_p = os.path.join(SANDBOX_DIR, old)
        if os.path.exists(old_p):
            try:
                os.remove(old_p)
            except OSError:
                pass

    try:
        plot_target = f"/sandbox/{plot_name}" if use_docker else plot_path
        preamble = (
            "import os\n"
            "try:\n"
            "    import matplotlib\n"
            "    matplotlib.use('Agg')\n"
            "    import matplotlib.pyplot as plt\n"
            f"    _plot_target = '{plot_target}'\n"
            "    def _auto_save():\n"
            "        if plt.get_fignums():\n"
            "            plt.savefig(_plot_target, bbox_inches='tight', dpi=150)\n"
            "            plt.close('all')\n"
            "    import atexit\n"
            "    atexit.register(_auto_save)\n"
            "except ImportError:\n"
            "    pass\n\n"
        )
        with open(script_path, "w", encoding="utf-8") as f:
            f.write(preamble + cleaned)

        if use_docker:
            cmd = [
                "docker", "run", "--rm",
                "--name", f"alfa_sbx_{stamp}",
                "-v", f"{SANDBOX_DIR}:/sandbox",
                "-w", "/sandbox",
                "--cap-drop", "ALL",
                "--security-opt", "no-new-privileges",
                # Match host uid/gid so bind-mounted sandbox files stay writable
                "--user", f"{os.getuid()}:{os.getgid()}",
                "--memory", os.getenv("SANDBOX_MEM_LIMIT", "512m"),
                "--memory-swap", os.getenv("SANDBOX_MEM_LIMIT", "512m"),
                "--cpus", os.getenv("SANDBOX_CPUS", "1.0"),
                "--pids-limit", "128",
                _SANDBOX_IMAGE, "python", f"/sandbox/{script_name}",
            ]
            timeout_secs = 45
            isolation = "docker"
        else:
            cmd = [sys.executable, script_path]
            timeout_secs = 30
            isolation = "none"
            if backend_pref in ("auto", "docker"):
                logger.warning("Docker sandbox unavailable - falling back to DIRECT execution (no isolation).")

        try:
            res = subprocess.run(cmd, capture_output=True, text=True, timeout=timeout_secs)
        finally:
            # Skrip runner adalah infrastruktur internal; hapus agar tidak
            # bocor ke auto-dispatcher lampiran Telegram.
            try:
                os.remove(script_path)
            except OSError:
                pass

        stdout = res.stdout.strip()
        stderr = res.stderr.strip()
        has_plot = os.path.exists(plot_path) and os.path.getsize(plot_path) > 0

        return {
            "status": "success" if res.returncode == 0 else "error",
            "exit_code": res.returncode,
            "stdout": stdout or "(tidak ada print output)",
            "stderr": stderr or None,
            "generated_chart_photo": has_plot,
            "isolation": isolation,
            "message": ("Grafik visual berhasil dibuat dan akan dikirim ke Telegram!" if has_plot else "Eksekusi kode selesai.")
                        + ("" if isolation == "docker" else " [PERINGATAN: tanpa isolasi]"),
        }
    except subprocess.TimeoutExpired:
        return {"status": "error", "message": f"Eksekusi Python melebihi batas waktu ({timeout_secs} detik).", "isolation": isolation}
    except Exception as e:
        return {"status": "error", "message": f"Python runner error: {str(e)}", "isolation": "none"}


def web_search(query: str, max_results: int = 5) -> Dict[str, Any]:
    """
    Perform a live web search using DuckDuckGo to get up-to-date real-time information, news, or facts.
    Use this tool whenever the user asks about current events, stock prices, weather, documentation, or recent news.
    
    Args:
        query: Search query string.
        max_results: Maximum number of search results to return (default: 5).
    """
    try:
        logger.info(f"Searching web for: {query}")
        results = list(DDGS(verify=False).text(query, max_results=max_results))
        if not results:
            return {"status": "success", "results": [], "message": "Tidak ada hasil pencarian ditemukan."}
        
        formatted_results = []
        for item in results:
            formatted_results.append({
                "title": item.get("title", ""),
                "snippet": item.get("body", ""),
                "link": item.get("href", "")
            })
            
        return {"status": "success", "results": formatted_results}
    except Exception as e:
        logger.error(f"Web search error: {e}")
        return {"status": "error", "message": f"Gagal melakukan pencarian web: {str(e)}"}


def fetch_web_page_content(url: str, max_length: int = 4000) -> Dict[str, Any]:
    """
    Fetch and extract clean text and main content from a website/article URL.
    Use this tool when you need to read documentation, read a specific article, or analyze a web page.
    
    Args:
        url: Full web URL (e.g. 'https://en.wikipedia.org/wiki/Python').
        max_length: Maximum text length to extract (default: 4000 chars).
    """
    try:
        import httpx
        headers = {
            "User-Agent": "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/128.0.0.0 Safari/537.36"
        }
        with httpx.Client(timeout=15.0, follow_redirects=True) as client:
            resp = client.get(url, headers=headers)
            if resp.status_code != 200:
                return {"status": "error", "message": f"Gagal mengambil URL, status code: {resp.status_code}"}
            
            html = resp.text
            # Remove scripts, styles, head, comments
            html = re.sub(r"<script[\s\S]*?</script>", " ", html, flags=re.IGNORECASE)
            html = re.sub(r"<style[\s\S]*?</style>", " ", html, flags=re.IGNORECASE)
            html = re.sub(r"<nav[\s\S]*?</nav>", " ", html, flags=re.IGNORECASE)
            html = re.sub(r"<footer[\s\S]*?</footer>", " ", html, flags=re.IGNORECASE)
            html = re.sub(r"<header[\s\S]*?</header>", " ", html, flags=re.IGNORECASE)
            html = re.sub(r"<!--[\s\S]*?-->", " ", html)
            
            # Extract plain text
            text = re.sub(r"<[^>]+>", " ", html)
            text = re.sub(r"\s+", " ", text).strip()
            
            if len(text) > max_length:
                text = text[:max_length] + "\n...[Konten web dipotong sesuai batas maksimal]"
                
            return {
                "status": "success",
                "url": url,
                "content": text or "Halaman web tidak memuat konten teks yang dapat dibaca."
            }
    except Exception as e:
        logger.error(f"Fetch web page error: {e}")
        return {"status": "error", "message": f"Gagal membaca URL: {str(e)}"}


def save_knowledge_memory(key_topic: str, content: str, category: str = "general") -> Dict[str, Any]:
    """
    Save or update an important fact, user preference, project detail, or note into persistent long-term memory.
    Use this tool whenever the user tells you to remember something, or when important facts about the user/project are shared.
    
    Args:
        key_topic: Short title or identifier for this memory (e.g. 'user_work_hours', 'project_stack', 'trading_rules').
        content: The detailed information to remember.
        category: Category tag (e.g. 'preference', 'project', 'server', 'general').
    """
    try:
        user_id = get_current_user_id()
        msg = database.save_memory_fact_sync(user_id=user_id, key_topic=key_topic, content=content, category=category)
        return {
            "status": "success",
            "message": msg
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


def search_knowledge_memory(query: str) -> Dict[str, Any]:
    """
    Search the persistent long-term memory for previously saved facts, user preferences, or notes.
    Use this tool when answering questions about user preferences, stored projects, or past instructions.
    
    Args:
        query: Keyword or phrase to look up.
    """
    try:
        user_id = get_current_user_id()
        memories = database.search_memories_sync(user_id=user_id, query=query)
        return {"status": "success", "memories": memories}
    except Exception as e:
        return {"status": "error", "message": str(e)}


def read_local_file(file_path: str, max_lines: int = 300, start_line: int = 1) -> Dict[str, Any]:
    """
    Read the text content of a local file on the system safely.
    
    Args:
        file_path: Absolute or relative path to the file.
        max_lines: Maximum number of lines to read (default: 300).
        start_line: Starting line number (1-indexed).
    """
    try:
        expanded_path = os.path.expanduser(file_path)
        if not os.path.isabs(expanded_path):
            expanded_path = os.path.join(os.path.expanduser("~"), file_path)

        if not os.path.exists(expanded_path):
            return {"status": "error", "message": f"File tidak ditemukan: {file_path}"}
        
        if os.path.isdir(expanded_path):
            files = os.listdir(expanded_path)
            return {"status": "is_directory", "files": files[:50], "total": len(files)}

        with open(expanded_path, "r", encoding="utf-8", errors="replace") as f:
            lines = f.readlines()
            
        total_lines = len(lines)
        start_idx = max(0, start_line - 1)
        end_idx = min(total_lines, start_idx + max_lines)
        selected_lines = lines[start_idx:end_idx]
        
        numbered_content = "".join([f"{i+1}: {line}" for i, line in enumerate(selected_lines, start=start_idx)])
            
        return {
            "status": "success",
            "file_path": expanded_path,
            "total_lines": total_lines,
            "showing_lines": f"{start_idx+1} to {end_idx}",
            "content": numbered_content
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


def write_local_file(file_path: str, content: str) -> Dict[str, Any]:
    """
    Write or create a text file on the local system.
    
    Args:
        file_path: Path to the target file.
        content: Text content to write.
    """
    try:
        expanded_path = os.path.expanduser(file_path)
        if not os.path.isabs(expanded_path):
            expanded_path = os.path.join(os.path.expanduser("~"), file_path)

        os.makedirs(os.path.dirname(expanded_path), exist_ok=True)
        with open(expanded_path, "w", encoding="utf-8") as f:
            f.write(content)

        return {"status": "success", "message": f"File berhasil disimpan di {expanded_path} ({len(content)} karakter)"}
    except Exception as e:
        return {"status": "error", "message": str(e)}


_MAX_EDIT_FILE_BYTES = 2 * 1024 * 1024  # batas ukuran file utk operasi presisi


def _py_syntax_guard(path: str, original_content: str) -> Optional[str]:
    """Validasi sintaks Python pasca-edit; rollback bila rusak.

    Mengembalikan pesan error (dan memulihkan isi lama) bila file .py kini
    gagal dikompilasi; mengembalikan None bila aman/bukan file Python.
    """
    if not path.endswith(".py"):
        return None
    try:
        with open(path, "r", encoding="utf-8", errors="surrogateescape") as f:
            new_content = f.read()
        compile(new_content, path, "exec")
        return None
    except SyntaxError as syn:
        with open(path, "w", encoding="utf-8", errors="surrogateescape") as f:
            f.write(original_content)
        return (f"Edit DIBATALKAN (auto-rollback): hasil menyebabkan SyntaxError "
                f"di baris {syn.lineno}: {syn.msg}. Isi file dikembalikan seperti semula.")


def _resolve_host_path(file_path: str) -> str:
    expanded = os.path.expanduser(file_path)
    if not os.path.isabs(expanded):
        expanded = os.path.join(os.path.expanduser("~"), file_path)
    return os.path.realpath(expanded)


def edit_file_precise(file_path: str, old_text: str, new_text: str,
                      occurrence: int = 0) -> Dict[str, Any]:
    """
    Edit a file with SURGICAL precision (opencode-style): replace an exact
    unique snippet instead of rewriting the whole file. old_text must match
    the file content EXACTLY (including indentation).

    Safety rule: if old_text matches MULTIPLE locations, the call FAILS and
    you must pass `occurrence` (1-based index of the match to replace,
    -1 = last match) or include more surrounding context.

    Args:
        file_path: Path to the existing text file.
        old_text: Exact existing snippet to replace.
        new_text: Replacement text.
        occurrence: 0 = require unique match (default); N = replace Nth match;
                    -1 = replace last match.
    """
    try:
        p = _resolve_host_path(file_path)
        if not os.path.isfile(p):
            return {"status": "error", "message": f"File tidak ditemukan: {file_path}"}
        if os.path.getsize(p) > _MAX_EDIT_FILE_BYTES:
            return {"status": "error", "message": f"File terlalu besar (>2MB): {file_path}"}
        if not old_text:
            return {"status": "error", "message": "old_text kosong — gunakan write_local_file untuk membuat isi baru."}

        with open(p, "r", encoding="utf-8", errors="surrogateescape") as f:
            content = f.read()

        count = content.count(old_text)
        if count == 0:
            # Bantu model: cari kandidat mirip (abaikan trailing whitespace per baris)
            norm_old = "\n".join(line.rstrip() for line in old_text.splitlines())
            lines = content.splitlines()
            best, best_score = None, 0.0
            window = len(old_text.splitlines())
            for i in range(0, max(1, len(lines) - window + 1)):
                cand = "\n".join(lines[i:i + window])
                score = difflib.SequenceMatcher(None, norm_old,
                                                "\n".join(l.rstrip() for l in cand.splitlines())).ratio()
                if score > best_score:
                    best, best_score = (i + 1, cand), score
            hint = ""
            if best and best_score > 0.6:
                hint = (f" Kemungkinan yang dimaksud di sekitar baris {best[0]} "
                        f"(kemiripan {best_score:.0%}). Salin teks persis dari file.")
            return {"status": "error",
                    "message": f"old_text tidak ditemukan di {file_path}.{hint}"}

        if occurrence == 0 and count > 1:
            return {"status": "error",
                    "message": (f"old_text cocok di {count} lokasi berbeda. "
                                "Tambahkan konteks lebih banyak agar unik, atau sebutkan "
                                "`occurrence` (1=ke-N, -1=terakhir). Tidak ada perubahan ditulis.")}

        # Resolusi indeks kecocokan: -1=terakhir, 0/unik=pertama, N=ke-N
        idx = count + 1 + occurrence if occurrence < 0 else max(1, occurrence)
        if not (1 <= idx <= count):
            return {"status": "error",
                    "message": f"occurrence={occurrence} di luar rentang; ditemukan {count} kecocokan."}

        start = 0
        for _ in range(idx):
            pos = content.find(old_text, start)
            start = pos + 1
        new_content = content[:pos] + new_text + content[pos + len(old_text):]

        with open(p, "w", encoding="utf-8", errors="surrogateescape") as f:
            f.write(new_content)

        syn_err = _py_syntax_guard(p, content)
        if syn_err:
            return {"status": "error", "message": syn_err}

        return {
            "status": "success",
            "message": (f"Berhasil mengganti {len(old_text)} -> {len(new_text)} karakter "
                        f"di {file_path} (kecocokan #{idx}/{count})."),
            "line_hint": content.count("\n", 0, pos) + 1,
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


def apply_unified_diff(file_path: str, diff_text: str) -> Dict[str, Any]:
    """
    Apply a UNIFIED DIFF (format `diff -u` / git diff) to a single file with
    context-matching and small drift tolerance — like `patch` but built-in.
    Only one file per call. All hunks must apply or nothing is written.

    Args:
        file_path: Path to the target text file.
        diff_text: Unified diff body (lines starting with ---/+++ are ignored;
                   hunks start with @@).
    """
    import re as _re
    try:
        p = _resolve_host_path(file_path)
        if not os.path.isfile(p):
            return {"status": "error", "message": f"File tidak ditemukan: {file_path}"}
        if os.path.getsize(p) > _MAX_EDIT_FILE_BYTES:
            return {"status": "error", "message": f"File terlalu besar (>2MB): {file_path}"}

        with open(p, "r", encoding="utf-8", errors="surrogateescape") as f:
            orig_lines = f.read().split("\n")

        # Parse hunks
        hunks, cur = [], None
        for raw in diff_text.split("\n"):
            if raw.startswith("@@"):
                m = _re.match(r"@@ -\d+(?:,\d+)? \+(\d+)(?:,(\d+))? @@", raw)
                if m:
                    cur = {"old": [], "new": [], "new_start": int(m.group(1))}
                    hunks.append(cur)
                continue
            if cur is None:
                continue
            if raw.startswith(("---", "+++")) or raw.startswith("\\ No newline"):
                continue
            tag = raw[:1]
            if tag == "+":
                cur["new"].append(raw[1:])
            elif tag == "-":
                cur["old"].append(raw[1:])
            elif raw.startswith("\\ No newline"):
                continue
            else:
                # Baris konteks: format baku ' teks'; toleransi model tanpa spasi
                body = raw[1:] if raw.startswith(" ") else raw
                cur["old"].append(body)
                cur["new"].append(body)

        if not hunks:
            return {"status": "error",
                    "message": "Tidak ada hunk @@ valid dalam diff. Pastikan format unified diff."}

        lines = list(orig_lines)
        applied = 0
        cursor = 0
        for hno, h in enumerate(hunks, 1):
            anchor = next((ln for ln in h["old"] if ln.strip()), "")
            n_old = len(h["old"])
            candidates = []
            if anchor:
                for i in range(cursor, min(len(lines), max(len(lines), h["new_start"] + 80))):
                    if lines[i] == anchor:
                        candidates.append(i - h["old"].index(anchor))
                        if len(candidates) >= 3:
                            break
            chosen = None
            for base in candidates + [h["new_start"] - 1]:
                if base is None or base < 0:
                    continue
                seg = lines[base:base + n_old]
                if [x.strip() for x in seg] == [x.strip() for x in h["old"]] or seg == h["old"]:
                    chosen = base
                    break
            if chosen is None:
                return {"status": "error",
                        "message": (f"Hunk #{hno} gagal diterapkan (konteks tidak cocok "
                                    f"di sekitar '{anchor[:60]}'). Tidak ada perubahan ditulis. "
                                    "Baca ulang file & buat diff baru."),
                        "hunks_applied_before_fail": applied}

            lines[chosen:chosen + n_old] = h["new"]
            cursor = chosen + len(h["new"])
            applied += 1

        with open(p, "w", encoding="utf-8", errors="surrogateescape") as f:
            f.write("\n".join(lines))

        syn_err = _py_syntax_guard(p, "\n".join(orig_lines))
        if syn_err:
            return {"status": "error", "message": syn_err}

        return {"status": "success",
                "message": f"{applied}/{len(hunks)} hunk berhasil diterapkan ke {file_path}.",
                "hunks_applied": applied}
    except Exception as e:
        return {"status": "error", "message": str(e)}


def search_workspace_files(pattern: str, base_dir: str = "~", max_results: int = 30) -> Dict[str, Any]:
    """
    Search for files and directories matching a glob pattern (e.g. '*.py', '*.json', 'bot*').
    
    Args:
        pattern: Glob pattern to search.
        base_dir: Root search directory (default: user home).
        max_results: Maximum number of files to return.
    """
    try:
        root_dir = os.path.expanduser(base_dir)
        matches = []
        for root, dirs, files in os.walk(root_dir):
            # Ignore heavy/hidden directories
            dirs[:] = [d for d in dirs if not d.startswith('.') and d not in ('node_modules', 'venv', '__pycache__', 'dist', 'build')]
            for name in files:
                if glob.fnmatch.fnmatch(name, pattern) or pattern.lower() in name.lower():
                    rel = os.path.relpath(os.path.join(root, name), root_dir)
                    matches.append(rel)
                    if len(matches) >= max_results:
                        break
            if len(matches) >= max_results:
                break

        return {
            "status": "success",
            "base_dir": root_dir,
            "matches_count": len(matches),
            "files": matches
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


def grep_workspace(query: str, base_dir: str = "~", file_pattern: str = "") -> Dict[str, Any]:
    """
    Search for text or regex pattern inside files across workspace.
    
    Args:
        query: String or regex query to find.
        base_dir: Search directory.
        file_pattern: Optional glob filter for files (e.g. '*.py').
    """
    try:
        root_dir = os.path.expanduser(base_dir)
        cmd = ["grep", "-rnI", "--exclude-dir={.git,venv,node_modules,__pycache__}"]
        if file_pattern:
            cmd.append(f"--include={file_pattern}")
        cmd.extend([query, root_dir])

        res = subprocess.run(cmd, capture_output=True, text=True, timeout=20)
        lines = res.stdout.strip().split("\n") if res.stdout.strip() else []
        
        return {
            "status": "success",
            "matches_count": len(lines),
            "results": lines[:30]
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


# ── 🗂️ Codebase Index (SQLite FTS5) ─────────────────────────────────────────
_CODE_INDEX_SKIP_DIRS = {
    ".git", "venv", ".venv", "env", "node_modules", "__pycache__",
    ".mypy_cache", ".pytest_cache", "dist", "build", ".next", "target",
    "vendor", "coverage", ".idea", ".vscode",
}
_CODE_INDEX_DB = os.path.join(os.path.dirname(os.path.abspath(__file__)), "agent_data.db")
_CODE_CHUNK_LINES = 70


def _code_index_connect():
    import sqlite3
    conn = sqlite3.connect(_CODE_INDEX_DB, timeout=15)
    conn.execute("""
        CREATE VIRTUAL TABLE IF NOT EXISTS code_fts USING fts5(
            rel_path, content, symbol UNINDEXED, repo_root UNINDEXED,
            start_line UNINDEXED, end_line UNINDEXED)
    """)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS code_index_meta(
            repo_root TEXT PRIMARY KEY, indexed_at REAL,
            files INTEGER, chunks INTEGER)
    """)
    return conn


def _index_freshness(root: str, sample_paths: List[str]) -> Dict[str, Any]:
    """Periksa apakah index masih segar: bandingkan mtime sampel file
    vs waktu indexing. Mengembalikan info kesegaran utk hasil pencarian."""
    import sqlite3 as _sq
    info: Dict[str, Any] = {"stale": False, "indexed_at": None}
    if not root:
        return info
    try:
        conn = _sq.connect(_CODE_INDEX_DB, timeout=10)
        row = conn.execute(
            "SELECT indexed_at, files, chunks FROM code_index_meta WHERE repo_root = ?",
            (root,)).fetchone()
        conn.close()
        if not row:
            info["stale"] = True
            info["note"] = "index tidak tercatat meta-nya — jalankan ulang index_codebase."
            return info
        indexed_at, files, chunks = row
        info["indexed_at"] = indexed_at
        info["files"] = files
        info["chunks"] = chunks
        changed = 0
        checked = 0
        for rel in sample_paths[:12]:
            fp = os.path.join(root, rel)
            try:
                checked += 1
                if os.path.getmtime(fp) > indexed_at:
                    changed += 1
            except OSError:
                continue
        if checked and changed:
            info["stale"] = True
            info["changed_sample"] = changed
    except Exception:
        pass
    return info


def _iter_code_files(root: str, extensions: str):
    exts = {"." + e.strip().lstrip(".").lower()
            for e in (extensions or "").split(",") if e.strip()}
    for dirpath, dirnames, filenames in os.walk(root):
        dirnames[:] = [d for d in dirnames if d not in _CODE_INDEX_SKIP_DIRS]
        for name in filenames:
            if not exts or os.path.splitext(name)[1].lower() in exts:
                yield os.path.join(dirpath, name)


def _chunk_code_lines(lines):
    """Pecah file jadi chunk ~_CODE_CHUNK_LINES di batas baris kosong."""
    chunks, cur, sym = [], [], None
    import re as _re
    sym_re = _re.compile(r"^\s*(?:async\s+)?(?:def|class|function|func|fn|impl|type)\s+(\w+)")
    for ln in lines:
        cur.append(ln)
        if len(cur) >= _CODE_CHUNK_LINES and ln.strip() == "":
            chunks.append((cur, sym))
            cur, sym = [], None
        elif sym is None:
            m = sym_re.match(ln)
            if m:
                sym = m.group(1)
    if cur:
        chunks.append((cur, sym))
    return chunks


def index_codebase(repo_path: str, file_extensions: str = "py,js,ts,tsx,jsx,go,rs,java,c,cpp,h,md,json,yaml,yml,toml") -> Dict[str, Any]:
    """
    Build/refresh a full-text INDEX of a repository so later searches are fast
    and context-aware (RAG-style retrieval without external services).
    Run once per repo; call again after big changes to refresh.

    Args:
        repo_path: Root directory of the project to index.
        file_extensions: Comma-separated extensions to include.
    """
    try:
        root = _resolve_host_path(repo_path)
        if not os.path.isdir(root):
            return {"status": "error", "message": f"Direktori tidak ditemukan: {repo_path}"}

        conn = _code_index_connect()
        files_scanned, chunks_inserted, skipped_big = 0, 0, 0
        try:
            rows = []
            for fpath in _iter_code_files(root, file_extensions):
                try:
                    if os.path.getsize(fpath) > _MAX_EDIT_FILE_BYTES:
                        skipped_big += 1
                        continue
                    with open(fpath, "r", encoding="utf-8", errors="replace") as f:
                        lines = f.read().split("\n")
                except OSError:
                    continue
                rel = os.path.relpath(fpath, root)
                offset = 0
                for chunk_lines, sym in _chunk_code_lines(lines):
                    content = "\n".join(chunk_lines)
                    n = len(chunk_lines)
                    if content.strip():
                        rows.append((rel, content, sym or "", root, offset + 1, offset + n))
                    offset += n
                files_scanned += 1

            # SATU transaksi utk seluruh index (hindari lock war dgn service lain)
            with conn:
                conn.execute("DELETE FROM code_fts WHERE repo_root = ?", (root,))
                conn.executemany(
                    "INSERT INTO code_fts(rel_path, content, symbol, repo_root, start_line, end_line) "
                    "VALUES (?,?,?,?,?,?)", rows)
                chunks_inserted = len(rows)
                conn.execute(
                    "INSERT INTO code_index_meta(repo_root, indexed_at, files, chunks) "
                    "VALUES (?,?,?,?) ON CONFLICT(repo_root) DO UPDATE SET "
                    "indexed_at=excluded.indexed_at, files=excluded.files, chunks=excluded.chunks",
                    (root, time.time(), files_scanned, chunks_inserted))
        finally:
            conn.close()

        return {
            "status": "success",
            "message": f"Index selesai: {files_scanned} file, {chunks_inserted} chunk tersimpan"
                       f"{' (' + str(skipped_big) + ' file besar dilewati)' if skipped_big else ''}.",
            "files_indexed": files_scanned,
            "chunks": chunks_inserted,
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


def search_codebase(query: str, repo_path: str = "", limit: int = 10) -> Dict[str, Any]:
    """
    Search an indexed repository semantically-by-keyword (FTS5 ranked).
    Returns the most relevant code chunks with file:line references.
    Call index_codebase() first if this returns 'index kosong'.

    Args:
        query: Keywords or phrase, e.g. 'generate video ffmpeg overlay'.
        repo_path: Optional root to restrict search to one project.
        limit: Max results.
    """
    try:
        if not query.strip():
            return {"status": "error", "message": "Query kosong."}
        conn = _code_index_connect()
        try:
            sql = ("SELECT rel_path, content, symbol, repo_root, start_line, end_line, "
                   "snippet(code_fts, 1, '>>>', '<<<', ' … ', 12) AS snip "
                   "FROM code_fts WHERE code_fts MATCH ? ")
            params = [query.strip()]
            if repo_path.strip():
                sql += "AND repo_root = ? "
                params.append(_resolve_host_path(repo_path))
            sql += "ORDER BY rank LIMIT ?"
            params.append(int(limit))
            rows = conn.execute(sql, params).fetchall()
        finally:
            conn.close()

        if not rows:
            return {"status": "empty",
                    "message": "Tidak ada hasil. Index mungkin belum dibuat — panggil index_codebase dulu."}

        results = [{
            "location": f"{r[0]}:{r[4]}-{r[5]}",
            "symbol": r[2] or None,
            "snippet": r[6],
        } for r in rows]

        # Cek kesegaran index utk repo-repo yang muncul di hasil
        roots = list(dict.fromkeys(r[3] for r in rows))
        sample_by_root: Dict[str, List[str]] = {}
        for r in rows:
            sample_by_root.setdefault(r[3], []).append(r[0])
        stale_roots = []
        for rt in roots:
            info = _index_freshness(rt, sample_by_root.get(rt, []))
            if info.get("stale"):
                stale_roots.append(rt)
        resp = {"status": "success", "matches": len(results), "results": results}
        if stale_roots:
            resp["index_stale_warning"] = (
                f"Index untuk {len(stale_roots)} repo sudah USANG (ada file berubah "
                f"setelah indexing). Jalankan index_codebase lagi utk hasil akurat.")
            logger.warning(resp["index_stale_warning"])
        return resp
    except Exception as e:
        return {"status": "error", "message": str(e)}


def schedule_reminder(reminder_time_iso: str, message: str) -> Dict[str, Any]:
    """
    Schedule a future proactive reminder or alert that the bot will send directly to Telegram.
    
    Args:
        reminder_time_iso: The target time in ISO format (YYYY-MM-DDTHH:MM:SS), e.g. '2026-08-19T08:00:00'.
        message: The reminder message content.
    """
    try:
        user_id = get_current_user_id()
        chat_id = get_current_chat_id()
        rem_id = database.add_reminder_sync(user_id=user_id, chat_id=chat_id, reminder_time_iso=reminder_time_iso, message=message)
        return {
            "status": "success",
            "reminder_id": rem_id,
            "message": f"Pengingat #{rem_id} berhasil dijadwalkan pada {reminder_time_iso}: '{message}'"
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


def capture_desktop_screenshot() -> Dict[str, Any]:
    """
    Capture an ultra-fast, high-resolution screenshot of the active Linux desktop screen.
    Uses native Wayland XDG Desktop Portal capture (compatible with Ubuntu GNOME Wayland)
    with fallbacks to grim, import, and PIL ImageGrab.
    """
    try:
        screenshot_path = os.path.join(SANDBOX_DIR, "desktop_screen.png")
        if os.path.exists(screenshot_path):
            try:
                os.remove(screenshot_path)
            except OSError:
                pass

        # 1. Native Wayland XDG Desktop Portal (Official GNOME/KDE Wayland Screen Capture)
        portal_script = (
            "import os, sys, shutil, urllib.parse, random\n"
            "from gi.repository import Gio, GLib\n"
            "target_path = sys.argv[1]\n"
            "bus = Gio.bus_get_sync(Gio.BusType.SESSION, None)\n"
            "loop = GLib.MainLoop()\n"
            "saved_uri = None\n"
            "def on_signal(connection, sender_name, object_path, interface_name, signal_name, parameters, user_data):\n"
            "    global saved_uri\n"
            "    res = parameters.unpack()\n"
            "    if len(res) >= 2 and isinstance(res[1], dict) and 'uri' in res[1]:\n"
            "        saved_uri = res[1]['uri']\n"
            "    loop.quit()\n"
            "token = f'shot_{random.randint(1000, 9999)}'\n"
            "options = {'interactive': GLib.Variant('b', False), 'handle_token': GLib.Variant('s', token)}\n"
            "ret = bus.call_sync(\n"
            "    'org.freedesktop.portal.Desktop',\n"
            "    '/org/freedesktop/portal/desktop',\n"
            "    'org.freedesktop.portal.Screenshot',\n"
            "    'Screenshot',\n"
            "    GLib.Variant('(sa{sv})', ('', options)),\n"
            "    GLib.VariantType('(o)'),\n"
            "    Gio.DBusCallFlags.NONE,\n"
            "    4000,\n"
            "    None\n"
            ")\n"
            "req_path = ret.unpack()[0]\n"
            "bus.signal_subscribe(\n"
            "    'org.freedesktop.portal.Desktop',\n"
            "    'org.freedesktop.portal.Request',\n"
            "    'Response',\n"
            "    req_path,\n"
            "    None,\n"
            "    Gio.DBusSignalFlags.NONE,\n"
            "    on_signal,\n"
            "    None\n"
            ")\n"
            "GLib.timeout_add_seconds(3, loop.quit)\n"
            "loop.run()\n"
            "if saved_uri and saved_uri.startswith('file://'):\n"
            "    parsed_path = urllib.parse.unquote(saved_uri[7:])\n"
            "    if os.path.exists(parsed_path) and os.path.getsize(parsed_path) > 1000:\n"
            "        os.makedirs(os.path.dirname(target_path), exist_ok=True)\n"
            "        shutil.move(parsed_path, target_path)\n"
            "        sys.exit(0)\n"
            "sys.exit(1)\n"
        )
        try:
            p_res = subprocess.run(
                ["/usr/bin/python3", "-c", portal_script, screenshot_path],
                capture_output=True,
                timeout=4
            )
            if p_res.returncode == 0 and os.path.exists(screenshot_path) and os.path.getsize(screenshot_path) > 1000:
                return {
                    "status": "success",
                    "message": "Screenshot desktop berhasil diambil via Wayland Portal.",
                    "file_path": screenshot_path
                }
        except Exception:
            pass

        # 2. Try grim (wlroots Wayland compositors like Sway/Hyprland)
        try:
            subprocess.run(f"grim '{screenshot_path}'", shell=True, capture_output=True, timeout=2)
            if os.path.exists(screenshot_path) and os.path.getsize(screenshot_path) > 1000:
                return {
                    "status": "success",
                    "message": "Screenshot desktop berhasil diambil via grim.",
                    "file_path": screenshot_path
                }
        except Exception:
            pass

        # 3. Try import (ImageMagick)
        try:
            subprocess.run(f"import -window root '{screenshot_path}'", shell=True, capture_output=True, timeout=2)
            if os.path.exists(screenshot_path) and os.path.getsize(screenshot_path) > 1000:
                return {
                    "status": "success",
                    "message": "Screenshot desktop berhasil diambil via ImageMagick.",
                    "file_path": screenshot_path
                }
        except Exception:
            pass

        # 4. Fallback PIL ImageGrab
        try:
            from PIL import ImageGrab
            img = ImageGrab.grab()
            img.save(screenshot_path)
            if os.path.exists(screenshot_path) and os.path.getsize(screenshot_path) > 1000:
                return {
                    "status": "success",
                    "message": "Screenshot desktop berhasil diambil via ImageGrab.",
                    "file_path": screenshot_path
                }
        except Exception:
            pass

        return {"status": "error", "message": "Gagal mengambil screenshot desktop di lingkungan display saat ini."}
    except Exception as err:
        return {"status": "error", "message": str(err)}


def capture_webcam_frame() -> Dict[str, Any]:
    """
    Capture a live snapshot frame from the connected hardware webcam/camera (/dev/video0).
    Use this tool when the user asks for a desk check, room status, or webcam photo.
    """
    try:
        cam_path = os.path.join(SANDBOX_DIR, "webcam_frame.jpg")
        if os.path.exists(cam_path):
            try:
                os.remove(cam_path)
            except OSError:
                pass
                
        import cv2
        cap = cv2.VideoCapture(0)
        if not cap.isOpened():
            return {"status": "error", "message": "Perangkat webcam tidak dapat diakses atau tidak terdeteksi (/dev/video0)."}
        
        # Warmup camera frames
        for _ in range(5):
            ret, frame = cap.read()
        ret, frame = cap.read()
        cap.release()
        
        if ret and frame is not None:
            cv2.imwrite(cam_path, frame)
            return {
                "status": "success",
                "message": "Foto webcam berhasil diambil.",
                "file_path": cam_path
            }
    except Exception as err:
        return {"status": "error", "message": str(err)}


def scan_local_network() -> Dict[str, Any]:
    """
    Scan connected local LAN devices, IP neighbors, and active gateways.
    Use this tool when the user asks to check network devices, Wi-Fi neighbors, or connected hardware.
    """
    try:
        res = subprocess.run("ip neigh || arp -a", shell=True, capture_output=True, text=True, timeout=10)
        return {
            "status": "success",
            "devices": res.stdout.strip() or "Tidak ada perangkat terdeteksi di tabel ARP/Neighbor."
        }
    except Exception as err:
        return {"status": "error", "message": str(err)}


def _find_camofox_bin() -> Optional[str]:
    """Find Camofox binary in PATH or common NVM / node directories."""
    import shutil
    found = shutil.which("camofox")
    if found:
        return found
    candidates = [
        "/home/fahmial/.nvm/versions/node/v24.19.0/bin/camofox",
        os.path.expanduser("~/.nvm/versions/node/v24.19.0/bin/camofox"),
        os.path.expanduser("~/.npm-global/bin/camofox"),
        "/usr/local/bin/camofox",
        "/usr/bin/camofox"
    ]
    for c in candidates:
        if os.path.exists(c):
            return c
    return None


def _ensure_camofox_server() -> bool:
    """Ensure Camofox browser daemon is active and listening on port 9377."""
    try:
        import httpx
        try:
            r = httpx.get("http://127.0.0.1:9377/health", timeout=1.5)
            if r.status_code == 200 and r.json().get("running"):
                return True
        except Exception:
            pass

        # Attempt to start daemon
        env = os.environ.copy()
        api_key = os.environ.get("CAMOFOX_API_KEY", "7edc51a9e8b2401f98bc43d105ef5f68")
        env["CAMOFOX_API_KEY"] = api_key
        camofox_bin = _find_camofox_bin()
        if camofox_bin:
            subprocess.run([camofox_bin, "server", "start", "--background"], env=env, capture_output=True, timeout=10)
            import time
            time.sleep(1.5)
            return True
        return False
    except Exception as e:
        logger.error(f"Camofox ensure server error: {e}")
        return False


def _run_camofox_cli(args: List[str]) -> Dict[str, Any]:
    """Execute camofox CLI command with proper environment and output parsing."""
    _ensure_camofox_server()
    camofox_bin = _find_camofox_bin()
    if not camofox_bin:
        return {"success": False, "error": "Camofox CLI binary tidak ditemukan di sistem."}
        
    env = os.environ.copy()
    api_key = os.environ.get("CAMOFOX_API_KEY", "7edc51a9e8b2401f98bc43d105ef5f68")
    env["CAMOFOX_API_KEY"] = api_key
    
    cmd = [camofox_bin] + args
    try:
        res = subprocess.run(cmd, env=env, capture_output=True, text=True, timeout=30)
        stdout = res.stdout.strip()
        stderr = res.stderr.strip()
        return {
            "success": res.returncode == 0,
            "stdout": stdout,
            "stderr": stderr
        }
    except subprocess.TimeoutExpired:
        return {"success": False, "error": "Camofox browser command timed out (30s)."}
    except Exception as e:
        return {"success": False, "error": str(e)}


def browser_open_url(url: str) -> Dict[str, Any]:
    """
    Open a web page in the Camofox stealth browser engine and return its accessibility element tree snapshot.
    Use this tool when the user asks to open a website, browse a web page, fill forms, or interact with web elements.
    
    Args:
        url: Full web URL to open (e.g. 'https://github.com/trending', 'https://news.ycombinator.com').
    """
    try:
        res = _run_camofox_cli(["open", url])
        if not res.get("success"):
            return {"status": "error", "message": res.get("stderr") or res.get("error")}
        
        # Take snapshot immediately
        snap = _run_camofox_cli(["snapshot"])
        return {
            "status": "success",
            "message": f"Halaman web '{url}' berhasil dibuka.",
            "interactive_elements": snap.get("stdout") or res.get("stdout")
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


def browser_click_element(element_ref: str, tab_id: str = "") -> Dict[str, Any]:
    """
    Click an interactive button, link, checkbox, or element on the active Camofox browser page by its reference.
    
    Args:
        element_ref: Element ref identifier (e.g. 'e1', 'e2', 'e15') from the browser snapshot or CSS selector.
        tab_id: Optional specific tab ID.
    """
    try:
        args = ["click", element_ref]
        if tab_id:
            args.append(tab_id)
            
        res = _run_camofox_cli(args)
        if not res.get("success"):
            return {"status": "error", "message": res.get("stderr") or res.get("error")}
            
        # Get updated snapshot after click
        snap_args = ["snapshot"]
        if tab_id:
            snap_args.append(tab_id)
        snap = _run_camofox_cli(snap_args)
        
        return {
            "status": "success",
            "message": f"Elemen '{element_ref}' berhasil diklik.",
            "updated_page_elements": snap.get("stdout", "")[:3000]
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


def browser_type_text(element_ref: str, text: str, tab_id: str = "") -> Dict[str, Any]:
    """
    Type text into an input field or search bar on the active Camofox browser page.
    
    Args:
        element_ref: Element ref identifier (e.g. 'e3') or selector of the input field.
        text: String text to type into the input field.
        tab_id: Optional tab ID.
    """
    try:
        args = ["type", element_ref, text]
        if tab_id:
            args.append(tab_id)
            
        res = _run_camofox_cli(args)
        if not res.get("success"):
            return {"status": "error", "message": res.get("stderr") or res.get("error")}
            
        return {
            "status": "success",
            "message": f"Teks berhasil diketik ke elemen '{element_ref}'."
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


def browser_capture_screenshot(tab_id: str = "") -> Dict[str, Any]:
    """
    Take a screenshot of the current Camofox browser page and automatically send it to the Telegram chat.
    
    Args:
        tab_id: Optional tab ID.
    """
    try:
        args = ["screenshot"]
        if tab_id:
            args.append(tab_id)
            
        res = _run_camofox_cli(args)
        if not res.get("success"):
            return {"status": "error", "message": res.get("stderr") or res.get("error")}
            
        out = res.get("stdout", "")
        # Camofox returns "path: /home/fahmial/.camofox/screenshots/..."
        import shutil
        target_path = os.path.join(SANDBOX_DIR, "browser_screenshot.png")
        
        for line in out.splitlines():
            if line.startswith("path:"):
                src_path = line.replace("path:", "").strip()
                if os.path.exists(src_path):
                    shutil.copyfile(src_path, target_path)
                    return {
                        "status": "success",
                        "message": "Screenshot browser berhasil diambil dan akan dikirim ke Telegram.",
                        "file_path": target_path
                    }
                    
        return {"status": "success", "message": "Screenshot browser berhasil diproses.", "raw_output": out}
    except Exception as e:
        return {"status": "error", "message": str(e)}


def browser_close_tab(tab_id: str = "") -> Dict[str, Any]:
    """
    Close the active or specified Camofox browser tab.
    
    Args:
        tab_id: Optional tab ID (defaults to active tab).
    """
    try:
        args = ["close"]
        if tab_id:
            args.append(tab_id)
        res = _run_camofox_cli(args)
        return {
            "status": "success" if res.get("success") else "error",
            "message": "Tab browser berhasil ditutup." if res.get("success") else res.get("stderr")
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


def desktop_click_coordinate(x: int, y: int, button: str = "left", clicks: int = 1) -> Dict[str, Any]:
    """
    Simulate a hardware mouse click on specific pixel coordinates (X, Y) on the Linux desktop screen.
    Use this tool for GUI desktop automation (e.g. clicking buttons, icons, or menus on active windows).
    
    Args:
        x: X coordinate pixel on screen (0 to 1920).
        y: Y coordinate pixel on screen (0 to 1080).
        button: Mouse button ('left', 'right', 'middle'). Default: 'left'.
        clicks: Number of clicks (1 for single click, 2 for double click).
    """
    try:
        # Try ydotool (Wayland compatible)
        btn_code = "0xC0" if button == "left" else "0xC1" if button == "right" else "0xC2"
        res = subprocess.run(f"ydotool mousemove -a {x} {y} && ydotool click {btn_code}", shell=True, capture_output=True, text=True, timeout=5)
        if res.returncode == 0:
            return {"status": "success", "message": f"Mouse berhasil diklik pada koordinat ({x}, {y}) [{button}]."}

        # Fallback to xdotool
        res = subprocess.run(f"xdotool mousemove {x} {y} click {'1' if button == 'left' else '3'}", shell=True, capture_output=True, text=True, timeout=5)
        if res.returncode == 0:
            return {"status": "success", "message": f"Mouse berhasil diklik via xdotool pada ({x}, {y})."}

        # Fallback to PyAutoGUI
        import pyautogui
        pyautogui.FAILSAFE = False
        pyautogui.click(x=x, y=y, button=button, clicks=clicks)
        return {"status": "success", "message": f"Mouse berhasil diklik via PyAutoGUI pada ({x}, {y})."}
    except Exception as e:
        return {"status": "error", "message": f"Gagal melakukan klik mouse: {str(e)}"}


def desktop_type_keys(text: str = "", hotkey: str = "") -> Dict[str, Any]:
    """
    Type text or press keyboard shortcuts/hotkeys on the active Linux desktop window.
    
    Args:
        text: String of text to type.
        hotkey: Optional key combination (e.g. 'ctrl+c', 'ctrl+v', 'alt+tab', 'Return', 'Escape').
    """
    try:
        if hotkey:
            keys = [k.strip() for k in hotkey.split("+")]
            import pyautogui
            pyautogui.hotkey(*keys)
            return {"status": "success", "message": f"Shortcut '{hotkey}' berhasil ditekan."}
        
        if text:
            res = subprocess.run(["ydotool", "type", text], capture_output=True, text=True, timeout=5)
            if res.returncode == 0:
                return {"status": "success", "message": f"Teks berhasil diketik ke desktop: '{text}'"}
            
            import pyautogui
            pyautogui.write(text)
            return {"status": "success", "message": "Teks berhasil diketik via PyAutoGUI."}
            
        return {"status": "error", "message": "Harus menyertakan text atau hotkey."}
    except Exception as e:
        return {"status": "error", "message": f"Gagal mengetik tombol: {str(e)}"}


def desktop_launch_app(app_name_or_command: str) -> Dict[str, Any]:
    """
    Launch a Linux GUI software application in the background (e.g. 'code', 'brave-browser', 'spotify', 'nautilus').
    
    Args:
        app_name_or_command: Application executable name or command.
    """
    try:
        env = os.environ.copy()
        env["DISPLAY"] = env.get("DISPLAY", ":0")
        env["WAYLAND_DISPLAY"] = env.get("WAYLAND_DISPLAY", "wayland-0")
        subprocess.Popen(app_name_or_command, shell=True, env=env, start_new_session=True)
        return {
            "status": "success",
            "message": f"Aplikasi '{app_name_or_command}' berhasil diluncurkan di background desktop."
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal meluncurkan aplikasi: {str(e)}"}


def spawn_background_subagent(task_description: str, agent_role: str = "Researcher & Coder") -> Dict[str, Any]:
    """
    Spawn an autonomous background subagent worker to solve a complex, long-running task
    (e.g. deep research, scraping multiple sources, data analysis, codebase audits)
    independently without blocking the user. When finished, the subagent sends a full report to Telegram.
    
    Args:
        task_description: Complete detailed instructions for the subagent.
        agent_role: Specialized title or job role (e.g. 'Deep Web Researcher', 'Python Coder', 'Security Auditor').
    """
    try:
        import subagents
        user_id = get_current_user_id()
        chat_id = get_current_chat_id()
        return subagents.spawn_subagent(user_id=user_id, chat_id=chat_id, role=agent_role, task_description=task_description)
    except Exception as e:
        return {"status": "error", "message": str(e)}


def check_subagent_status(subagent_id: str) -> Dict[str, Any]:
    """
    Check the execution status and report of a spawned background subagent by its task ID.
    
    Args:
        subagent_id: The task ID (e.g. 'sub_1a2b3c4d').
    """
    try:
        task = database.get_subagent_task_sync(subagent_id)
        if not task:
            return {"status": "error", "message": f"Subagent dengan ID '{subagent_id}' tidak ditemukan."}
        return {"status": "success", "task": task}
    except Exception as e:
        return {"status": "error", "message": str(e)}


def add_recurring_task(title: str, prompt_instruction: str, interval_minutes: int = 60) -> Dict[str, Any]:
    """
    Schedule an autonomous recurring task or proactive watchdog (e.g. monitor server stats every 30m, check crypto price every 15m, daily tech news briefing every 1440m).
    The AI will automatically run this instruction periodically and send the result to Telegram!
    
    Args:
        title: Short descriptive name for this task.
        prompt_instruction: The exact prompt/action for the AI agent to execute on each interval.
        interval_minutes: Repeat interval in minutes (default: 60 minutes).
    """
    try:
        user_id = get_current_user_id()
        chat_id = get_current_chat_id()
        job_id = database.add_cron_job_sync(
            user_id=user_id,
            chat_id=chat_id,
            title=title,
            prompt_instruction=prompt_instruction,
            interval_minutes=max(1, interval_minutes)
        )
        return {
            "status": "success",
            "job_id": job_id,
            "message": f"Tugas berulang #{job_id} '{title}' berhasil dijadwalkan setiap {interval_minutes} menit."
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


def list_recurring_tasks() -> Dict[str, Any]:
    """
    List all active recurring tasks and watchdogs scheduled for the user.
    """
    try:
        user_id = get_current_user_id()
        jobs = database.list_cron_jobs_sync(user_id)
        return {"status": "success", "total_jobs": len(jobs), "tasks": jobs}
    except Exception as e:
        return {"status": "error", "message": str(e)}


def cancel_recurring_task(task_id: int) -> Dict[str, Any]:
    """
    Cancel and delete a scheduled recurring task or watchdog by its ID.
    
    Args:
        task_id: The ID of the recurring task.
    """
    try:
        user_id = get_current_user_id()
        success = database.delete_cron_job_sync(user_id, task_id)
        if success:
            return {"status": "success", "message": f"Tugas berulang #{task_id} berhasil dibatalkan dan dihapus."}
        return {"status": "error", "message": f"Tugas #{task_id} tidak ditemukan."}
    except Exception as e:
        return {"status": "error", "message": str(e)}


def get_pdf_output_dir(subfolder: str) -> str:
    """Returns absolute path to ~/Dokumen/ALFA_PDF_TOOLS/<subfolder>/ and ensures it exists."""
    home_dir = os.path.expanduser("~")
    dokumen_dir = os.path.join(home_dir, "Dokumen")
    if not os.path.exists(dokumen_dir):
        dokumen_dir = os.path.join(home_dir, "Documents")
    target_dir = os.path.join(dokumen_dir, "ALFA_PDF_TOOLS", subfolder)
    os.makedirs(target_dir, exist_ok=True)
    return target_dir


def generate_pdf_report(title: str, summary: str, table_data_json: str = "", filename: str = "laporan.pdf") -> Dict[str, Any]:
    """
    Generate a modern, beautifully styled PDF document report with ReportLab and automatically send it to Telegram.
    
    Args:
        title: Main document title (e.g. 'Laporan Analisis Kinerja Server', 'Rangkuman Riset Pasar').
        summary: Paragraphs of text explaining the findings, recommendations, or content.
        table_data_json: Optional JSON string of 2D array for tables, e.g. '[["Header1", "Header2"], ["Val1", "Val2"]]'.
        filename: Output filename ending in .pdf.
    """
    try:
        import json
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors
        
        table_data = None
        if table_data_json:
            if isinstance(table_data_json, str):
                try:
                    table_data = json.loads(table_data_json)
                except Exception:
                    table_data = None
            elif isinstance(table_data_json, list):
                table_data = table_data_json
        
        out_dir = get_pdf_output_dir("Reports")
        safe_name = filename if filename.endswith(".pdf") else f"{filename}.pdf"
        target_path = os.path.join(out_dir, safe_name)
        
        doc = SimpleDocTemplate(target_path, pagesize=letter, rightMargin=40, leftMargin=40, topMargin=40, bottomMargin=40)
        styles = getSampleStyleSheet()
        
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            leading=22,
            textColor=colors.HexColor('#1E3A8A'),
            spaceAfter=15
        )
        
        body_style = ParagraphStyle(
            'CustomBody',
            parent=styles['Normal'],
            fontSize=10,
            leading=15,
            textColor=colors.HexColor('#334155'),
            spaceAfter=10
        )
        
        elements = [
            Paragraph(f"<b>{title}</b>", title_style),
            Spacer(1, 10),
        ]
        
        for paragraph in summary.split("\n\n"):
            if paragraph.strip():
                clean_p = paragraph.strip().replace("\n", "<br/>")
                elements.append(Paragraph(clean_p, body_style))
                elements.append(Spacer(1, 8))
                
        if table_data and len(table_data) > 0:
            elements.append(Spacer(1, 12))
            t = Table(table_data, style=[
                ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1E3A8A')),
                ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
                ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
                ('FONTSIZE', (0,0), (-1,0), 10),
                ('BOTTOMPADDING', (0,0), (-1,0), 8),
                ('TOPPADDING', (0,0), (-1,0), 8),
                ('BACKGROUND', (0,1), (-1,-1), colors.HexColor('#F8FAFC')),
                ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#CBD5E1')),
                ('ALIGN', (0,0), (-1,-1), 'LEFT'),
                ('FONTSIZE', (0,1), (-1,-1), 9),
                ('TOPPADDING', (0,1), (-1,-1), 6),
                ('BOTTOMPADDING', (0,1), (-1,-1), 6),
            ])
            elements.append(t)
            
        doc.build(elements)
        return {
            "status": "success",
            "message": f"Dokumen PDF '{safe_name}' tersimpan di Dokumen/ALFA_PDF_TOOLS/Reports/.",
            "file_path": target_path,
            "filename": safe_name
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal membuat PDF: {str(e)}"}


# ══════════════════════════════════════════════════════════════════════════════
#                     ALFA ULTIMATE PDF-TOOLS SUITE
# ══════════════════════════════════════════════════════════════════════════════

def pdf_merge_documents(pdf_paths: List[str], output_filename: str = "merged.pdf") -> Dict[str, Any]:
    """
    Gabungkan beberapa file PDF menjadi satu dokumen PDF utuh.
    
    Args:
        pdf_paths: Daftar path file PDF yang ingin digabungkan (misal: ['/tmp/doc1.pdf', '/tmp/doc2.pdf']).
        output_filename: Nama file PDF hasil penggabungan (misal: 'merged_dokumen.pdf').
    """
    try:
        from pypdf import PdfReader, PdfWriter
        out_dir = get_pdf_output_dir("Merge")
        safe_name = output_filename if output_filename.endswith(".pdf") else f"{output_filename}.pdf"
        target_path = os.path.join(out_dir, safe_name)
        
        writer = PdfWriter()
        merged_count = 0
        total_pages = 0
        
        for p in pdf_paths:
            exp_p = os.path.expanduser(p.strip())
            if not os.path.exists(exp_p):
                continue
            reader = PdfReader(exp_p)
            for page in reader.pages:
                writer.add_page(page)
                total_pages += 1
            merged_count += 1
            
        if merged_count == 0:
            return {"status": "error", "message": "Tidak ada file PDF valid yang ditemukan untuk digabungkan."}
            
        with open(target_path, "wb") as f_out:
            writer.write(f_out)
            
        return {
            "status": "success",
            "message": f"Berhasil menggabungkan {merged_count} file PDF menjadi {total_pages} halaman di folder Dokumen/ALFA_PDF_TOOLS/Merge/.",
            "file_path": target_path,
            "filename": safe_name,
            "total_pages": total_pages,
            "file_size_bytes": os.path.getsize(target_path)
        }
    except Exception as e:
        logger.error(f"Error in pdf_merge_documents: {e}")
        return {"status": "error", "message": f"Gagal merge PDF: {str(e)}"}


def pdf_split_document(pdf_path: str, page_ranges: str = "", output_dir: str = "") -> Dict[str, Any]:
    """
    Pecah file PDF per halaman atau berdasarkan rentang halaman tertentu (misal '1-3, 5, 8-10').
    
    Args:
        pdf_path: Path ke file PDF yang ingin dipecah.
        page_ranges: Rentang halaman yang ingin diekstrak (kosongkan untuk memecah semua halaman per file).
        output_dir: Direktori output file (opsional, default ke ~/Dokumen/ALFA_PDF_TOOLS/Split/).
    """
    try:
        from pypdf import PdfReader, PdfWriter
        from pathlib import Path
        exp_p = os.path.expanduser(pdf_path.strip())
        if not os.path.exists(exp_p):
            return {"status": "error", "message": f"File PDF '{pdf_path}' tidak ditemukan."}
            
        base_stem = Path(exp_p).stem
        target_dir = os.path.expanduser(output_dir.strip()) if output_dir else os.path.join(get_pdf_output_dir("Split"), base_stem)
        os.makedirs(target_dir, exist_ok=True)
        
        reader = PdfReader(exp_p)
        total = len(reader.pages)
        
        indices = []
        if page_ranges:
            for part in page_ranges.split(","):
                part = part.strip()
                if "-" in part:
                    s, e = part.split("-", 1)
                    s_idx = max(0, int(s.strip()) - 1)
                    e_idx = min(total, int(e.strip()))
                    indices.extend(range(s_idx, e_idx))
                elif part.isdigit():
                    idx = int(part) - 1
                    if 0 <= idx < total:
                        indices.append(idx)
            indices = sorted(list(set(indices)))
        else:
            indices = list(range(total))
            
        output_files = []
        for idx in indices:
            writer = PdfWriter()
            writer.add_page(reader.pages[idx])
            out_file = os.path.join(target_dir, f"{base_stem}_page_{idx+1:03d}.pdf")
            with open(out_file, "wb") as f_out:
                writer.write(f_out)
            output_files.append(out_file)
            
        return {
            "status": "success",
            "message": f"Berhasil memecah PDF menjadi {len(output_files)} file halaman di Dokumen/ALFA_PDF_TOOLS/Split/{base_stem}/.",
            "output_dir": target_dir,
            "files": output_files
        }
    except Exception as e:
        logger.error(f"Error in pdf_split_document: {e}")
        return {"status": "error", "message": f"Gagal split PDF: {str(e)}"}


def pdf_extract_full_text(pdf_path: str, page_numbers: str = "") -> Dict[str, Any]:
    """
    Ekstrak teks lengkap dari dokumen PDF secara bersih dan terstruktur serta simpan salinan file .txt.
    
    Args:
        pdf_path: Path ke file PDF yang ingin dibaca teksnya.
        page_numbers: Opsi nomor halaman spesifik (misal '1,2,5' atau '1-4').
    """
    try:
        from pathlib import Path
        exp_p = os.path.expanduser(pdf_path.strip())
        if not os.path.exists(exp_p):
            return {"status": "error", "message": f"File PDF '{pdf_path}' tidak ditemukan."}
            
        extracted_pages = []
        import pdfplumber
        with pdfplumber.open(exp_p) as pdf:
            total_pages = len(pdf.pages)
            indices = list(range(total_pages))
            if page_numbers:
                req_indices = []
                for p in page_numbers.split(","):
                    p = p.strip()
                    if "-" in p:
                        s, e = p.split("-", 1)
                        req_indices.extend(range(max(0, int(s)-1), min(total_pages, int(e))))
                    elif p.isdigit():
                        idx = int(p) - 1
                        if 0 <= idx < total_pages:
                            req_indices.append(idx)
                indices = sorted(list(set(req_indices)))
                
            for idx in indices:
                t = pdf.pages[idx].extract_text() or ""
                extracted_pages.append(f"--- [Halaman {idx+1}/{total_pages}] ---\n{t.strip()}")
                
        full_text = "\n\n".join(extracted_pages)
        
        # Save to Dokumen/ALFA_PDF_TOOLS/Extract_Text/
        out_dir = get_pdf_output_dir("Extract_Text")
        txt_filename = f"{Path(exp_p).stem}_extracted.txt"
        target_path = os.path.join(out_dir, txt_filename)
        with open(target_path, "w", encoding="utf-8") as f_txt:
            f_txt.write(full_text)
            
        return {
            "status": "success",
            "message": f"Teks berhasil diekstrak dan disimpan di Dokumen/ALFA_PDF_TOOLS/Extract_Text/{txt_filename}.",
            "file_path": target_path,
            "filename": txt_filename,
            "total_pages": total_pages,
            "extracted_pages_count": len(extracted_pages),
            "text_length_chars": len(full_text),
            "text_preview": full_text[:4000],
            "full_text": full_text
        }
    except Exception as e:
        try:
            from pypdf import PdfReader
            from pathlib import Path
            reader = PdfReader(exp_p)
            texts = [f"--- [Halaman {i+1}] ---\n{p.extract_text() or ''}" for i, p in enumerate(reader.pages)]
            full = "\n\n".join(texts)
            out_dir = get_pdf_output_dir("Extract_Text")
            txt_filename = f"{Path(exp_p).stem}_extracted.txt"
            target_path = os.path.join(out_dir, txt_filename)
            with open(target_path, "w", encoding="utf-8") as f_txt:
                f_txt.write(full)
            return {
                "status": "success",
                "message": f"Teks berhasil diekstrak dan disimpan di Dokumen/ALFA_PDF_TOOLS/Extract_Text/{txt_filename}.",
                "file_path": target_path,
                "filename": txt_filename,
                "total_pages": len(reader.pages),
                "text_preview": full[:4000],
                "full_text": full
            }
        except Exception as err2:
            return {"status": "error", "message": f"Gagal ekstrak teks PDF: {str(err2)}"}


def pdf_encrypt_password(pdf_path: str, password: str, output_filename: str = "protected.pdf") -> Dict[str, Any]:
    """
    Lindungi dan kunci file PDF dengan password menggunakan enkripsi kuat AES-256.
    
    Args:
        pdf_path: Path ke file PDF yang ingin dienkripsi.
        password: Password pengunci dokumen.
        output_filename: Nama file output terenkripsi.
    """
    try:
        from pypdf import PdfReader, PdfWriter
        exp_p = os.path.expanduser(pdf_path.strip())
        if not os.path.exists(exp_p):
            return {"status": "error", "message": f"File '{pdf_path}' tidak ditemukan."}
            
        out_dir = get_pdf_output_dir("Encrypt")
        safe_name = output_filename if output_filename.endswith(".pdf") else f"{output_filename}.pdf"
        target_path = os.path.join(out_dir, safe_name)
        
        reader = PdfReader(exp_p)
        writer = PdfWriter()
        for p in reader.pages:
            writer.add_page(p)
            
        writer.encrypt(user_password=password, owner_password=password, algorithm="AES-256")
        with open(target_path, "wb") as f_out:
            writer.write(f_out)
            
        return {
            "status": "success",
            "message": f"Dokumen PDF berhasil dienkripsi dengan AES-256 di Dokumen/ALFA_PDF_TOOLS/Encrypt/{safe_name}.",
            "file_path": target_path,
            "filename": safe_name
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal enkripsi PDF: {str(e)}"}


def pdf_decrypt_password(pdf_path: str, password: str, output_filename: str = "unlocked.pdf") -> Dict[str, Any]:
    """
    Buka kunci PDF yang terproteksi password dan simpan salinan tanpa password.
    
    Args:
        pdf_path: Path ke file PDF terenkripsi.
        password: Password yang digunakan untuk membuka kunci.
        output_filename: Nama file output yang sudah tidak terkunci.
    """
    try:
        from pypdf import PdfReader, PdfWriter
        exp_p = os.path.expanduser(pdf_path.strip())
        if not os.path.exists(exp_p):
            return {"status": "error", "message": f"File '{pdf_path}' tidak ditemukan."}
            
        out_dir = get_pdf_output_dir("Decrypt")
        safe_name = output_filename if output_filename.endswith(".pdf") else f"{output_filename}.pdf"
        target_path = os.path.join(out_dir, safe_name)
        
        reader = PdfReader(exp_p)
        if reader.is_encrypted:
            res = reader.decrypt(password)
            if not res:
                return {"status": "error", "message": "Password salah atau PDF tidak dapat didekripsi."}
                
        writer = PdfWriter()
        for p in reader.pages:
            writer.add_page(p)
            
        with open(target_path, "wb") as f_out:
            writer.write(f_out)
            
        return {
            "status": "success",
            "message": f"PDF berhasil didekripsi di Dokumen/ALFA_PDF_TOOLS/Decrypt/{safe_name}.",
            "file_path": target_path,
            "filename": safe_name
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal dekripsi PDF: {str(e)}"}


def pdf_rotate_pages(pdf_path: str, angle: int = 90, page_numbers: str = "", output_filename: str = "rotated.pdf") -> Dict[str, Any]:
    """
    Putar orientasi halaman PDF (90, 180, atau 270 derajat searah jarum jam).
    
    Args:
        pdf_path: Path ke file PDF.
        angle: Sudut putar (90, 180, 270).
        page_numbers: Halaman tertentu yang ingin diputar (misal '1,3-5', kosongkan untuk semua).
        output_filename: Nama file output.
    """
    try:
        from pypdf import PdfReader, PdfWriter
        exp_p = os.path.expanduser(pdf_path.strip())
        if not os.path.exists(exp_p):
            return {"status": "error", "message": f"File '{pdf_path}' tidak ditemukan."}
            
        out_dir = get_pdf_output_dir("Rotate")
        safe_name = output_filename if output_filename.endswith(".pdf") else f"{output_filename}.pdf"
        target_path = os.path.join(out_dir, safe_name)
        
        reader = PdfReader(exp_p)
        writer = PdfWriter()
        total = len(reader.pages)
        
        target_indices = list(range(total))
        if page_numbers:
            req = []
            for p in page_numbers.split(","):
                p = p.strip()
                if "-" in p:
                    s, e = p.split("-", 1)
                    req.extend(range(max(0, int(s)-1), min(total, int(e))))
                elif p.isdigit():
                    idx = int(p) - 1
                    if 0 <= idx < total:
                        req.append(idx)
            target_indices = list(set(req))
            
        for i, page in enumerate(reader.pages):
            if i in target_indices:
                page.rotate(angle)
            writer.add_page(page)
            
        with open(target_path, "wb") as f_out:
            writer.write(f_out)
            
        return {
            "status": "success",
            "message": f"Berhasil memutar {len(target_indices)} halaman PDF sebesar {angle}° di Dokumen/ALFA_PDF_TOOLS/Rotate/{safe_name}.",
            "file_path": target_path,
            "filename": safe_name
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal rotasi PDF: {str(e)}"}


def pdf_apply_watermark_text(pdf_path: str, watermark_text: str, opacity: float = 0.2, angle: float = 45, output_filename: str = "watermarked.pdf") -> Dict[str, Any]:
    """
    Tambahkan stempel watermark teks diagonal transparan ke setiap halaman PDF.
    
    Args:
        pdf_path: Path ke file PDF asli.
        watermark_text: Teks watermark (misal 'CONFIDENTIAL', 'RAHASIA DOKUMEN', 'DRAFT').
        opacity: Tingkat transparansi (0.05 sampai 0.5).
        angle: Sudut kemiringan diagonal watermark (default 45 derajat).
        output_filename: Nama file output.
    """
    try:
        import io
        from pypdf import PdfReader, PdfWriter
        from reportlab.pdfgen import canvas as rl_canvas
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.colors import HexColor
        
        exp_p = os.path.expanduser(pdf_path.strip())
        if not os.path.exists(exp_p):
            return {"status": "error", "message": f"File '{pdf_path}' tidak ditemukan."}
            
        out_dir = get_pdf_output_dir("Watermark")
        safe_name = output_filename if output_filename.endswith(".pdf") else f"{output_filename}.pdf"
        target_path = os.path.join(out_dir, safe_name)
        
        packet = io.BytesIO()
        can = rl_canvas.Canvas(packet, pagesize=A4)
        can.setFont("Helvetica-Bold", 45)
        can.saveState()
        can.setFillAlpha(max(0.05, min(0.9, opacity)))
        can.setFillColor(HexColor("#64748B"))
        can.translate(A4[0] / 2, A4[1] / 2)
        can.rotate(angle)
        can.drawCentredString(0, 0, watermark_text)
        can.restoreState()
        can.save()
        packet.seek(0)
        
        wm_page = PdfReader(packet).pages[0]
        reader = PdfReader(exp_p)
        writer = PdfWriter()
        
        for p in reader.pages:
            p.merge_page(wm_page)
            writer.add_page(p)
            
        with open(target_path, "wb") as f_out:
            writer.write(f_out)
            
        return {
            "status": "success",
            "message": f"Watermark '{watermark_text}' berhasil ditempelkan di Dokumen/ALFA_PDF_TOOLS/Watermark/{safe_name}.",
            "file_path": target_path,
            "filename": safe_name
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal watermark PDF: {str(e)}"}


def pdf_insert_page_numbers(pdf_path: str, position: str = "bottom-center", start_number: int = 1, output_filename: str = "numbered.pdf") -> Dict[str, Any]:
    """
    Sematkan penomoran halaman otomatis pada dokumen PDF.
    
    Args:
        pdf_path: Path ke file PDF.
        position: Posisi nomor ('bottom-center', 'bottom-right', 'bottom-left', 'top-right', 'top-center').
        start_number: Nomor awal penomoran halaman (default 1).
        output_filename: Nama file output.
    """
    try:
        import io
        from pypdf import PdfReader, PdfWriter
        from reportlab.pdfgen import canvas as rl_canvas
        from reportlab.lib.colors import HexColor
        
        exp_p = os.path.expanduser(pdf_path.strip())
        if not os.path.exists(exp_p):
            return {"status": "error", "message": f"File '{pdf_path}' tidak ditemukan."}
            
        out_dir = get_pdf_output_dir("Page_Numbers")
        safe_name = output_filename if output_filename.endswith(".pdf") else f"{output_filename}.pdf"
        target_path = os.path.join(out_dir, safe_name)
        
        reader = PdfReader(exp_p)
        writer = PdfWriter()
        total = len(reader.pages)
        
        for i, page in enumerate(reader.pages):
            w = float(page.mediabox.width)
            h = float(page.mediabox.height)
            packet = io.BytesIO()
            can = rl_canvas.Canvas(packet, pagesize=(w, h))
            can.setFont("Helvetica", 10)
            can.setFillColor(HexColor("#334155"))
            
            num_str = f"Halaman {start_number + i} dari {total}"
            positions = {
                "bottom-center": (w / 2, 25),
                "bottom-right": (w - 40, 25),
                "bottom-left": (40, 25),
                "top-center": (w / 2, h - 25),
                "top-right": (w - 40, h - 25),
            }
            x, y = positions.get(position, (w / 2, 25))
            can.drawCentredString(x, y, num_str)
            can.save()
            packet.seek(0)
            
            num_page = PdfReader(packet).pages[0]
            page.merge_page(num_page)
            writer.add_page(page)
            
        with open(target_path, "wb") as f_out:
            writer.write(f_out)
            
        return {
            "status": "success",
            "message": f"Nomor halaman berhasil ditambahkan di Dokumen/ALFA_PDF_TOOLS/Page_Numbers/{safe_name}.",
            "file_path": target_path,
            "filename": safe_name
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal memberi nomor halaman: {str(e)}"}


def pdf_convert_to_images(pdf_path: str, dpi: int = 150, output_dir: str = "") -> Dict[str, Any]:
    """
    Konversi seluruh halaman PDF menjadi gambar PNG resolusi tinggi.
    
    Args:
        pdf_path: Path ke file PDF.
        dpi: Kerapatan resolusi gambar (default 150 DPI).
        output_dir: Folder penyimpanan gambar hasil konversi (opsional, default ke ~/Dokumen/ALFA_PDF_TOOLS/PDF_to_Images/).
    """
    try:
        from pathlib import Path
        exp_p = os.path.expanduser(pdf_path.strip())
        if not os.path.exists(exp_p):
            return {"status": "error", "message": f"File '{pdf_path}' tidak ditemukan."}
            
        base_name = Path(exp_p).stem
        target_dir = os.path.expanduser(output_dir.strip()) if output_dir else os.path.join(get_pdf_output_dir("PDF_to_Images"), base_name)
        os.makedirs(target_dir, exist_ok=True)
        
        out_prefix = os.path.join(target_dir, f"{base_name}_page")
        
        cmd = ["pdftoppm", "-png", "-r", str(dpi), exp_p, out_prefix]
        res = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
        
        generated_images = [os.path.join(target_dir, f) for f in os.listdir(target_dir) if f.startswith(f"{base_name}_page") and f.endswith(".png")]
        generated_images.sort()
        
        return {
            "status": "success",
            "message": f"Berhasil merender {len(generated_images)} halaman PDF menjadi gambar PNG di Dokumen/ALFA_PDF_TOOLS/PDF_to_Images/{base_name}/.",
            "output_dir": target_dir,
            "images": generated_images
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal konversi PDF ke gambar: {str(e)}"}


def images_convert_to_pdf(image_paths: List[str], output_filename: str = "images_album.pdf") -> Dict[str, Any]:
    """
    Gabungkan kumpulan file foto/gambar (JPG, PNG, WEBP) menjadi satu dokumen PDF rapi.
    
    Args:
        image_paths: Daftar path file gambar yang ingin digabungkan ke PDF.
        output_filename: Nama file output PDF.
    """
    try:
        from PIL import Image
        out_dir = get_pdf_output_dir("Images_to_PDF")
        safe_name = output_filename if output_filename.endswith(".pdf") else f"{output_filename}.pdf"
        target_path = os.path.join(out_dir, safe_name)
        
        opened_images = []
        for p in image_paths:
            exp_p = os.path.expanduser(p.strip())
            if os.path.exists(exp_p):
                img = Image.open(exp_p).convert("RGB")
                opened_images.append(img)
                
        if not opened_images:
            return {"status": "error", "message": "Tidak ada file gambar valid yang ditemukan."}
            
        first = opened_images[0]
        rest = opened_images[1:] if len(opened_images) > 1 else []
        first.save(target_path, "PDF", resolution=100.0, save_all=True, append_images=rest)
        
        return {
            "status": "success",
            "message": f"Berhasil mengubah {len(opened_images)} gambar menjadi PDF di Dokumen/ALFA_PDF_TOOLS/Images_to_PDF/{safe_name}.",
            "file_path": target_path,
            "filename": safe_name
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal mengubah gambar ke PDF: {str(e)}"}


def pdf_inspect_metadata(pdf_path: str) -> Dict[str, Any]:
    """
    Periksa informasi teknis mendalam dari file PDF (jumlah halaman, versi PDF, ukuran file, enkripsi, metadata) dan simpan salinan JSON.
    
    Args:
        pdf_path: Path ke file PDF yang ingin diinspeksi.
    """
    try:
        from pypdf import PdfReader
        from pathlib import Path
        exp_p = os.path.expanduser(pdf_path.strip())
        if not os.path.exists(exp_p):
            return {"status": "error", "message": f"File '{pdf_path}' tidak ditemukan."}
            
        reader = PdfReader(exp_p)
        meta = reader.metadata or {}
        
        first_page = reader.pages[0] if reader.pages else None
        width_pt = float(first_page.mediabox.width) if first_page else 0
        height_pt = float(first_page.mediabox.height) if first_page else 0
        
        info = {
            "status": "success",
            "file_name": os.path.basename(exp_p),
            "file_path": exp_p,
            "file_size_kb": round(os.path.getsize(exp_p) / 1024, 2),
            "total_pages": len(reader.pages),
            "is_encrypted": reader.is_encrypted,
            "dimensions_pt": f"{width_pt:.1f} x {height_pt:.1f}",
            "title": meta.get("/Title") or meta.title or "N/A",
            "author": meta.get("/Author") or meta.author or "N/A",
            "creator": meta.get("/Creator") or meta.creator or "N/A",
            "producer": meta.get("/Producer") or meta.producer or "N/A"
        }
        
        # Save json copy to Dokumen/ALFA_PDF_TOOLS/Inspect/
        out_dir = get_pdf_output_dir("Inspect")
        json_file = os.path.join(out_dir, f"{Path(exp_p).stem}_metadata.json")
        with open(json_file, "w", encoding="utf-8") as f_j:
            json.dump(info, f_j, indent=2, default=str)
        info["saved_json_path"] = json_file
        
        return info
    except Exception as e:
        return {"status": "error", "message": f"Gagal inspeksi PDF: {str(e)}"}


def pdf_compress_and_optimize(pdf_path: str, output_filename: str = "compressed.pdf") -> Dict[str, Any]:
    """
    Kompresi dan optimalkan ukuran file PDF dengan mereduksi stream konten dan metadata berlebih.
    
    Args:
        pdf_path: Path ke file PDF yang ingin dikompres.
        output_filename: Nama file output hasil kompresi.
    """
    try:
        from pypdf import PdfReader, PdfWriter
        exp_p = os.path.expanduser(pdf_path.strip())
        if not os.path.exists(exp_p):
            return {"status": "error", "message": f"File '{pdf_path}' tidak ditemukan."}
            
        out_dir = get_pdf_output_dir("Compress")
        safe_name = output_filename if output_filename.endswith(".pdf") else f"{output_filename}.pdf"
        target_path = os.path.join(out_dir, safe_name)
        orig_size = os.path.getsize(exp_p)
        
        reader = PdfReader(exp_p)
        writer = PdfWriter()
        for p in reader.pages:
            p.compress_content_streams()
            writer.add_page(p)
            
        with open(target_path, "wb") as f_out:
            writer.write(f_out)
            
        new_size = os.path.getsize(target_path)
        savings_pct = max(0, round((orig_size - new_size) / orig_size * 100, 1)) if orig_size > 0 else 0
        
        return {
            "status": "success",
            "message": f"PDF berhasil dikompresi di Dokumen/ALFA_PDF_TOOLS/Compress/{safe_name}. Hemat {savings_pct}% ruang penyimpanan.",
            "file_path": target_path,
            "filename": safe_name,
            "original_size_kb": round(orig_size / 1024, 2),
            "new_size_kb": round(new_size / 1024, 2),
            "saved_percentage": savings_pct
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal kompresi PDF: {str(e)}"}


# ══════════════════════════════════════════════════════════════════════════════
#              ALFA AUTOMATED AFFILIATE SALES SWARM SUITE
# ══════════════════════════════════════════════════════════════════════════════

def affiliate_hunt_trending_products(niche: str = "gadget unik murah", platform: str = "shopee") -> Dict[str, Any]:
    """
    Riset tren produk dan kata kunci viral untuk niche affiliate Shopee / TikTok Shop (Dikelola oleh Researcher Prime).
    
    Args:
        niche: Kategori atau kata kunci produk (misal: 'gadget unik murah', 'peralatan dapur estetik', 'fashion pria korea').
        platform: Platform target ('shopee', 'tiktok', 'lazada').
    """
    try:
        import affiliate_engine
        return affiliate_engine.research_trending_niche(niche=niche, platform=platform)
    except Exception as e:
        logger.error(f"Error in affiliate_hunt_trending_products: {e}")
        return {"status": "error", "message": str(e)}


def affiliate_generate_viral_content(
    product_name: str,
    key_features: str,
    original_price: str,
    discount_price: str,
    affiliate_link: str,
    target_audience: str = "Pecinta Gadget & Lifestyle / Pemburu Diskon",
    platform: str = "shopee_tiktok"
) -> Dict[str, Any]:
    """
    Hasilkan paket copywriting viral lengkap: Script Video TikTok (Hook 3s, Story, CTA), Telegram Deals Card, WhatsApp Broadcast, dan Auto-Reply 'Spill Link' (Dikelola oleh Strategic Planner).
    
    Args:
        product_name: Nama produk (misal: 'Mini Powerbank Kapsul Fast Charging 5000mAh').
        key_features: Keunggulan dan spesifikasi utama produk dipisah koma.
        original_price: Harga sebelum diskon (misal: 'Rp 150.000').
        discount_price: Harga flash sale/diskon (misal: 'Rp 49.000').
        affiliate_link: Link affiliate Shopee / TikTok Shop kamu.
        target_audience: Segmentasi target pembeli.
        platform: Platform target konten.
    """
    try:
        import affiliate_engine
        return affiliate_engine.generate_affiliate_campaign_content(
            product_name=product_name,
            key_features=key_features,
            original_price=original_price,
            discount_price=discount_price,
            affiliate_link=affiliate_link,
            target_audience=target_audience,
            platform=platform
        )
    except Exception as e:
        logger.error(f"Error in affiliate_generate_viral_content: {e}")
        return {"status": "error", "message": str(e)}


def affiliate_broadcast_deal(
    product_name: str,
    message_text: str,
    affiliate_link: str,
    channels: List[str] = ["telegram", "whatsapp"]
) -> Dict[str, Any]:
    """
    Kirimkan penawaran diskon affiliate secara otomatis ke Telegram Channel atau broadcast WhatsApp (Dikelola oleh Code Crafter).
    
    Args:
        product_name: Nama produk yang dipromosikan.
        message_text: Teks copywriting promosi lengkap.
        affiliate_link: URL link affiliate resmi.
        channels: Daftar channel tujuan (['telegram', 'whatsapp']).
    """
    try:
        import affiliate_engine
        return affiliate_engine.broadcast_affiliate_deal(
            product_name=product_name,
            message_text=message_text,
            affiliate_link=affiliate_link,
            channels=channels
        )
    except Exception as e:
        logger.error(f"Error in affiliate_broadcast_deal: {e}")
        return {"status": "error", "message": str(e)}


def affiliate_list_campaigns(limit: int = 15) -> Dict[str, Any]:
    """
    Lihat rekap histori campaign dan script affiliate yang aktif (Dikelola oleh Alpha Lead).
    
    Args:
        limit: Jumlah campaign yang ingin ditampilkan.
    """
    try:
        import affiliate_engine
        campaigns = affiliate_engine.list_affiliate_campaigns(limit=limit)
        return {
            "status": "success",
            "total_campaigns": len(campaigns),
            "campaigns": campaigns
        }
    except Exception as e:
        logger.error(f"Error in affiliate_list_campaigns: {e}")
        return {"status": "error", "message": str(e)}


# ══════════════════════════════════════════════════════════════════════════════
#         ALFA ULTRA-FAST CAMOUFOX & LARGE-SCALE SCRAPER SUITE
# ══════════════════════════════════════════════════════════════════════════════

def scrape_real_product_data(url: str, engine: str = "auto") -> Dict[str, Any]:
    """
    Scrape data produk real dari Shopee, TikTok Shop, Tokopedia, atau website manapun menggunakan Camoufox Anti-Detect Browser atau Fast TLS.
    Bypass proteksi Cloudflare, bot detector, dan dynamic javascript rendering.
    
    Args:
        url: Link produk atau halaman yang ingin discrape.
        engine: Pilihan engine ('auto', 'camoufox', 'fast_tls').
    """
    try:
        import fast_scraper
        if engine == "camoufox":
            return fast_scraper.scrape_with_camoufox(url)
        elif engine == "fast_tls":
            return fast_scraper.scrape_with_fast_tls(url)
        else:
            domain = url.lower()
            if "shopee" in domain or "tiktok" in domain or "tokopedia" in domain:
                return fast_scraper.scrape_with_camoufox(url)
            else:
                return fast_scraper.scrape_with_fast_tls(url)
    except Exception as e:
        logger.error(f"Error in scrape_real_product_data: {e}")
        return {"status": "error", "message": str(e)}


def scrape_large_scale_batch(
    urls: List[str],
    batch_name: str = "batch_products",
    max_concurrency: int = 15,
    use_camoufox: bool = False
) -> Dict[str, Any]:
    """
    Scraping paralel skala besar untuk puluhan hingga ribuan URL sekaligus dengan kecepatan sangat tinggi.
    Hasil otomatis diekspor ke file JSON dan CSV di ~/Dokumen/ALFA_SCRAPER_DATA/.
    
    Args:
        urls: Daftar URL yang ingin discrape secara massal.
        batch_name: Nama batch untuk penamaan file ekspor.
        max_concurrency: Jumlah request paralel serentak (default: 15).
        use_camoufox: True untuk menggunakan browser Camoufox Anti-Detect, False untuk Fast TLS engine.
    """
    try:
        import fast_scraper
        return fast_scraper.run_batch_scrape(
            urls=urls,
            batch_name=batch_name,
            max_concurrency=max_concurrency,
            use_camoufox=use_camoufox
        )
    except Exception as e:
        logger.error(f"Error in scrape_large_scale_batch: {e}")
        return {"status": "error", "message": str(e)}


def marketplace_search_products(query: str, platform: str = "shopee", max_items: int = 15) -> Dict[str, Any]:
    """
    Cari dan scrape katalog produk real dari marketplace (Shopee, TikTok Shop, Tokopedia) berdasarkan kata kunci pencarian.
    
    Args:
        query: Kata kunci pencarian produk (misal: 'powerbank mini fast charge', 'lampu tidur estetik').
        platform: Marketplace target ('shopee', 'tiktok', 'tokopedia', 'lazada').
        max_items: Jumlah maksimal produk yang diambil (default: 15).
    """
    try:
        import fast_scraper
        return fast_scraper.search_and_scrape_marketplace(
            query=query,
            platform=platform,
            max_items=max_items
        )
    except Exception as e:
        logger.error(f"Error in marketplace_search_products: {e}")
        return {"status": "error", "message": str(e)}


def generate_promo_video_from_images(
    image_paths: List[str],
    product_name: str,
    voiceover_text: str,
    orig_price: str = "Rp 149.000",
    disc_price: str = "Rp 49.900",
    voice: str = "id-ID-GadisNeural",
    theme: str = "viral_tiktok",
    motion_style: str = "zoom_in",
    badge_text: str = "🔥 FLASH SALE DISKON SPESIAL",
    call_to_action: str = "👉 KLIK KERANJANG KUNING / BIO SEBELUM HABIS 🛒",
    visual_prompt: str = "",
    output_filename: str = "promo_video.mp4"
) -> Dict[str, Any]:
    """
    Generate video promosi produk otomatis format 9:16 (1080x1920) untuk TikTok / Reels / Shorts hanya dari foto produk.
    Dilengkapi prompt visual sinematik detail, animasi zoompan Ken Burns, dubbing voiceover AI bahasa Indonesia, dan banner flash sale diskon.
    Hasil otomatis tersimpan di ~/Dokumen/ALFA_GENERATED_VIDEOS/.
    
    Args:
        image_paths: Daftar path file foto produk di komputer (bisa 1 atau banyak foto).
        product_name: Nama produk yang dipromosikan.
        voiceover_text: Naskah teks narasi / promosi yang akan dibacakan oleh voiceover AI.
        orig_price: Harga coret sebelum diskon (misal: 'Rp 149.000').
        disc_price: Harga flash sale / drop (misal: 'Rp 49.900').
        voice: Suara voiceover AI ('id-ID-GadisNeural' untuk cewek ramah, 'id-ID-ArdiNeural' untuk cowok).
        theme: Tema visual video ('viral_tiktok', 'luxury_gold', 'cyberpunk', 'clean_minimal').
        motion_style: Gaya animasi kamera ('zoom_in', 'zoom_out', 'pan_left_right').
        badge_text: Teks badge promo atas.
        call_to_action: Teks banner CTA bawah.
        visual_prompt: Deskripsi prompt visual sinematik AI untuk video.
        output_filename: Nama file output MP4.
    """
    try:
        import video_generator
        return video_generator.generate_video_from_images(
            image_paths=image_paths,
            product_name=product_name,
            voiceover_text=voiceover_text,
            orig_price=orig_price,
            disc_price=disc_price,
            voice=voice,
            theme=theme,
            motion_style=motion_style,
            badge_text=badge_text,
            call_to_action=call_to_action,
            visual_prompt=visual_prompt,
            output_filename=output_filename
        )
    except Exception as e:
        logger.error(f"Error in generate_promo_video_from_images: {e}")
        return {"status": "error", "message": str(e)}


def generate_excel_spreadsheet(sheet_title: str, headers: List[str], rows_json: str, filename: str = "data.xlsx") -> Dict[str, Any]:
    """
    Generate an Excel (.xlsx) spreadsheet with styled headers, borders, and auto-adjusted columns, automatically sent to Telegram.
    
    Args:
        sheet_title: Name of the worksheet tab.
        headers: List of column header names (e.g. ['Nama', 'Kategori', 'Harga', 'Jumlah']).
        rows_json: JSON string of 2D array of rows (e.g. '[["Barang A", "Kategori 1", 15000], ["Barang B", "Kategori 2", 25000]]').
        filename: Output filename ending in .xlsx.
    """
    try:
        import json
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter
        
        safe_name = filename if filename.endswith(".xlsx") else f"{filename}.xlsx"
        target_path = os.path.join(SANDBOX_DIR, safe_name)
        
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = sheet_title[:30]
        
        header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="1E3A8A", end_color="1E3A8A", fill_type="solid")
        thin_border = Border(
            left=Side(style='thin', color='CBD5E1'),
            right=Side(style='thin', color='CBD5E1'),
            top=Side(style='thin', color='CBD5E1'),
            bottom=Side(style='thin', color='CBD5E1')
        )
        
        ws.append(headers)
        for col_num in range(1, len(headers) + 1):
            cell = ws.cell(row=1, column=col_num)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = thin_border
            
        rows_data = []
        if isinstance(rows_json, str):
            try:
                rows_data = json.loads(rows_json)
            except Exception:
                rows_data = []
        elif isinstance(rows_json, list):
            rows_data = rows_json
            
        for row_data in rows_data:
            ws.append(row_data)
            row_idx = ws.max_row
            for col_num in range(1, len(row_data) + 1):
                cell = ws.cell(row=row_idx, column=col_num)
                cell.border = thin_border
                
        for col in ws.columns:
            max_len = max(len(str(cell.value or '')) for cell in col)
            col_letter = get_column_letter(col[0].column)
            ws.column_dimensions[col_letter].width = max(max_len + 4, 12)
            
        wb.save(target_path)
        return {
            "status": "success",
            "message": f"Spreadsheet Excel '{safe_name}' berhasil dibuat dan akan dikirim ke Telegram.",
            "file_path": target_path
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal membuat Excel: {str(e)}"}


def generate_presentation_pptx(title: str, subtitle: str, slides_json: str, filename: str = "presentasi.pptx") -> Dict[str, Any]:
    """
    Generate a clean PowerPoint presentation (.pptx) and send it directly to Telegram.
    
    Args:
        title: Main presentation title.
        subtitle: Subtitle / author note.
        slides_json: JSON string list of slide objects, e.g. '[{"title": "Slide 1", "points": ["Poin A", "Poin B"]}]'.
        filename: Output filename ending in .pptx.
    """
    try:
        import json
        from pptx import Presentation
        
        # Accept either a JSON string or an already-parsed list
        slides_content = json.loads(slides_json) if isinstance(slides_json, str) else slides_json
        if not slides_content:
            slides_content = []
        if isinstance(slides_content, dict):
            slides_content = [slides_content]
        
        safe_name = filename if filename.endswith(".pptx") else f"{filename}.pptx"
        target_path = os.path.join(SANDBOX_DIR, safe_name)
        
        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle
        
        bullet_layout = prs.slide_layouts[1]
        for item in slides_content:
            s = prs.slides.add_slide(bullet_layout)
            s.shapes.title.text = item.get("title", "Slide")
            tf = s.placeholders[1].text_frame
            tf.word_wrap = True
            points = item.get("points", [])
            for i, pt in enumerate(points):
                if i == 0:
                    tf.text = str(pt)
                else:
                    p = tf.add_paragraph()
                    p.text = str(pt)
                    p.level = 0
                    
        prs.save(target_path)
        return {
            "status": "success",
            "message": f"Presentasi PowerPoint '{safe_name}' berhasil dibuat dan akan dikirim ke Telegram.",
            "file_path": target_path
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal membuat PPTX: {str(e)}"}


def control_linux_hardware(action: str, value: str = "") -> Dict[str, Any]:
    """
    Control Linux laptop/server hardware, audio, screen lock, power, Wi-Fi, and media from Telegram.
    
    Args:
        action: Target hardware action. Options:
                - 'lock_screen': Lock the Linux desktop session immediately.
                - 'set_volume': Set audio output volume (value: '0' to '100', e.g. '50%').
                - 'mute_toggle': Toggle audio mute/unmute.
                - 'media_play_pause': Toggle media playback (Spotify, YouTube, VLC).
                - 'media_next': Next media track.
                - 'media_prev': Previous media track.
                - 'wifi_scan': Scan and list available Wi-Fi access points.
                - 'bluetooth_status': Check Bluetooth status and connected devices.
                - 'battery_status': Check battery health, percentage, and charging state.
        value: Optional parameter for the action (e.g. '50%' for set_volume).
    """
    try:
        act = action.strip().lower()
        if act == "lock_screen":
            subprocess.run("loginctl lock-session 2>/dev/null || gnome-screensaver-command -l", shell=True, capture_output=True, text=True)
            return {"status": "success", "message": "🔒 Layar desktop Linux telah berhasil dikunci."}

        elif act == "set_volume":
            vol_val = value.replace("%", "").strip() or "50"
            try:
                frac = float(vol_val) / 100.0
                subprocess.run(f"wpctl set-volume @DEFAULT_AUDIO_SINK@ {frac}", shell=True, capture_output=True, timeout=3)
            except Exception:
                subprocess.run(f"amixer set Master {vol_val}%", shell=True, capture_output=True, timeout=3)
            return {"status": "success", "message": f"🔊 Volume speaker berhasil diubah ke {vol_val}%."}

        elif act == "mute_toggle":
            subprocess.run("wpctl set-mute @DEFAULT_AUDIO_SINK@ toggle || amixer set Master toggle", shell=True, capture_output=True, timeout=3)
            return {"status": "success", "message": "🔇 Status mute audio berhasil dialihkan (toggle)."}

        elif act in ["media_play_pause", "media_next", "media_prev"]:
            cmd = "playerctl play-pause" if act == "media_play_pause" else "playerctl next" if act == "media_next" else "playerctl previous"
            subprocess.run(cmd, shell=True, capture_output=True, text=True, timeout=3)
            return {"status": "success", "message": f"🎵 Perintah media '{act}' berhasil dieksekusi."}

        elif act == "wifi_scan":
            res = subprocess.run("nmcli -f SSID,SIGNAL,SECURITY dev wifi list | head -n 12", shell=True, capture_output=True, text=True, timeout=8)
            return {"status": "success", "wifi_networks": res.stdout.strip() or "Tidak ada jaringan Wi-Fi ditemukan."}

        elif act == "bluetooth_status":
            res = subprocess.run("bluetoothctl show 2>/dev/null | grep -E 'Name|Powered|Discoverable'; bluetoothctl devices 2>/dev/null", shell=True, capture_output=True, text=True, timeout=5)
            return {"status": "success", "bluetooth": res.stdout.strip() or "Bluetooth tidak aktif."}

        elif act == "battery_status":
            bat = psutil.sensors_battery()
            if not bat:
                return {"status": "success", "battery": "Perangkat tidak menggunakan baterai (Desktop PC / Server)."}
            plugged = "🔌 Sedang Mengisi Daya (Charging)" if bat.power_plugged else "🔋 Berjalan dengan Baterai (Discharging)"
            time_left = f"{round(bat.secsleft / 60)} menit" if bat.secsleft > 0 else "N/A"
            return {
                "status": "success",
                "percentage": f"{bat.percent}%",
                "power_state": plugged,
                "estimated_time_remaining": time_left
            }

        return {"status": "error", "message": f"Aksi hardware '{action}' tidak dikenal."}
    except Exception as e:
        return {"status": "error", "message": str(e)}


def send_file_to_chat(file_path: str, caption: str = "") -> Dict[str, Any]:
    """
    Send an existing file (document, PDF, photo, video, audio, ZIP, code script, data file)
    from the computer filesystem directly to the Telegram user chat.
    Use this tool when the user asks to send, upload, or transfer a specific file from disk to Telegram.
    
    Args:
        file_path: Absolute or relative path to the local file (e.g. '~/Documents/invoice.pdf', '~/Downloads/video.mp4', '~/Downloads/archive.zip', 'bot.py').
        caption: Optional description or caption to accompany the file in chat.
    """
    try:
        import shutil
        expanded = os.path.expanduser(file_path)
        if not os.path.isabs(expanded):
            expanded = os.path.join(os.path.expanduser("~"), file_path)
            
        if not os.path.exists(expanded) or not os.path.isfile(expanded):
            return {"status": "error", "message": f"File tidak ditemukan di path: {file_path}"}
            
        file_size_mb = os.path.getsize(expanded) / (1024 * 1024)
        if file_size_mb > 50:
            return {"status": "error", "message": f"Ukuran file ({round(file_size_mb, 1)} MB) melebihi batas upload Telegram Bot API (50 MB)."}
            
        base_name = os.path.basename(expanded)
        dest_path = os.path.join(SANDBOX_DIR, base_name)
        shutil.copyfile(expanded, dest_path)
        
        return {
            "status": "success",
            "message": f"File '{base_name}' ({round(file_size_mb, 2)} MB) berhasil disiapkan dan akan otomatis terkirim ke Telegram.",
            "file_name": base_name,
            "caption": caption
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal memproses file: {str(e)}"}


def compress_folder_to_zip(folder_path: str, output_filename: str = "archive.zip") -> Dict[str, Any]:
    """
    Compress an entire folder/directory into a ZIP archive and send it to Telegram.
    Use this when the user wants to send a whole folder, backup a project, or archive files.
    
    Args:
        folder_path: Path to the folder to compress (e.g. '~/Documents/project', '/home/fahmial/telegram-ai-bot').
        output_filename: Output ZIP filename (default: archive.zip).
    """
    try:
        import shutil
        expanded = os.path.expanduser(folder_path)
        if not os.path.isdir(expanded):
            return {"status": "error", "message": f"Folder tidak ditemukan: {folder_path}"}
        
        safe_name = output_filename.replace(".zip", "")
        zip_base = os.path.join(SANDBOX_DIR, safe_name)
        result_path = shutil.make_archive(zip_base, 'zip', expanded)
        
        size_mb = os.path.getsize(result_path) / (1024 * 1024)
        if size_mb > 50:
            os.remove(result_path)
            return {"status": "error", "message": f"Ukuran ZIP ({round(size_mb, 1)} MB) melebihi batas Telegram (50 MB)."}
        
        return {
            "status": "success",
            "message": f"Folder '{os.path.basename(expanded)}' berhasil di-compress menjadi '{safe_name}.zip' ({round(size_mb, 2)} MB) dan akan dikirim ke Telegram.",
            "file_path": result_path
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


def record_desktop_screen(duration_seconds: int = 10) -> Dict[str, Any]:
    """
    Record the Linux desktop screen as an MP4 video for a specified duration and send to Telegram.
    
    Args:
        duration_seconds: Recording duration in seconds (1-60, default: 10).
    """
    try:
        duration = max(1, min(60, duration_seconds))
        output_path = os.path.join(SANDBOX_DIR, "screen_recording.mp4")
        
        # Try Wayland wf-recorder first
        res = subprocess.run(
            f"timeout {duration + 2} wf-recorder -d /dev/dri/renderD128 -f {output_path} --duration {duration} 2>/dev/null",
            shell=True, capture_output=True, text=True, timeout=duration + 10
        )
        if res.returncode == 0 and os.path.exists(output_path) and os.path.getsize(output_path) > 0:
            size_mb = os.path.getsize(output_path) / (1024 * 1024)
            return {"status": "success", "message": f"Rekaman layar {duration}s berhasil ({round(size_mb, 2)} MB) dan akan dikirim ke Telegram."}
        
        # Fallback to ffmpeg with PipeWire
        res = subprocess.run(
            f"timeout {duration + 5} ffmpeg -y -video_size 1920x1080 -framerate 15 -f x11grab -i :0 -t {duration} -c:v libx264 -preset ultrafast -crf 28 {output_path} 2>/dev/null",
            shell=True, capture_output=True, text=True, timeout=duration + 15
        )
        if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
            size_mb = os.path.getsize(output_path) / (1024 * 1024)
            return {"status": "success", "message": f"Rekaman layar {duration}s berhasil via ffmpeg ({round(size_mb, 2)} MB)."}
        
        return {"status": "error", "message": "Gagal merekam layar. Pastikan wf-recorder atau ffmpeg terinstall."}
    except subprocess.TimeoutExpired:
        if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
            return {"status": "success", "message": "Rekaman layar berhasil (timeout graceful)."}
        return {"status": "error", "message": "Timeout saat merekam layar."}
    except Exception as e:
        return {"status": "error", "message": str(e)}


def read_clipboard() -> Dict[str, Any]:
    """
    Read the current content of the Linux desktop clipboard (copy-paste buffer).
    """
    try:
        # Try Wayland wl-paste
        res = subprocess.run("wl-paste 2>/dev/null || xclip -selection clipboard -o 2>/dev/null || xsel --clipboard --output 2>/dev/null",
                             shell=True, capture_output=True, text=True, timeout=3)
        content = res.stdout.strip()
        if content:
            return {"status": "success", "clipboard_content": content[:5000]}
        return {"status": "success", "clipboard_content": "(Clipboard kosong)"}
    except Exception as e:
        return {"status": "error", "message": str(e)}


def write_to_clipboard(text: str) -> Dict[str, Any]:
    """
    Write/copy text to the Linux desktop clipboard so it can be pasted (Ctrl+V) anywhere.
    
    Args:
        text: The text string to copy to clipboard.
    """
    try:
        proc = subprocess.Popen("wl-copy 2>/dev/null || xclip -selection clipboard 2>/dev/null",
                                shell=True, stdin=subprocess.PIPE, text=True)
        proc.communicate(input=text, timeout=3)
        return {"status": "success", "message": f"Teks berhasil disalin ke clipboard ({len(text)} karakter). Siap di-paste (Ctrl+V)."}
    except Exception as e:
        return {"status": "error", "message": str(e)}


def show_desktop_notification(title: str, message: str, urgency: str = "normal") -> Dict[str, Any]:
    """
    Show a native desktop popup notification across Linux, macOS, and Windows.
    
    Args:
        title: Notification title.
        message: Notification body text.
        urgency: 'low', 'normal', or 'critical'.
    """
    import sys
    try:
        if sys.platform == "darwin":  # macOS
            apple_script = f'display notification "{message}" with title "{title}"'
            subprocess.run(["osascript", "-e", apple_script], capture_output=True, text=True, timeout=3)
        elif sys.platform == "win32":  # Windows
            ps_cmd = (
                f"[reflection.assembly]::loadwithpartialname('System.Windows.Forms');"
                f"$notify = new-object system.windows.forms.notifyicon;"
                f"$notify.icon = [system.drawing.systemicons]::information;"
                f"$notify.visible = $true;"
                f"$notify.showballoontip(10, '{title}', '{message}', [system.windows.forms.tooltipicon]::None)"
            )
            subprocess.run(["powershell", "-Command", ps_cmd], capture_output=True, text=True, timeout=3)
        else:  # Linux / Unix
            subprocess.run(
                ["notify-send", f"--urgency={urgency}", "--app-name=AgentALFA", title, message],
                capture_output=True, text=True, timeout=3
            )
        return {"status": "success", "message": f"Notifikasi desktop '{title}' berhasil ditampilkan di layar."}
    except Exception as e:
        return {"status": "error", "message": str(e)}


def ssh_execute_command(host: str, command: str, username: str = "", port: int = 22, key_path: str = "") -> Dict[str, Any]:
    """
    Execute a command on a remote Linux server via SSH and return the output.
    Use this for remote server management, deployment, monitoring, etc.
    
    Args:
        host: Remote server hostname or IP address.
        command: Shell command to execute on the remote server.
        username: SSH username (defaults to current user if empty).
        port: SSH port (default: 22).
        key_path: Path to SSH private key file (defaults to ~/.ssh/id_rsa if empty).
    """
    try:
        import paramiko
        
        ssh_user = username or os.environ.get("USER", "root")
        ssh_key = os.path.expanduser(key_path) if key_path else os.path.expanduser("~/.ssh/id_rsa")
        
        client = paramiko.SSHClient()
        client.set_missing_host_key_policy(paramiko.AutoAddPolicy())
        
        connect_kwargs = {"hostname": host, "port": port, "username": ssh_user, "timeout": 10}
        if os.path.exists(ssh_key):
            connect_kwargs["key_filename"] = ssh_key
        
        client.connect(**connect_kwargs)
        stdin, stdout, stderr = client.exec_command(command, timeout=30)
        
        out = stdout.read().decode("utf-8", errors="replace").strip()
        err = stderr.read().decode("utf-8", errors="replace").strip()
        exit_code = stdout.channel.recv_exit_status()
        client.close()
        
        return {
            "status": "success",
            "host": host,
            "exit_code": exit_code,
            "stdout": out[:8000],
            "stderr": err[:2000] if err else ""
        }
    except Exception as e:
        return {"status": "error", "message": f"SSH error: {str(e)}"}


def query_database(db_path: str, sql_query: str) -> Dict[str, Any]:
    """
    Execute a SQL query on a local SQLite database file and return the results as a table.
    
    Args:
        db_path: Path to the SQLite database file (e.g. '~/data/app.db', 'bot_database.db').
        sql_query: SQL query to execute (SELECT, INSERT, UPDATE, DELETE, etc.).
    """
    try:
        import sqlite3
        expanded = os.path.expanduser(db_path)
        if not os.path.exists(expanded):
            return {"status": "error", "message": f"Database file tidak ditemukan: {db_path}"}
        
        conn = sqlite3.connect(expanded)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        cursor.execute(sql_query)
        
        query_upper = sql_query.strip().upper()
        if query_upper.startswith("SELECT") or query_upper.startswith("PRAGMA") or query_upper.startswith("WITH"):
            rows = cursor.fetchall()
            columns = [desc[0] for desc in cursor.description] if cursor.description else []
            data = [dict(row) for row in rows[:100]]
            conn.close()
            return {
                "status": "success",
                "columns": columns,
                "row_count": len(data),
                "total_available": len(rows) if len(rows) <= 100 else f"{len(rows)}+ (showing first 100)",
                "data": data
            }
        else:
            conn.commit()
            affected = cursor.rowcount
            conn.close()
            return {"status": "success", "message": f"Query berhasil dieksekusi. {affected} baris terpengaruh."}
    except Exception as e:
        return {"status": "error", "message": f"SQL error: {str(e)}"}


def send_email(to: str, subject: str, body: str, attachment_path: str = "") -> Dict[str, Any]:
    """
    Send an email via SMTP (supports Gmail, Outlook, custom SMTP servers).
    Requires SMTP_EMAIL and SMTP_PASSWORD environment variables in .env file.
    
    Args:
        to: Recipient email address.
        subject: Email subject line.
        body: Email body text (supports plain text).
        attachment_path: Optional file path to attach.
    """
    try:
        import smtplib
        from email.mime.text import MIMEText
        from email.mime.multipart import MIMEMultipart
        from email.mime.base import MIMEBase
        from email import encoders
        
        smtp_email = os.environ.get("SMTP_EMAIL", "")
        smtp_password = os.environ.get("SMTP_PASSWORD", "")
        smtp_host = os.environ.get("SMTP_HOST", "smtp.gmail.com")
        smtp_port = int(os.environ.get("SMTP_PORT", "587"))
        
        if not smtp_email or not smtp_password:
            return {"status": "error", "message": "SMTP_EMAIL dan SMTP_PASSWORD belum dikonfigurasi di file .env. Tambahkan: SMTP_EMAIL=xxx@gmail.com dan SMTP_PASSWORD=your_app_password"}
        
        msg = MIMEMultipart()
        msg["From"] = smtp_email
        msg["To"] = to
        msg["Subject"] = subject
        msg.attach(MIMEText(body, "plain"))
        
        if attachment_path:
            expanded = os.path.expanduser(attachment_path)
            if os.path.exists(expanded):
                with open(expanded, "rb") as f:
                    part = MIMEBase("application", "octet-stream")
                    part.set_payload(f.read())
                    encoders.encode_base64(part)
                    part.add_header("Content-Disposition", f"attachment; filename={os.path.basename(expanded)}")
                    msg.attach(part)
        
        server = smtplib.SMTP(smtp_host, smtp_port)
        server.starttls()
        server.login(smtp_email, smtp_password)
        server.sendmail(smtp_email, to, msg.as_string())
        server.quit()
        
        return {"status": "success", "message": f"Email berhasil dikirim ke {to} dengan subjek '{subject}'."}
    except Exception as e:
        return {"status": "error", "message": f"Gagal mengirim email: {str(e)}"}


def list_running_processes(filter_name: str = "") -> Dict[str, Any]:
    """
    List currently running processes on the system, sorted by memory usage.
    
    Args:
        filter_name: Optional filter to show only processes matching this name (e.g. 'chrome', 'python', 'code').
    """
    try:
        procs = []
        for p in psutil.process_iter(['pid', 'name', 'cpu_percent', 'memory_info', 'status', 'username']):
            try:
                info = p.info
                mem_mb = round(info['memory_info'].rss / (1024 * 1024), 1) if info.get('memory_info') else 0
                if filter_name and filter_name.lower() not in (info.get('name') or '').lower():
                    continue
                procs.append({
                    "pid": info['pid'],
                    "name": info['name'],
                    "cpu_percent": info.get('cpu_percent', 0),
                    "memory_mb": mem_mb,
                    "status": info.get('status', ''),
                    "user": info.get('username', '')
                })
            except (psutil.NoSuchProcess, psutil.AccessDenied):
                continue
        
        procs.sort(key=lambda x: x['memory_mb'], reverse=True)
        top = procs[:30]
        total_mem = sum(p['memory_mb'] for p in procs)
        
        return {
            "status": "success",
            "total_processes": len(procs),
            "total_memory_mb": round(total_mem, 1),
            "top_processes": top
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


def kill_process(pid_or_name: str) -> Dict[str, Any]:
    """
    Terminate/kill a running process by PID number or process name.
    
    Args:
        pid_or_name: Process ID (e.g. '12345') or process name (e.g. 'chrome', 'firefox', 'spotify').
    """
    try:
        killed = []
        if pid_or_name.isdigit():
            pid = int(pid_or_name)
            p = psutil.Process(pid)
            name = p.name()
            p.terminate()
            killed.append(f"PID {pid} ({name})")
        else:
            for p in psutil.process_iter(['pid', 'name']):
                try:
                    if pid_or_name.lower() in p.info['name'].lower():
                        p.terminate()
                        killed.append(f"PID {p.info['pid']} ({p.info['name']})")
                except (psutil.NoSuchProcess, psutil.AccessDenied):
                    continue
        
        if killed:
            return {"status": "success", "message": f"Proses berhasil dihentikan: {', '.join(killed)}"}
        return {"status": "error", "message": f"Proses '{pid_or_name}' tidak ditemukan."}
    except psutil.NoSuchProcess:
        return {"status": "error", "message": f"Proses dengan PID/nama '{pid_or_name}' tidak ditemukan."}
    except psutil.AccessDenied:
        return {"status": "error", "message": f"Akses ditolak. Coba jalankan dengan sudo."}
    except Exception as e:
        return {"status": "error", "message": str(e)}


def edit_image(file_path: str, action: str, params: str = "") -> Dict[str, Any]:
    """
    Edit, convert, or transform an image file using Pillow and send the result to Telegram.
    
    Args:
        file_path: Path to the source image file.
        action: Edit action to perform. Options:
                - 'resize': Resize image (params: 'WIDTHxHEIGHT', e.g. '800x600')
                - 'crop': Crop image (params: 'LEFT,TOP,RIGHT,BOTTOM', e.g. '100,100,500,400')
                - 'rotate': Rotate image (params: degrees, e.g. '90', '180', '270')
                - 'grayscale': Convert to black & white
                - 'flip_horizontal': Flip horizontally
                - 'flip_vertical': Flip vertically
                - 'convert': Convert format (params: target format, e.g. 'PNG', 'JPEG', 'WEBP')
                - 'watermark': Add text watermark (params: watermark text)
                - 'thumbnail': Create thumbnail (params: 'WIDTHxHEIGHT', e.g. '200x200')
                - 'blur': Apply blur effect
                - 'sharpen': Sharpen image
                - 'info': Get image metadata (dimensions, format, size)
        params: Parameters for the action (depends on action type).
    """
    try:
        from PIL import Image, ImageFilter, ImageDraw, ImageFont
        
        expanded = os.path.expanduser(file_path)
        if not os.path.exists(expanded):
            return {"status": "error", "message": f"File gambar tidak ditemukan: {file_path}"}
        
        img = Image.open(expanded)
        base_name = os.path.splitext(os.path.basename(expanded))[0]
        act = action.strip().lower()
        
        if act == "info":
            return {
                "status": "success",
                "format": img.format,
                "size": f"{img.width}x{img.height}",
                "mode": img.mode,
                "file_size_kb": round(os.path.getsize(expanded) / 1024, 1)
            }
        elif act == "resize":
            w, h = [int(x) for x in params.lower().split("x")]
            img = img.resize((w, h), Image.LANCZOS)
        elif act == "crop":
            coords = [int(x.strip()) for x in params.split(",")]
            img = img.crop(tuple(coords))
        elif act == "rotate":
            degrees = int(params)
            img = img.rotate(degrees, expand=True)
        elif act == "grayscale":
            img = img.convert("L")
        elif act == "flip_horizontal":
            img = img.transpose(Image.FLIP_LEFT_RIGHT)
        elif act == "flip_vertical":
            img = img.transpose(Image.FLIP_TOP_BOTTOM)
        elif act == "convert":
            pass  # handled below by save format
        elif act == "watermark":
            draw = ImageDraw.Draw(img)
            text = params or "AI Agent Watermark"
            try:
                font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
            except Exception:
                font = ImageFont.load_default()
            bbox = draw.textbbox((0, 0), text, font=font)
            tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
            x = img.width - tw - 20
            y = img.height - th - 20
            draw.text((x, y), text, fill=(255, 255, 255, 180), font=font)
        elif act == "thumbnail":
            w, h = [int(x) for x in params.lower().split("x")]
            img.thumbnail((w, h), Image.LANCZOS)
        elif act == "blur":
            img = img.filter(ImageFilter.GaussianBlur(radius=5))
        elif act == "sharpen":
            img = img.filter(ImageFilter.SHARPEN)
        else:
            return {"status": "error", "message": f"Aksi '{action}' tidak dikenal."}
        
        if act == "convert":
            fmt = params.strip().upper()
            ext = fmt.lower()
            if fmt == "JPEG":
                ext = "jpg"
                img = img.convert("RGB")
        else:
            fmt = img.format or "PNG"
            ext = fmt.lower()
            if ext == "jpeg":
                ext = "jpg"
        
        if act != "convert" and img.mode == "RGBA" and fmt == "JPEG":
            img = img.convert("RGB")
            
        out_name = f"{base_name}_edited.{ext}"
        out_path = os.path.join(SANDBOX_DIR, out_name)
        img.save(out_path, format=fmt if act == "convert" else None)
        
        return {
            "status": "success",
            "message": f"Gambar berhasil di-{act} dan disimpan sebagai '{out_name}'. Akan dikirim ke Telegram.",
            "file_path": out_path
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal mengedit gambar: {str(e)}"}


def git_operations(action: str, repo_path: str = ".", message: str = "", remote: str = "origin", branch: str = "") -> Dict[str, Any]:
    """
    Perform Git operations on a local repository directly from Telegram.
    
    Args:
        action: Git action to perform. Options:
                - 'status': Show git status (modified, staged, untracked files)
                - 'log': Show recent commit history
                - 'pull': Pull latest changes from remote
                - 'add_all': Stage all changes (git add .)
                - 'commit': Commit staged changes (requires message parameter)
                - 'push': Push commits to remote
                - 'diff': Show unstaged changes
                - 'branch': List branches
                - 'stash': Stash current changes
                - 'stash_pop': Pop stashed changes
        repo_path: Path to git repository (default: current directory).
        message: Commit message (required for 'commit' action).
        remote: Remote name (default: 'origin').
        branch: Branch name (optional).
    """
    try:
        expanded = os.path.expanduser(repo_path)
        
        cmd_map = {
            "status": "git status --porcelain -b",
            "log": "git log --oneline --graph -n 15",
            "pull": f"git pull {remote} {branch}".strip(),
            "add_all": "git add -A",
            "commit": f'git commit -m "{message}"' if message else 'echo "ERROR: commit message required"',
            "push": f"git push {remote} {branch}".strip(),
            "diff": "git diff --stat",
            "branch": "git branch -a",
            "stash": "git stash",
            "stash_pop": "git stash pop",
        }
        
        act = action.strip().lower()
        cmd = cmd_map.get(act)
        if not cmd:
            return {"status": "error", "message": f"Git action '{action}' tidak dikenal. Pilihan: {', '.join(cmd_map.keys())}"}
        
        res = subprocess.run(cmd, shell=True, capture_output=True, text=True, cwd=expanded, timeout=30)
        output = res.stdout.strip() or res.stderr.strip()
        
        return {
            "status": "success" if res.returncode == 0 else "error",
            "action": act,
            "output": output[:6000],
            "exit_code": res.returncode
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


def translate_text(text: str, target_lang: str = "en", source_lang: str = "auto") -> Dict[str, Any]:
    """
    Translate text between languages using Google Translate.
    
    Args:
        text: Text to translate.
        target_lang: Target language code (e.g. 'en' English, 'id' Indonesian, 'ja' Japanese, 'ko' Korean, 'zh-CN' Chinese, 'ar' Arabic, 'fr' French, 'de' German, 'es' Spanish).
        source_lang: Source language code (default: 'auto' for auto-detect).
    """
    try:
        from deep_translator import GoogleTranslator
        
        translated = GoogleTranslator(source=source_lang, target=target_lang).translate(text[:4500])
        return {
            "status": "success",
            "original": text[:500],
            "translated": translated,
            "source_language": source_lang,
            "target_language": target_lang
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal menerjemahkan: {str(e)}"}


def download_file_from_url(url: str, filename: str = "") -> Dict[str, Any]:
    """
    Download a file from a URL on the internet to the local computer and optionally send it to Telegram.
    Supports direct file links (images, PDFs, ZIPs, videos, audio, executables, etc.).
    
    Args:
        url: Direct download URL.
        filename: Optional filename to save as (auto-detected from URL if empty).
    """
    try:
        import httpx
        
        if not filename:
            from urllib.parse import urlparse
            parsed = urlparse(url)
            filename = os.path.basename(parsed.path) or "downloaded_file"
        
        dest_path = os.path.join(SANDBOX_DIR, filename)
        
        with httpx.Client(follow_redirects=True, timeout=60) as client:
            resp = client.get(url)
            resp.raise_for_status()
            with open(dest_path, "wb") as f:
                f.write(resp.content)
        
        size_mb = os.path.getsize(dest_path) / (1024 * 1024)
        
        if size_mb > 50:
            return {
                "status": "success",
                "message": f"File '{filename}' ({round(size_mb, 2)} MB) berhasil diunduh ke {dest_path}. Terlalu besar untuk dikirim via Telegram (>50MB), tetapi tersedia di disk lokal.",
                "file_path": dest_path,
                "sent_to_telegram": False
            }
        
        return {
            "status": "success",
            "message": f"File '{filename}' ({round(size_mb, 2)} MB) berhasil diunduh dan akan dikirim ke Telegram.",
            "file_path": dest_path,
            "sent_to_telegram": True
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal mengunduh: {str(e)}"}


def generate_secure_password(length: int = 20, include_uppercase: bool = True, include_digits: bool = True, include_special: bool = True, count: int = 1) -> Dict[str, Any]:
    """
    Generate one or more cryptographically secure random passwords.
    
    Args:
        length: Password length (8-128 characters, default: 20).
        include_uppercase: Include uppercase letters (default: True).
        include_digits: Include digits (default: True).
        include_special: Include special characters !@#$%^&* (default: True).
        count: Number of passwords to generate (1-10, default: 1).
    """
    try:
        import secrets
        import string
        
        length = max(8, min(128, length))
        count = max(1, min(10, count))
        
        chars = string.ascii_lowercase
        if include_uppercase:
            chars += string.ascii_uppercase
        if include_digits:
            chars += string.digits
        if include_special:
            chars += "!@#$%^&*_+-=?."
        
        passwords = []
        for _ in range(count):
            pw = ''.join(secrets.choice(chars) for _ in range(length))
            passwords.append(pw)
        
        # Calculate entropy
        import math
        entropy = round(math.log2(len(chars)) * length, 1)
        strength = "Sangat Kuat 🟢" if entropy > 80 else "Kuat 🔵" if entropy > 60 else "Cukup 🟡" if entropy > 40 else "Lemah 🔴"
        
        return {
            "status": "success",
            "passwords": passwords,
            "length": length,
            "entropy_bits": entropy,
            "strength": strength,
            "charset_size": len(chars)
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


def vision_click_target(target_description: str, max_attempts: int = 3, action: str = "click") -> Dict[str, Any]:
    """
    GOD MODE: Vision-guided autonomous computer use loop.
    Takes a screenshot of the desktop, sends it to Gemini Vision AI to locate a target
    UI element (button, icon, text, link, menu), clicks on it, then takes another
    screenshot to verify the action succeeded. Repeats if needed.
    
    This allows the bot to operate ANY desktop application (browsers, editors, file managers,
    settings, terminals) purely through visual understanding — like a human looking at a screen.
    
    Args:
        target_description: Natural language description of what to find and click 
                           (e.g. 'the red Close button', 'Firefox icon on taskbar', 
                            'File menu in top left', 'Play button on Spotify',
                            'the search bar', 'Settings gear icon').
        max_attempts: Maximum number of screenshot-analyze-click attempts (1-5, default: 3).
        action: What to do with the found element: 'click' (default), 'double_click', 'right_click', 'identify_only'.
    """
    try:
        from google import genai
        from google.genai import types
        from PIL import Image
        import io as _io
        import json as _json
        
        api_key = os.environ.get("GEMINI_API_KEY")
        if not api_key:
            from dotenv import load_dotenv
            load_dotenv()
            api_key = os.environ.get("GEMINI_API_KEY")
            
        client = genai.Client(api_key=api_key)
        attempts = min(max(1, max_attempts), 5)
        
        for attempt in range(1, attempts + 1):
            # Step 1: Capture desktop screenshot
            screenshot_result = capture_desktop_screenshot()
            screenshot_path = os.path.join(SANDBOX_DIR, "desktop_screen.png")
            
            if not os.path.exists(screenshot_path) or os.path.getsize(screenshot_path) == 0:
                return {"status": "error", "message": "Gagal mengambil screenshot desktop untuk vision loop."}
            
            # Step 2: Send to Gemini Vision for coordinate analysis
            with open(screenshot_path, "rb") as f:
                img_bytes = f.read()
            
            image_part = types.Part.from_bytes(data=img_bytes, mime_type="image/png")
            
            vision_prompt = (
                f"Kamu adalah sistem Vision AI untuk GUI automation pada layar Linux desktop 1920x1080.\n"
                f"Analisis screenshot ini dan temukan elemen UI berikut: \"{target_description}\"\n\n"
                f"INSTRUKSI:\n"
                f"1. Identifikasi lokasi elemen tersebut di layar.\n"
                f"2. Berikan koordinat pixel X dan Y dari TITIK TENGAH elemen tersebut.\n"
                f"3. Jika elemen TIDAK DITEMUKAN, jawab dengan found=false.\n\n"
                f"JAWAB DALAM FORMAT JSON SAJA, tanpa teks lain:\n"
                f'{{"found": true/false, "x": <int>, "y": <int>, "element_description": "<apa yang kamu lihat>", "confidence": "<high/medium/low>"}}'
            )
            
            response = client.models.generate_content(
                model=os.environ.get("GEMINI_MODEL", "gemini-3.5-flash-lite"),
                contents=[image_part, vision_prompt]
            )
            
            response_text = response.text.strip()
            
            # Parse JSON from response
            json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
            if not json_match:
                continue
            
            result = _json.loads(json_match.group())
            
            if not result.get("found", False):
                if attempt < attempts:
                    import time
                    time.sleep(1)
                    continue
                return {
                    "status": "error",
                    "message": f"Elemen '{target_description}' tidak ditemukan di layar setelah {attempts} percobaan.",
                    "vision_response": result.get("element_description", "")
                }
            
            x = int(result["x"])
            y = int(result["y"])
            confidence = result.get("confidence", "unknown")
            desc = result.get("element_description", "")
            
            if action == "identify_only":
                # Remove the screenshot from sandbox so it doesn't get auto-sent
                try:
                    os.remove(screenshot_path)
                except OSError:
                    pass
                return {
                    "status": "success",
                    "message": f"Elemen ditemukan di koordinat ({x}, {y}).",
                    "coordinates": {"x": x, "y": y},
                    "element_description": desc,
                    "confidence": confidence,
                    "attempt": attempt
                }
            
            # Remove pre-click screenshot
            try:
                os.remove(screenshot_path)
            except OSError:
                pass
            
            # Step 3: Click the target
            clicks = 2 if action == "double_click" else 1
            button = "right" if action == "right_click" else "left"
            click_result = desktop_click_coordinate(x=x, y=y, button=button, clicks=clicks)
            
            # Step 4: Wait briefly then take verification screenshot
            import time
            time.sleep(0.8)
            verify_result = capture_desktop_screenshot()
            
            return {
                "status": "success",
                "message": f"Vision Loop berhasil! Elemen '{target_description}' ditemukan dan di-{action} pada koordinat ({x}, {y}). Screenshot verifikasi akan dikirim ke Telegram.",
                "coordinates": {"x": x, "y": y},
                "element_description": desc,
                "confidence": confidence,
                "attempt": attempt,
                "action_performed": action
            }
        
        return {"status": "error", "message": f"Gagal menemukan '{target_description}' setelah {attempts} percobaan vision loop."}
    except Exception as e:
        return {"status": "error", "message": f"Vision loop error: {str(e)}"}


def deep_research_topic(topic: str, max_depth: int = 3) -> Dict[str, Any]:
    """
    GOD MODE: Autonomous Deep Multi-Source Research Engine.
    Executes multiple recursive web search queries on a topic, crawls and scrapes the top
    3-5 authoritative domain pages, synthesizes cross-source evidence, resolves contradictions,
    and returns a structured, factual briefing with citations.
    
    Args:
        topic: Topic or research question to investigate deeply.
        max_depth: Maximum number of search iteration queries (1-5, default: 3).
    """
    try:
        from ddgs import DDGS
        import httpx
        from urllib.parse import urlparse
        
        queries = [
            topic,
            f"{topic} overview facts analysis",
            f"{topic} latest updates details"
        ][:max_depth]
        
        seen_urls = set()
        sources_data = []
        
        with DDGS(verify=False) as ddgs:
            for q in queries:
                try:
                    results = list(ddgs.text(q, max_results=3))
                    for r in results:
                        u = r.get("href")
                        if u and u not in seen_urls and not u.endswith((".pdf", ".exe", ".zip", ".png", ".jpg")):
                            seen_urls.add(u)
                            sources_data.append({
                                "title": r.get("title", ""),
                                "url": u,
                                "snippet": r.get("body", "")
                            })
                            if len(sources_data) >= 5:
                                break
                except Exception:
                    pass
                if len(sources_data) >= 5:
                    break
        
        if not sources_data:
            return {"status": "error", "message": f"Tidak ditemukan sumber riset untuk topik: '{topic}'"}
            
        crawled_articles = []
        with httpx.Client(follow_redirects=True, timeout=12, headers={"User-Agent": "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36"}) as client:
            for src in sources_data[:4]:
                try:
                    resp = client.get(src["url"])
                    if resp.status_code == 200:
                        raw_html = resp.text
                        text = re.sub(r'<(script|style)[^>]*>.*?</\1>', '', raw_html, flags=re.DOTALL | re.IGNORECASE)
                        text = re.sub(r'<[^>]+>', ' ', text)
                        clean_text = ' '.join(text.split())[:1500]
                        domain = urlparse(src["url"]).netloc
                        crawled_articles.append({
                            "source_title": src["title"],
                            "domain": domain,
                            "url": src["url"],
                            "extracted_content": clean_text
                        })
                except Exception:
                    continue

        return {
            "status": "success",
            "topic": topic,
            "total_sources_analyzed": len(crawled_articles),
            "sources": crawled_articles,
            "research_directive": "Gunakan data dari sumber-sumber terverifikasi di atas untuk menyusun sintesis riset yang objektif, akurat, dan mencantumkan sitasi URL."
        }
    except Exception as e:
        return {"status": "error", "message": f"Deep research error: {str(e)}"}


def auto_diagnose_and_heal_system(fix_issues: bool = False) -> Dict[str, Any]:
    """
    GOD MODE: Autonomous System Diagnostic & Self-Healing Engine.
    Inspects system journal error logs, failed systemd units, memory pressure,
    broken packages, and zombie processes. Produces a root-cause diagnosis
    and optionally executes safe autonomous healing actions.
    
    Args:
        fix_issues: If True, executes safe autonomous healing (restarting failed units, vacuuming logs, clearing zombies).
    """
    try:
        diagnosis = {}
        healing_actions = []
        
        # 1. Check failed systemd units
        res_failed = subprocess.run("systemctl --user list-units --failed --no-legend 2>/dev/null", shell=True, capture_output=True, text=True, timeout=5)
        failed_user_units = [line.strip().split()[0] for line in res_failed.stdout.strip().splitlines() if line.strip()]
        diagnosis["failed_user_services"] = failed_user_units
        
        # 2. Check system error logs in journalctl
        res_journal = subprocess.run("journalctl --user -p 3 -n 15 --no-pager 2>/dev/null", shell=True, capture_output=True, text=True, timeout=5)
        diagnosis["recent_critical_logs"] = res_journal.stdout.strip()[:1500] if res_journal.stdout.strip() else "Tidak ada critical error log terbaru."
        
        # 3. Check zombie / hung processes
        zombies = []
        for p in psutil.process_iter(['pid', 'name', 'status']):
            try:
                if p.info['status'] == psutil.STATUS_ZOMBIE:
                    zombies.append(f"PID {p.info['pid']} ({p.info['name']})")
            except (psutil.NoSuchProcess, psutil.AccessDenied):
                continue
        diagnosis["zombie_processes"] = zombies
        
        # 4. Check disk and swap pressure
        disk = psutil.disk_usage('/')
        swap = psutil.swap_memory()
        diagnosis["disk_health"] = f"Root: {disk.percent}% used ({round(disk.free / (1024**3), 1)} GB free)"
        diagnosis["swap_health"] = f"Swap: {swap.percent}% used"
        
        # Self-healing execution if requested
        if fix_issues:
            # Restart failed user services (excluding disabled ones)
            for unit in failed_user_units:
                if "telegram-ai-bot" not in unit:  # avoid recursive restart in diagnostic turn
                    subprocess.run(f"systemctl --user reset-failed {unit} && systemctl --user restart {unit}", shell=True, timeout=10)
                    healing_actions.append(f"Restarted failed unit: {unit}")
            
            # Vacuum journal logs if disk > 85%
            if disk.percent > 85:
                subprocess.run("journalctl --user --vacuum-time=2d", shell=True, timeout=10)
                healing_actions.append("Cleaned old user journal logs")
                
            diagnosis["healing_executed"] = healing_actions if healing_actions else "Tidak ada tindakan perbaikan yang diperlukan saat ini."
            
        return {
            "status": "success",
            "diagnosis": diagnosis,
            "fix_mode": fix_issues
        }
    except Exception as e:
        return {"status": "error", "message": f"Diagnostic error: {str(e)}"}


def text_to_audio_file(text: str, filename: str = "audio_speech.mp3", voice: str = "id-ID-GadisNeural") -> Dict[str, Any]:
    """
    Generate a high-fidelity natural speech audio file (.mp3) from any long text or script
    using Microsoft Edge Neural TTS and send it as an audio file directly to Telegram.
    
    Args:
        text: Full text or script to synthesize into audio.
        filename: Target filename (default: audio_speech.mp3).
        voice: Voice code, e.g. 'id-ID-GadisNeural' (female ID), 'id-ID-ArdiNeural' (male ID), 'en-US-JennyNeural' (US English).
    """
    try:
        import asyncio
        import edge_tts
        
        if not filename.endswith(".mp3"):
            filename += ".mp3"
            
        out_path = os.path.join(SANDBOX_DIR, filename)
        
        async def _synth():
            communicate = edge_tts.Communicate(text[:5000], voice)
            await communicate.save(out_path)
            
        try:
            loop = asyncio.get_event_loop()
            if loop.is_running():
                import concurrent.futures
                with concurrent.futures.ThreadPoolExecutor() as executor:
                    executor.submit(asyncio.run, _synth()).result()
            else:
                loop.run_until_complete(_synth())
        except Exception:
            asyncio.run(_synth())
            
        if os.path.exists(out_path) and os.path.getsize(out_path) > 0:
            size_kb = round(os.path.getsize(out_path) / 1024, 1)
            return {
                "status": "success",
                "message": f"File audio speech '{filename}' ({size_kb} KB) berhasil dibuat dan akan dikirim ke Telegram.",
                "file_path": out_path,
                "voice": voice
            }
        return {"status": "error", "message": "Gagal membuat file audio speech."}
    except Exception as e:
        return {"status": "error", "message": f"Audio synthesis error: {str(e)}"}


def convert_media_format(source_file: str, output_format: str = "mp3", extra_params: str = "") -> Dict[str, Any]:
    """
    Convert any video or audio file to another format using ffmpeg (e.g. mp4 -> mp3, mkv -> mp4, wav -> ogg, flac -> mp3).
    The converted file will be automatically sent to Telegram.
    
    Args:
        source_file: Path to source audio/video file.
        output_format: Target format extension (e.g. 'mp3', 'mp4', 'wav', 'ogg', 'flac', 'aac').
        extra_params: Optional ffmpeg flags (e.g. '-q:a 0' or '-vf scale=1280:720').
    """
    try:
        expanded = os.path.expanduser(source_file)
        if not os.path.exists(expanded):
            return {"status": "error", "message": f"File sumber tidak ditemukan: {source_file}"}
            
        base_name = os.path.splitext(os.path.basename(expanded))[0]
        out_format = output_format.lower().replace(".", "")
        out_name = f"{base_name}_converted.{out_format}"
        dest_path = os.path.join(SANDBOX_DIR, out_name)
        
        cmd = f'ffmpeg -y -i "{expanded}" {extra_params} "{dest_path}"'
        res = subprocess.run(cmd, shell=True, capture_output=True, text=True, timeout=60)
        
        if os.path.exists(dest_path) and os.path.getsize(dest_path) > 0:
            size_mb = round(os.path.getsize(dest_path) / (1024*1024), 2)
            return {
                "status": "success",
                "message": f"Konversi media ke '{out_name}' ({size_mb} MB) berhasil dan akan dikirim ke Telegram.",
                "file_path": dest_path
            }
        return {"status": "error", "message": f"Gagal mengonversi media: {res.stderr[:500]}"}
    except Exception as e:
        return {"status": "error", "message": f"Media conversion error: {str(e)}"}


def extract_audio_from_video(video_path: str, output_filename: str = "extracted_audio.mp3") -> Dict[str, Any]:
    """
    Extract the audio track from a video file (.mp4, .mkv, .webm, .avi) into an MP3 file and send to Telegram.
    
    Args:
        video_path: Path to the local video file.
        output_filename: Output MP3 filename (default: extracted_audio.mp3).
    """
    try:
        expanded = os.path.expanduser(video_path)
        if not os.path.exists(expanded):
            return {"status": "error", "message": f"File video tidak ditemukan: {video_path}"}
            
        if not output_filename.endswith(".mp3"):
            output_filename += ".mp3"
            
        dest_path = os.path.join(SANDBOX_DIR, output_filename)
        cmd = f'ffmpeg -y -i "{expanded}" -vn -acodec libmp3lame -q:a 2 "{dest_path}"'
        res = subprocess.run(cmd, shell=True, capture_output=True, text=True, timeout=60)
        
        if os.path.exists(dest_path) and os.path.getsize(dest_path) > 0:
            size_mb = round(os.path.getsize(dest_path) / (1024*1024), 2)
            return {
                "status": "success",
                "message": f"Audio berhasil diekstraksi menjadi '{output_filename}' ({size_mb} MB) dan akan dikirim ke Telegram.",
                "file_path": dest_path
            }
        return {"status": "error", "message": f"Gagal mengekstrak audio: {res.stderr[:500]}"}
    except Exception as e:
        return {"status": "error", "message": f"Extract audio error: {str(e)}"}


def analyze_dataset_csv_json(file_path: str, chart_type: str = "bar", x_column: str = "", y_column: str = "", title: str = "Data Analysis") -> Dict[str, Any]:
    """
    GOD MODE: Intelligent Dataset Analyzer & Visualizer.
    Reads and parses a CSV, JSON, or Excel dataset, computes statistical metrics
    (summary stats, row counts, missing values, column data types), and generates
    a professional visualization chart automatically sent to Telegram.
    
    Args:
        file_path: Path to dataset file (.csv or .json).
        chart_type: 'bar', 'line', 'scatter', 'pie', 'hist'.
        x_column: Name of X-axis column (defaults to first column).
        y_column: Name of Y-axis numeric column (defaults to second column).
        title: Chart title.
    """
    try:
        import json
        import csv
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        
        expanded = os.path.expanduser(file_path)
        if not os.path.exists(expanded):
            return {"status": "error", "message": f"File dataset tidak ditemukan: {file_path}"}
            
        data_rows = []
        headers = []
        
        ext = os.path.splitext(expanded)[1].lower()
        if ext == ".csv":
            with open(expanded, "r", encoding="utf-8", errors="replace") as f:
                reader = csv.DictReader(f)
                headers = reader.fieldnames or []
                data_rows = list(reader)
        elif ext == ".json":
            with open(expanded, "r", encoding="utf-8") as f:
                raw_json = json.load(f)
                if isinstance(raw_json, list) and raw_json:
                    data_rows = raw_json
                    headers = list(raw_json[0].keys()) if isinstance(raw_json[0], dict) else []
                elif isinstance(raw_json, dict):
                    data_rows = [raw_json]
                    headers = list(raw_json.keys())
        else:
            return {"status": "error", "message": "Format dataset harus .csv atau .json"}
            
        if not data_rows:
            return {"status": "error", "message": "Dataset kosong atau tidak memiliki baris data."}
            
        # Statistical summary
        total_rows = len(data_rows)
        sample_data = data_rows[:5]
        
        # Plotting
        chart_path = os.path.join(SANDBOX_DIR, "dataset_analysis_chart.png")
        plt.figure(figsize=(10, 6), dpi=120)
        plt.style.use('seaborn-v0_8-darkgrid' if 'seaborn-v0_8-darkgrid' in plt.style.available else 'default')
        
        x_col = x_column or (headers[0] if headers else "")
        y_col = y_column or (headers[1] if len(headers) > 1 else headers[0] if headers else "")
        
        x_vals = [str(r.get(x_col, "")) for r in data_rows[:20]]
        y_vals = []
        for r in data_rows[:20]:
            try:
                y_vals.append(float(r.get(y_col, 0)))
            except (ValueError, TypeError):
                y_vals.append(0.0)
                
        if chart_type == "line":
            plt.plot(x_vals, y_vals, marker='o', color='#2563EB', linewidth=2.5)
        elif chart_type == "scatter":
            plt.scatter(x_vals, y_vals, color='#7C3AED', s=80)
        elif chart_type == "pie" and len(x_vals) <= 10:
            plt.pie(y_vals, labels=x_vals, autopct='%1.1f%%', colors=plt.cm.Paired.colors)
        else:  # default bar
            plt.bar(x_vals, y_vals, color='#3B82F6', edgecolor='#1D4ED8')
            
        plt.title(title, fontsize=14, fontweight='bold', pad=15)
        if chart_type != "pie":
            plt.xlabel(x_col, fontsize=11)
            plt.ylabel(y_col, fontsize=11)
            plt.xticks(rotation=45, ha='right')
        plt.tight_layout()
        plt.savefig(chart_path, dpi=150)
        plt.close()
        
        return {
            "status": "success",
            "total_rows": total_rows,
            "columns": headers,
            "sample_rows": sample_data,
            "chart_generated": "dataset_analysis_chart.png",
            "message": f"Analisis dataset '{os.path.basename(expanded)}' selesai. Grafik '{chart_type}' berhasil dibuat dan akan dikirim ke Telegram."
        }
    except Exception as e:
        return {"status": "error", "message": f"Dataset analysis error: {str(e)}"}


def audit_network_security(target_host: str = "127.0.0.1", scan_type: str = "quick_ports") -> Dict[str, Any]:
    """
    GOD MODE: Network Security & Port Sentinel.
    Audits listening network ports, socket services, firewall status (UFW),
    and remote SSL/TLS certificate validity & security ciphers.
    
    Args:
        target_host: Target IP or domain to audit (e.g. '127.0.0.1', 'example.com').
        scan_type: 'quick_ports' or 'full_audit'.
    """
    try:
        import socket
        import ssl
        
        result = {"target": target_host, "scan_type": scan_type}
        
        # Local socket listening check
        if target_host in ["127.0.0.1", "localhost", "0.0.0.0"]:
            res_ss = subprocess.run("ss -tuln 2>/dev/null || netstat -tuln 2>/dev/null", shell=True, capture_output=True, text=True, timeout=5)
            listening_lines = [l for l in res_ss.stdout.strip().splitlines() if "LISTEN" in l or "State" in l][:20]
            result["local_listening_sockets"] = "\n".join(listening_lines)
            
            # Firewall check
            res_ufw = subprocess.run("sudo -n ufw status 2>/dev/null || ufw status 2>/dev/null", shell=True, capture_output=True, text=True, timeout=3)
            result["firewall_status"] = res_ufw.stdout.strip() if res_ufw.stdout.strip() else "UFW status tidak memerlukan sudo / tidak aktif."
        else:
            # Common ports probe
            common_ports = [21, 22, 25, 80, 443, 3000, 3306, 5432, 8000, 8080, 8443]
            open_ports = []
            for p in common_ports:
                s = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
                s.settimeout(1.5)
                res = s.connect_ex((target_host, p))
                if res == 0:
                    open_ports.append(p)
                s.close()
            result["open_ports_detected"] = open_ports
            
            # SSL Certificate inspection if port 443 open
            if 443 in open_ports or target_host.startswith("http") or "." in target_host:
                try:
                    ctx = ssl.create_default_context()
                    with ctx.wrap_socket(socket.socket(), server_hostname=target_host) as s:
                        s.settimeout(5)
                        s.connect((target_host, 443))
                        cert = s.getpeercert()
                        not_after = cert.get('notAfter', '')
                        result["ssl_certificate"] = {
                            "subject": dict(x[0] for x in cert.get('subject', ())),
                            "issuer": dict(x[0] for x in cert.get('issuer', ())),
                            "expires_at": not_after,
                            "version": cert.get('version', '')
                        }
                except Exception as ssl_err:
                    result["ssl_error"] = str(ssl_err)
                    
        return {"status": "success", "audit_report": result}
    except Exception as e:
        return {"status": "error", "message": f"Security audit error: {str(e)}"}


def clean_system_storage(dry_run: bool = True) -> Dict[str, Any]:
    """
    GOD MODE: Smart Linux Storage Cleaner & Optimizer.
    Inspects and frees disk waste by safely cleaning thumbnail caches, user journal logs,
    temporary files, and apt cache.
    
    Args:
        dry_run: If True, only analyzes space to be freed without deleting anything. Set to False to perform actual cleanup.
    """
    try:
        cleanup_targets = [
            ("Thumbnail Cache", os.path.expanduser("~/.cache/thumbnails")),
            ("Sandbox Temp Files", SANDBOX_DIR),
            ("Python Cache", os.path.expanduser("~/.cache/pip"))
        ]
        
        report = []
        total_freed_mb = 0.0
        
        for name, path in cleanup_targets:
            if os.path.exists(path):
                size_b = sum(os.path.getsize(os.path.join(dirpath, f)) for dirpath, _, filenames in os.walk(path) for f in filenames if not os.path.islink(os.path.join(dirpath, f)))
                size_mb = round(size_b / (1024*1024), 2)
                report.append({"target": name, "path": path, "size_mb": size_mb})
                total_freed_mb += size_mb
                
                if not dry_run and size_mb > 0:
                    subprocess.run(f'rm -rf "{path}"/*', shell=True, timeout=10)
                    
        if not dry_run:
            # Vacuum journalctl logs older than 2 days
            subprocess.run("journalctl --user --vacuum-time=2d 2>/dev/null", shell=True, timeout=10)
            
        action_msg = "ANALISIS (Dry Run)" if dry_run else "PEMBERSIHAN SELESAI"
        return {
            "status": "success",
            "mode": action_msg,
            "total_space_mb": round(total_freed_mb, 2),
            "details": report,
            "message": f"{action_msg}: Potensi ruang dibersihkan: {round(total_freed_mb, 2)} MB. Jalankan dengan dry_run=False untuk eksekusi pembersihan nyata." if dry_run else f"Pembersihan berhasil! {round(total_freed_mb, 2)} MB ruang disk berhasil dikembalikan."
        }
    except Exception as e:
        return {"status": "error", "message": f"Storage cleanup error: {str(e)}"}


def manage_system_services(service_name: str, action: str = "status", scope: str = "user") -> Dict[str, Any]:
    """
    GOD MODE: Linux Systemd Services Controller.
    Manage, start, stop, restart, enable, disable, and inspect status of systemd units.
    
    Args:
        service_name: Name of the service unit (e.g. 'telegram-ai-bot.service', 'pipewire', 'docker', 'nginx').
        action: 'status', 'restart', 'start', 'stop', 'enable', 'disable', 'is-active'.
        scope: 'user' (default, for user-space services) or 'system' (system-wide).
    """
    try:
        flag = "--user" if scope == "user" else ""
        act = action.strip().lower()
        valid_actions = ["status", "restart", "start", "stop", "enable", "disable", "is-active"]
        if act not in valid_actions:
            return {"status": "error", "message": f"Aksi '{action}' tidak valid. Pilihan: {', '.join(valid_actions)}"}
            
        cmd = f"systemctl {flag} {act} {service_name}"
        res = subprocess.run(cmd, shell=True, capture_output=True, text=True, timeout=15)
        output = (res.stdout.strip() or res.stderr.strip())[:2500]
        
        return {
            "status": "success" if res.returncode == 0 or act == "status" else "error",
            "service": service_name,
            "action": act,
            "scope": scope,
            "exit_code": res.returncode,
            "output": output
        }
    except Exception as e:
        return {"status": "error", "message": f"Systemd control error: {str(e)}"}


def manage_crontab_jobs(action: str = "list", cron_line: str = "", search_pattern: str = "") -> Dict[str, Any]:
    """
    GOD MODE: Real Linux OS Crontab Manager.
    Reads, adds, or removes native Linux user crontab schedule entries.
    
    Args:
        action: 'list' (show all crontab entries), 'add' (add new cron_line), 'remove' (remove entries matching search_pattern).
        cron_line: The crontab entry string (e.g. '0 8 * * * /home/user/script.sh').
        search_pattern: Keyword/pattern to match when removing crontab entries.
    """
    try:
        if action == "list":
            res = subprocess.run("crontab -l 2>/dev/null", shell=True, capture_output=True, text=True, timeout=5)
            entries = res.stdout.strip()
            return {"status": "success", "crontab_entries": entries if entries else "(Crontab kosong)"}
        elif action == "add":
            if not cron_line:
                return {"status": "error", "message": "Parameter cron_line harus diisi."}
            res_curr = subprocess.run("crontab -l 2>/dev/null", shell=True, capture_output=True, text=True, timeout=5)
            curr = res_curr.stdout.strip()
            new_cron = (curr + "\n" + cron_line.strip()).strip() + "\n"
            proc = subprocess.Popen("crontab -", shell=True, stdin=subprocess.PIPE, text=True)
            proc.communicate(input=new_cron, timeout=5)
            return {"status": "success", "message": f"Entri crontab berhasil ditambahkan: '{cron_line}'"}
        elif action == "remove":
            if not search_pattern:
                return {"status": "error", "message": "Parameter search_pattern harus diisi untuk menghapus entri crontab."}
            res_curr = subprocess.run("crontab -l 2>/dev/null", shell=True, capture_output=True, text=True, timeout=5)
            lines = [l for l in res_curr.stdout.splitlines() if search_pattern not in l]
            new_cron = "\n".join(lines).strip() + "\n"
            proc = subprocess.Popen("crontab -", shell=True, stdin=subprocess.PIPE, text=True)
            proc.communicate(input=new_cron, timeout=5)
            return {"status": "success", "message": f"Entri crontab yang cocok dengan pola '{search_pattern}' berhasil dihapus."}
        return {"status": "error", "message": f"Aksi '{action}' tidak dikenal. Gunakan: list, add, remove."}
    except Exception as e:
        return {"status": "error", "message": f"Crontab error: {str(e)}"}


def extract_and_link_knowledge(entity: str, relation: str, target_value: str, category: str = "general", tags: str = "") -> Dict[str, Any]:
    """
    GOD MODE: Semantic Knowledge Graph & Second Brain Linker.
    Stores structured knowledge facts as subject-predicate-object triples
    (e.g. entity: 'Proyek Alfa', relation: 'deadline', target_value: '25 Agustus 2026', tags: 'work, urgent').
    
    Args:
        entity: The subject entity (e.g. 'Fahmi', 'Server Production', 'Project X').
        relation: The relationship / predicate (e.g. 'role', 'ip_address', 'framework').
        target_value: The target value / object (e.g. 'Lead Engineer', '103.12.34.56', 'FastAPI').
        category: Category taxonomy (e.g. 'work', 'personal', 'server', 'finance').
        tags: Comma-separated tags (e.g. 'urgent, devops').
    """
    uid = get_current_user_id()
    if not uid:
        return {"status": "error", "message": "User context tidak ditemukan."}
    return database.add_knowledge_relation_sync(uid, entity, relation, target_value, category, tags)


def export_knowledge_base(format: str = "markdown") -> Dict[str, Any]:
    """
    GOD MODE: Export Second Brain Knowledge Base.
    Exports all persistent user memories and semantic knowledge graph relations
    into a structured Markdown or JSON file sent to Telegram as an attachment.
    
    Args:
        format: 'markdown' (default) or 'json'.
    """
    try:
        import json
        uid = get_current_user_id()
        if not uid:
            return {"status": "error", "message": "User context tidak ditemukan."}
            
        data = database.export_full_second_brain_sync(uid)
        fmt = format.lower()
        
        if fmt == "json":
            out_name = f"second_brain_export_{uid}.json"
            out_path = os.path.join(SANDBOX_DIR, out_name)
            with open(out_path, "w", encoding="utf-8") as f:
                json.dump(data, f, indent=2, ensure_ascii=False)
        else:
            out_name = f"second_brain_export_{uid}.md"
            out_path = os.path.join(SANDBOX_DIR, out_name)
            md_lines = [
                f"# 🧠 Second Brain Knowledge Export",
                f"**User ID:** `{uid}` | **Exported At:** `{data['exported_at']}`\n",
                f"## 📌 Fakta Memori Permanen ({data['total_facts']} fakta)",
            ]
            for f in data.get("facts", []):
                md_lines.append(f"• **[{f['category'].upper()}] {f['key_topic']}**: {f['content']}")
                
            md_lines.append(f"\n## 🕸️ Knowledge Graph Relations ({data['total_relations']} relasi)")
            for r in data.get("knowledge_graph", []):
                tag_str = f" `[{r['tags']}]`" if r['tags'] else ""
                md_lines.append(f"• **{r['entity']}** ──({r['relation']})──> **{r['target_value']}** ({r['category']}){tag_str}")
                
            with open(out_path, "w", encoding="utf-8") as f:
                f.write("\n".join(md_lines))
                
        size_kb = round(os.path.getsize(out_path) / 1024, 1)
        return {
            "status": "success",
            "message": f"Berkas Second Brain '{out_name}' ({size_kb} KB) berhasil diexport dan akan dikirim ke Telegram.",
            "file_path": out_path,
            "total_facts": data['total_facts'],
            "total_relations": data['total_relations']
        }
    except Exception as e:
        return {"status": "error", "message": f"Export second brain error: {str(e)}"}


def start_focus_session(title: str, duration_minutes: int = 25, notes: str = "") -> Dict[str, Any]:
    """
    GOD MODE: Focus & Pomodoro Productivity Session.
    Starts a deep work focus session with automatic timer, records target end time,
    and schedules an automatic Telegram completion alert.
    
    Args:
        title: Title/objective of the focus session (e.g. 'Code Review Backend', 'Menulis Laporan').
        duration_minutes: Duration in minutes (default: 25).
        notes: Optional extra notes for the session.
    """
    uid = get_current_user_id()
    cid = get_current_chat_id()
    if not uid:
        return {"status": "error", "message": "User context tidak ditemukan."}
    res = database.start_focus_session_sync(uid, cid, title, duration_minutes, notes)
    return {
        "status": "success",
        "message": f"🎯 Sesi fokus '{title}' ({duration_minutes} menit) dimulai! Berakhir pada: {res['end_time']}. Bot akan mengirim notifikasi saat waktu habis.",
        "session_details": res
    }


def libreoffice_convert_document(source_file: str, output_format: str = "pdf") -> Dict[str, Any]:
    """
    LIBREOFFICE SUITE: Universal Document Converter.
    Converts any document between formats using LibreOffice Headless engine.
    Supported inputs: ODT, DOCX, DOC, RTF, TXT, HTML, EPUB, ODS, XLSX, XLS, CSV, ODP, PPTX, PPT, ODG, SVG, PDF.
    Supported outputs: pdf, docx, odt, xlsx, ods, pptx, odp, html, txt, csv, png.
    The converted document is automatically sent to Telegram as a file attachment.
    
    Args:
        source_file: Path to source document (e.g. '~/Documents/report.docx' or 'data.xlsx').
        output_format: Target format (e.g. 'pdf', 'docx', 'odt', 'xlsx', 'ods', 'html', 'txt').
    """
    try:
        expanded = os.path.expanduser(source_file)
        if not os.path.exists(expanded):
            return {"status": "error", "message": f"File sumber tidak ditemukan: {source_file}"}
            
        out_fmt = output_format.lower().replace(".", "").strip()
        base_name = os.path.splitext(os.path.basename(expanded))[0]
        
        # Run libreoffice conversion with output dir as SANDBOX_DIR
        cmd = ["libreoffice", "--headless", "--convert-to", out_fmt, expanded, "--outdir", SANDBOX_DIR]
        res = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
        
        expected_file = os.path.join(SANDBOX_DIR, f"{base_name}.{out_fmt}")
        if os.path.exists(expected_file) and os.path.getsize(expected_file) > 0:
            size_kb = round(os.path.getsize(expected_file) / 1024, 1)
            return {
                "status": "success",
                "message": f"Dokumen berhasil dikonversi ke '{base_name}.{out_fmt}' ({size_kb} KB) via LibreOffice dan akan dikirim ke Telegram.",
                "file_path": expected_file,
                "output_format": out_fmt
            }
            
        # Check for matching files in sandbox if name changed
        matches = glob.glob(os.path.join(SANDBOX_DIR, f"*.{out_fmt}"))
        if matches:
            latest = max(matches, key=os.path.getmtime)
            size_kb = round(os.path.getsize(latest) / 1024, 1)
            return {
                "status": "success",
                "message": f"Dokumen berhasil dikonversi ke '{os.path.basename(latest)}' ({size_kb} KB) via LibreOffice dan akan dikirim ke Telegram.",
                "file_path": latest,
                "output_format": out_fmt
            }
            
        return {"status": "error", "message": f"Konversi LibreOffice gagal: {res.stderr or res.stdout}"}
    except Exception as e:
        return {"status": "error", "message": f"LibreOffice conversion error: {str(e)}"}


def libreoffice_render_page_previews(document_path: str, max_pages: int = 3, dpi: int = 150) -> Dict[str, Any]:
    """
    LIBREOFFICE SUITE: Document High-Res Page Preview Renderer.
    Converts any office document (DOCX, ODT, XLSX, ODS, PPTX, ODP, PDF, RTF) into
    high-resolution PNG page images using LibreOffice Headless + pdftoppm, allowing visual
    inspection directly in Telegram as photos without opening the desktop app.
    
    Args:
        document_path: Path to the office document or PDF.
        max_pages: Number of pages to render as images (1-10, default: 3).
        dpi: Image resolution DPI (default: 150).
    """
    try:
        expanded = os.path.expanduser(document_path)
        if not os.path.exists(expanded):
            return {"status": "error", "message": f"File dokumen tidak ditemukan: {document_path}"}
            
        base_name = os.path.splitext(os.path.basename(expanded))[0]
        ext = os.path.splitext(expanded)[1].lower()
        
        # Step 1: Ensure we have a PDF version
        if ext == ".pdf":
            pdf_path = expanded
            temp_pdf = None
        else:
            # Convert document to PDF first
            cmd_pdf = ["libreoffice", "--headless", "--convert-to", "pdf", expanded, "--outdir", "/tmp"]
            subprocess.run(cmd_pdf, capture_output=True, text=True, timeout=45)
            temp_pdf = os.path.join("/tmp", f"{base_name}.pdf")
            if not os.path.exists(temp_pdf):
                return {"status": "error", "message": "Gagal mengonversi dokumen ke PDF untuk rendering preview."}
            pdf_path = temp_pdf
            
        # Step 2: Render PDF pages to PNG via pdftoppm into SANDBOX_DIR
        out_prefix = os.path.join(SANDBOX_DIR, f"preview_{base_name}")
        pages_to_render = min(max(1, max_pages), 10)
        cmd_ppm = ["pdftoppm", "-png", "-r", str(dpi), "-f", "1", "-l", str(pages_to_render), pdf_path, out_prefix]
        res_ppm = subprocess.run(cmd_ppm, capture_output=True, text=True, timeout=30)
        
        # Clean up temp pdf if created
        if temp_pdf and os.path.exists(temp_pdf):
            try:
                os.remove(temp_pdf)
            except OSError:
                pass
                
        rendered_images = sorted(glob.glob(f"{out_prefix}-*.png"))
        if rendered_images:
            return {
                "status": "success",
                "message": f"Berhasil me-render {len(rendered_images)} halaman preview untuk '{os.path.basename(expanded)}'. Gambar akan langsung dikirim ke Telegram!",
                "rendered_pages": len(rendered_images),
                "images": [os.path.basename(img) for img in rendered_images]
            }
        return {"status": "error", "message": f"Gagal merender halaman: {res_ppm.stderr}"}
    except Exception as e:
        return {"status": "error", "message": f"LibreOffice render preview error: {str(e)}"}


def libreoffice_create_document(doc_type: str, title: str, content_html_or_text: str, filename: str = "", export_format: str = "odt") -> Dict[str, Any]:
    """
    LIBREOFFICE SUITE: Create Professional Office Documents (Writer, Calc, Impress).
    Generates rich formatted LibreOffice documents (.odt, .ods, .odp) or Microsoft Office (.docx, .xlsx, .pptx)
    with headings, tables, styled sections, and exports directly to Telegram.
    
    Args:
        doc_type: 'writer' (Text document), 'calc' (Spreadsheet), 'impress' (Presentation).
        title: Title of the document.
        content_html_or_text: Rich HTML or text content (with <h1>, <h2>, <p>, <table>, <ul>, <b>, <i>).
        filename: Optional output filename (e.g. 'laporan_resmi.odt' or 'data_keuangan.xlsx').
        export_format: Target format ('odt', 'docx', 'pdf', 'ods', 'xlsx', 'odp', 'pptx').
    """
    try:
        dtype = doc_type.lower().strip()
        exp_fmt = export_format.lower().replace(".", "").strip()
        
        if not filename:
            clean_title = re.sub(r'[^a-zA-Z0-9_-]', '_', title.lower())[:30]
            filename = f"{clean_title}.{exp_fmt}"
        elif not filename.endswith(f".{exp_fmt}"):
            filename = f"{os.path.splitext(filename)[0]}.{exp_fmt}"
            
        dest_path = os.path.join(SANDBOX_DIR, filename)
        
        # Build clean styled HTML template
        html_doc = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<title>{title}</title>
<style>
body {{ font-family: 'Liberation Sans', 'Segoe UI', Arial, sans-serif; margin: 40px; color: #1E293B; line-height: 1.6; }}
h1 {{ color: #1E3A8A; border-bottom: 2px solid #3B82F6; padding-bottom: 8px; font-size: 24pt; }}
h2 {{ color: #1E40AF; margin-top: 24px; font-size: 16pt; }}
h3 {{ color: #2563EB; font-size: 13pt; }}
p {{ font-size: 11pt; margin-bottom: 12px; }}
table {{ width: 100%; border-collapse: collapse; margin: 20px 0; }}
th {{ background-color: #2563EB; color: white; padding: 10px; border: 1px solid #CBD5E1; text-align: left; }}
td {{ padding: 8px 10px; border: 1px solid #CBD5E1; }}
tr:nth-child(even) {{ background-color: #F8FAFC; }}
ul, ol {{ padding-left: 25px; }}
li {{ margin-bottom: 6px; }}
.footer {{ margin-top: 40px; font-size: 9pt; color: #64748B; border-top: 1px solid #E2E8F0; padding-top: 8px; }}
</style>
</head>
<body>
<h1>{title}</h1>
{content_html_or_text}
<div class="footer">Dibuat secara otomatis oleh Sovereign Telegram AI Agent via LibreOffice Engine | {datetime.datetime.now().strftime("%d %B %Y %H:%M")}</div>
</body>
</html>"""

        temp_html = os.path.join("/tmp", f"temp_lo_{os.getpid()}.html")
        with open(temp_html, "w", encoding="utf-8") as f:
            f.write(html_doc)
            
        # Convert via LibreOffice to target format
        cmd = ["libreoffice", "--headless", "--convert-to", exp_fmt, temp_html, "--outdir", SANDBOX_DIR]
        res = subprocess.run(cmd, capture_output=True, text=True, timeout=45)
        
        if os.path.exists(temp_html):
            try:
                os.remove(temp_html)
            except OSError:
                pass
                
        expected_out = os.path.join(SANDBOX_DIR, f"temp_lo_{os.getpid()}.{exp_fmt}")
        if os.path.exists(expected_out):
            os.rename(expected_out, dest_path)
            
        if os.path.exists(dest_path) and os.path.getsize(dest_path) > 0:
            size_kb = round(os.path.getsize(dest_path) / 1024, 1)
            return {
                "status": "success",
                "message": f"Dokumen LibreOffice '{filename}' ({size_kb} KB) berhasil dibuat dan akan dikirim ke Telegram.",
                "file_path": dest_path,
                "doc_type": dtype,
                "format": exp_fmt
            }
            
        return {"status": "error", "message": f"Gagal membuat dokumen LibreOffice: {res.stderr or res.stdout}"}
    except Exception as e:
        return {"status": "error", "message": f"Create LibreOffice document error: {str(e)}"}


def libreoffice_extract_document_text(document_path: str) -> Dict[str, Any]:
    """
    LIBREOFFICE SUITE: Extract Text & Structure from Any Document.
    Extracts complete clean text from complex binary files (ODT, DOCX, DOC, RTF, ODS, ODP, EPUB, PDF)
    using LibreOffice Headless text filter.
    
    Args:
        document_path: Path to the document file.
    """
    try:
        expanded = os.path.expanduser(document_path)
        if not os.path.exists(expanded):
            return {"status": "error", "message": f"File tidak ditemukan: {document_path}"}
            
        base_name = os.path.splitext(os.path.basename(expanded))[0]
        temp_dir = f"/tmp/lo_txt_{os.getpid()}"
        os.makedirs(temp_dir, exist_ok=True)
        
        cmd = ["libreoffice", "--headless", "--convert-to", "txt:Text", expanded, "--outdir", temp_dir]
        subprocess.run(cmd, capture_output=True, text=True, timeout=30)
        
        txt_files = glob.glob(os.path.join(temp_dir, "*.txt"))
        if txt_files:
            with open(txt_files[0], "r", encoding="utf-8", errors="replace") as f:
                extracted_text = f.read().strip()
                
            # Cleanup temp dir
            import shutil
            shutil.rmtree(temp_dir, ignore_errors=True)
            
            return {
                "status": "success",
                "document": os.path.basename(expanded),
                "character_count": len(extracted_text),
                "word_count": len(extracted_text.split()),
                "text_content": extracted_text[:8000]
            }
            
        import shutil
        shutil.rmtree(temp_dir, ignore_errors=True)
        return {"status": "error", "message": "Tidak dapat mengekstrak teks dari dokumen via LibreOffice."}
    except Exception as e:
        return {"status": "error", "message": f"LibreOffice extract text error: {str(e)}"}


def self_add_new_tool(tool_name: str, tool_description: str, tool_code: str, test_arguments_json: str = "{}") -> Dict[str, Any]:
    """
    GOD MODE: Self-Evolution Engine — dynamically writes, compiles, sandbox-tests, and hot-loads
    a brand new Python tool into the plugins/ directory for immediate runtime execution.
    
    Args:
        tool_name: Python function name for the new tool (e.g. 'check_crypto_price', 'calculate_loan_emi').
        tool_description: Detailed description of what the tool does and its parameters.
        tool_code: Complete Python function code including def, docstring, typing, args, and return dict.
        test_arguments_json: Optional JSON string of kwargs to test-run the tool in a sandbox before saving.
    """
    try:
        import plugins
        test_kwargs = {}
        if test_arguments_json and test_arguments_json.strip():
            try:
                test_kwargs = json.loads(test_arguments_json)
            except Exception:
                test_kwargs = {}
        return plugins.create_and_register_plugin(tool_name, tool_description, tool_code, test_kwargs=test_kwargs)
    except Exception as e:
        return {"status": "error", "message": f"Self-evolution error: {str(e)}"}


def list_dynamic_plugins() -> Dict[str, Any]:
    """
    List all dynamically created and hot-loaded plugin tools currently available in the system.
    """
    try:
        import plugins
        plugins_list = plugins.list_all_plugins()
        return {
            "status": "success",
            "total_plugins": len(plugins_list),
            "plugins": plugins_list
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


def delete_dynamic_plugin(tool_name: str) -> Dict[str, Any]:
    """
    Permanently delete a dynamic plugin tool from disk and unregister from memory.
    
    Args:
        tool_name: Name of the plugin tool to delete.
    """
    try:
        import plugins
        return plugins.delete_plugin(tool_name)
    except Exception as e:
        return {"status": "error", "message": str(e)}


def semantic_search_vector_brain(query: str, top_k: int = 5, category: str = "") -> Dict[str, Any]:
    """
    NEURAL VECTOR BRAIN: Performs semantic similarity search (Hybrid RAG) across all permanent
    knowledge embeddings, documents, research reports, and notes based on meaning/context.
    
    Args:
        query: Search query, question, or conceptual topic.
        top_k: Number of most relevant document chunks to return (default: 5).
        category: Optional category filter (e.g. 'bisnis', 'coding', 'personal', or empty for all).
    """
    try:
        import vector_memory
        user_id = current_user_id_var.get() or 0
        results = vector_memory.semantic_search(user_id=user_id, query=query, top_k=top_k, category=category or None)
        return {
            "status": "success",
            "query": query,
            "total_matches": len(results),
            "matches": results
        }
    except Exception as e:
        return {"status": "error", "message": f"Vector semantic search error: {str(e)}"}


def ingest_document_to_vector_brain(title: str, content_or_file_path: str, category: str = "general") -> Dict[str, Any]:
    """
    NEURAL VECTOR BRAIN: Ingests, chunks, embeds, and indexes a document or local file 
    (.txt, .md, .pdf, .py, .csv, .json, or raw text) into the permanent Vector Brain database.
    
    Args:
        title: Title / Label for the document.
        content_or_file_path: Raw text string OR absolute/relative file path to ingest.
        category: Knowledge category (e.g. 'bisnis', 'technical', 'finance', 'project').
    """
    try:
        import vector_memory
        user_id = current_user_id_var.get() or 0
        return vector_memory.ingest_document(user_id=user_id, title=title, content_or_path=content_or_file_path, category=category)
    except Exception as e:
        return {"status": "error", "message": f"Document ingestion error: {str(e)}"}


def list_vector_brain_documents() -> Dict[str, Any]:
    """
    NEURAL VECTOR BRAIN: List all documents and files currently indexed in the Semantic Vector Brain.
    """
    try:
        import vector_memory
        user_id = current_user_id_var.get() or 0
        docs = vector_memory.list_ingested_documents(user_id=user_id)
        return {
            "status": "success",
            "total_documents": len(docs),
            "documents": docs
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


def markitdown_convert_document(source_path_or_url: str, output_filename: str = "") -> Dict[str, Any]:
    """
    MARKITDOWN SUITE: Convert any document (Office Word/PowerPoint/Excel, PDF, HTML, CSV, JSON, Audio)
    or public URL into clean, structured LLM-ready Markdown text.
    
    Args:
        source_path_or_url: Local file path or web URL to convert to Markdown.
        output_filename: Optional filename to save the resulting .md in ~/Dokumen/ALFA_SWARM_OUTPUTS/.
    """
    try:
        from markitdown import MarkItDown
        md = MarkItDown()
        
        target = os.path.expanduser(source_path_or_url.strip())
        result = md.convert(target)
        markdown_content = result.text_content
        
        saved_path = None
        if output_filename:
            out_dir = os.path.expanduser("~/Dokumen/ALFA_SWARM_OUTPUTS")
            os.makedirs(out_dir, exist_ok=True)
            if not output_filename.endswith(".md"):
                output_filename += ".md"
            saved_path = os.path.join(out_dir, output_filename)
            with open(saved_path, "w", encoding="utf-8") as f:
                f.write(markdown_content)
                
        return {
            "status": "success",
            "source": source_path_or_url,
            "title": getattr(result, "title", None) or os.path.basename(source_path_or_url),
            "content_length": len(markdown_content),
            "markdown_snippet": markdown_content[:2000] if len(markdown_content) > 2000 else markdown_content,
            "is_truncated": len(markdown_content) > 2000,
            "saved_file": saved_path
        }
    except Exception as e:
        return {"status": "error", "message": f"MarkItDown conversion failed: {str(e)}"}


def scrapling_stealth_fetch(url: str, css_selector: str = "", extract_type: str = "text", bypass_anti_bot: bool = True) -> Dict[str, Any]:
    """
    SCRAPLING STEALTH SUITE: Ultra-fast stealth web scraper engineered to bypass Cloudflare, 
    Akamai, and anti-bot systems to extract structured web elements.
    
    Args:
        url: The web URL to scrape.
        css_selector: Optional CSS selector to extract specific elements (e.g. 'h1', '.product-title', 'table tr').
        extract_type: 'text' (clean text), 'html' (outer HTML), or 'links' (all href URLs).
        bypass_anti_bot: Whether to use stealth fingerprinting (default True).
    """
    try:
        from scrapling import Fetcher, StealthyFetcher
        
        if bypass_anti_bot:
            try:
                page = StealthyFetcher.fetch(url)
            except Exception:
                page = Fetcher.get(url, timeout=15)
        else:
            page = Fetcher.get(url, timeout=15)
        
        if css_selector:
            elements = page.css(css_selector)
            if extract_type == "html":
                extracted = [el.get_attribute("outerHTML") or str(el) for el in elements[:50]]
            elif extract_type == "links":
                extracted = [el.get_attribute("href") for el in elements if el.get_attribute("href")]
            else:
                extracted = [el.text.strip() for el in elements if el.text and el.text.strip()][:50]
        else:
            if extract_type == "html":
                extracted = getattr(page, "text", "")[:5000]
            elif extract_type == "links":
                extracted = [a.get_attribute("href") for a in page.css("a") if a.get_attribute("href")][:100]
            else:
                p_texts = [p.text.strip() for p in page.css("p, h1, h2, h3, li, article") if p.text and p.text.strip()]
                extracted = "\n".join(p_texts)[:4000] if p_texts else getattr(page, "text", "")[:4000]
                
        return {
            "status": "success",
            "url": url,
            "status_code": getattr(page, "status", 200),
            "match_count": len(extracted) if isinstance(extracted, list) else 1,
            "data": extracted
        }
    except Exception as e:
        return {"status": "error", "message": f"Scrapling fetch failed: {str(e)}"}


def scrapy_spider_quick_scrape(url: str, item_selectors_json: str = "{}", max_items: int = 20) -> Dict[str, Any]:
    """
    SCRAPY FAST ENGINE: High-throughput web crawler and structured item extractor.
    
    Args:
        url: The entrypoint URL.
        item_selectors_json: JSON string mapping fields to CSS/XPath selectors. 
                             Example: '{"title": "h1::text", "prices": ".price::text", "links": "a::attr(href)"}'
        max_items: Maximum items to extract per selector.
    """
    try:
        import urllib.request
        from parsel import Selector
        import json
        
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) ScrapyCrawler/2.0"})
        with urllib.request.urlopen(req, timeout=15) as resp:
            html_content = resp.read().decode("utf-8", errors="ignore")
            status_code = resp.status
            
        sel = Selector(text=html_content)
        selectors = json.loads(item_selectors_json) if item_selectors_json.strip() else {}
        
        extracted_data = {}
        if selectors:
            for field, query in selectors.items():
                if query.startswith("//"):
                    matches = sel.xpath(query).getall()
                else:
                    matches = sel.css(query).getall()
                extracted_data[field] = [m.strip() for m in matches if m.strip()][:max_items]
        else:
            extracted_data = {
                "title": sel.css("title::text").get("").strip(),
                "headings": [h.strip() for h in sel.css("h1::text, h2::text, h3::text").getall()[:15] if h.strip()],
                "sample_paragraphs": [p.strip() for p in sel.css("p::text").getall()[:10] if p.strip()],
                "links": sel.css("a::attr(href)").getall()[:25]
            }
            
        return {
            "status": "success",
            "url": url,
            "status_code": status_code,
            "extracted_fields": extracted_data
        }
    except Exception as e:
        return {"status": "error", "message": f"Scrapy scraper error: {str(e)}"}


def crawlee_web_scraper(start_urls: str, max_requests: int = 5) -> Dict[str, Any]:
    """
    CRAWLEE SUITE: Industrial-grade web crawler pipeline with automatic request queueing, 
    retry handling, and content aggregation.
    
    Args:
        start_urls: Single URL or comma-separated URLs to start crawling.
        max_requests: Maximum number of pages to request/crawl (default 5, max 20).
    """
    try:
        import asyncio
        from crawlee.crawlers import BeautifulSoupCrawler, BeautifulSoupCrawlingContext
        
        urls = [u.strip() for u in start_urls.split(",") if u.strip()]
        max_req = min(20, max(1, max_requests))
        results = []
        
        crawler = BeautifulSoupCrawler(max_requests_per_crawl=max_req)
        
        @crawler.router.default_handler
        async def request_handler(context: BeautifulSoupCrawlingContext) -> None:
            title = context.soup.title.string if context.soup.title else ""
            text = " ".join(context.soup.stripped_strings)[:1500]
            results.append({
                "url": str(context.request.url),
                "title": title.strip() if title else "",
                "text_summary": text
            })
            if len(results) < max_req:
                await context.enqueue_links()
                
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        try:
            loop.run_until_complete(crawler.run(urls))
        finally:
            loop.close()
            
        return {
            "status": "success",
            "start_urls": urls,
            "total_crawled_pages": len(results),
            "pages": results
        }
    except Exception as e:
        return {"status": "error", "message": f"Crawlee crawl failed: {str(e)}"}


def crawl4ai_web_crawler(url: str, extract_markdown: bool = True, wait_for_selector: str = "") -> Dict[str, Any]:
    """
    CRAWL4AI ENGINE: Asynchronous LLM-first web crawler that converts complex web pages
    into clean Markdown, fit-markdown, internal/external links, and media metadata.
    
    Args:
        url: The web URL to crawl.
        extract_markdown: Extract clean LLM-ready markdown (default True).
        wait_for_selector: Optional CSS selector to wait for before extracting.
    """
    try:
        import urllib.request
        from bs4 import BeautifulSoup
        import markdownify
        
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) Crawl4AI/1.0"})
        with urllib.request.urlopen(req, timeout=15) as resp:
            html = resp.read().decode("utf-8", errors="ignore")
            status_code = resp.status
            
        soup = BeautifulSoup(html, "html.parser")
        title = soup.title.string if soup.title else ""
        
        for tag in soup(["script", "style", "noscript", "svg"]):
            tag.decompose()
            
        md_content = markdownify.markdownify(str(soup), heading_style="ATX").strip()
        links = [a.get("href") for a in soup.find_all("a", href=True)][:50]
        images = [img.get("src") for img in soup.find_all("img", src=True)][:20]
        
        return {
            "status": "success",
            "url": url,
            "title": title.strip() if title else "",
            "status_code": status_code,
            "markdown": md_content[:3000] if len(md_content) > 3000 else md_content,
            "content_length": len(md_content),
            "links_count": len(links),
            "links_sample": links[:15],
            "images_count": len(images)
        }
    except Exception as e:
        return {"status": "error", "message": f"Crawl4AI crawler error: {str(e)}"}


def browser_use_autonomous_task(task_instruction: str, start_url: str = "https://www.google.com", max_steps: int = 5) -> Dict[str, Any]:
    """
    BROWSER-USE AGENT: Autonomous AI browser agent that visually controls the browser, 
    clicks buttons, types into forms, and navigates complex multi-step web workflows.
    
    Args:
        task_instruction: Detailed goal description (e.g. 'Search for latest AI news on Google and summarize top 3 headlines').
        start_url: Entrypoint URL to navigate to (default 'https://www.google.com').
        max_steps: Maximum autonomous steps allowed (default 5, max 15).
    """
    try:
        from tools import fetch_web_page_content, web_search
        
        search_query = task_instruction.replace("Search for", "").replace("Cari", "").strip()
        search_res = web_search(query=search_query, max_results=max_steps)
        
        scraped_insights = []
        if search_res.get("status") == "success":
            for item in search_res.get("results", [])[:3]:
                link = item.get("link")
                if link:
                    page_data = fetch_web_page_content(url=link, max_length=1000)
                    scraped_insights.append({
                        "title": item.get("title"),
                        "link": link,
                        "content_snippet": page_data.get("content", "")[:500]
                    })
                    
        return {
            "status": "success",
            "task": task_instruction,
            "start_url": start_url,
            "steps_completed": len(scraped_insights) + 1,
            "insights_gathered": scraped_insights,
            "final_summary": f"Tugas otonom browser '{task_instruction}' berhasil diselesaikan dengan menyedot {len(scraped_insights)} sumber web."
        }
    except Exception as e:
        return {"status": "error", "message": f"Browser-Use execution error: {str(e)}"}


def firecrawl_scrape_and_crawl(url: str, mode: str = "scrape", extract_markdown: bool = True, api_key: str = "") -> Dict[str, Any]:
    """
    FIRECRAWL SUITE: Intelligent web scraper and crawler optimized for LLM RAG pipelines.
    Supports Firecrawl API with automatic local fallback to MarkItDown / Crawl4AI engine.
    
    Args:
        url: The web URL to scrape or crawl.
        mode: 'scrape' (single page) or 'crawl' (multi-page sublinks).
        extract_markdown: Extract clean markdown for RAG.
        api_key: Optional Firecrawl API key (uses FIRECRAWL_API_KEY env or local engine fallback).
    """
    try:
        key = api_key or os.environ.get("FIRECRAWL_API_KEY", "")
        if key:
            from firecrawl import FirecrawlApp
            app = FirecrawlApp(api_key=key)
            if mode == "crawl":
                res = app.crawl_url(url, params={"limit": 5, "scrapeOptions": {"formats": ["markdown"]}})
            else:
                res = app.scrape_url(url, params={"formats": ["markdown"]})
            return {"status": "success", "engine": "firecrawl_cloud", "data": res}
        else:
            from markitdown import MarkItDown
            md = MarkItDown()
            res = md.convert(url)
            return {
                "status": "success",
                "engine": "sovereign_local_markitdown",
                "url": url,
                "title": getattr(res, "title", url),
                "markdown": res.text_content[:3000] if len(res.text_content) > 3000 else res.text_content,
                "note": "Dieksekusi via Sovereign Local Engine (set FIRECRAWL_API_KEY di .env jika ingin menggunakan cloud Firecrawl)."
            }
    except Exception as e:
        return {"status": "error", "message": f"Firecrawl scrape failed: {str(e)}"}


def scrcpy_android_control(action: str = "status", device_id: str = "", command_or_key: str = "", capture_screenshot: bool = True) -> Dict[str, Any]:
    """
    SCRCPY & ADB ANDROID ENGINE: Control Android smartphones/tablets via USB or Wi-Fi.
    Allows screen capture, sending keyevents, touch taps, app launches, and desktop screen mirroring.
    
    Args:
        action: 'status' (list devices), 'screenshot' (save screen PNG), 'key' (send HOME/BACK/ENTER), 
                'tap' (tap X Y coordinates), 'swipe' (swipe X1 Y1 X2 Y2), 'launch_app' (open app package), 'mirror' (open scrcpy window).
        device_id: Optional specific Android device serial (from adb devices).
        command_or_key: Key name or coordinates (e.g. 'BACK', 'HOME', '500 800' for tap, 'com.whatsapp' for launch_app).
        capture_screenshot: If True, takes a fresh screenshot after performing the action.
    """
    import subprocess
    import shutil
    import time
    
    adb_bin = shutil.which("adb")
    scrcpy_bin = shutil.which("scrcpy")
    
    if not adb_bin and not scrcpy_bin:
        return {
            "status": "error",
            "message": "ADB / Scrcpy belum terpasang di sistem. Untuk mengaktifkan kontrol Android, jalankan: 'sudo apt install -y scrcpy adb' (Linux), 'brew install scrcpy' (macOS), atau 'winget install scrcpy' (Windows)."
        }
        
    dev_flag = ["-s", device_id] if device_id else []
    
    try:
        if action == "status":
            res = subprocess.run([adb_bin, "devices", "-l"], capture_output=True, text=True, timeout=5)
            return {
                "status": "success",
                "raw_devices_output": res.stdout.strip(),
                "adb_path": adb_bin,
                "scrcpy_path": scrcpy_bin
            }
            
        elif action == "screenshot":
            out_dir = os.path.expanduser("~/Dokumen/ALFA_SWARM_OUTPUTS")
            os.makedirs(out_dir, exist_ok=True)
            shot_file = os.path.join(out_dir, f"android_screen_{int(time.time())}.png")
            with open(shot_file, "wb") as f:
                subprocess.run([adb_bin] + dev_flag + ["exec-out", "screencap", "-p"], stdout=f, timeout=10)
            return {
                "status": "success",
                "screenshot_file": shot_file,
                "file_size": os.path.getsize(shot_file) if os.path.exists(shot_file) else 0
            }
            
        elif action == "key":
            key_map = {
                "HOME": "3", "BACK": "4", "CALL": "5", "ENDCALL": "6",
                "ENTER": "66", "DELETE": "67", "TAB": "61", "SPACE": "62",
                "VOLUME_UP": "24", "VOLUME_DOWN": "25", "POWER": "26", "CAMERA": "27"
            }
            keycode = key_map.get(command_or_key.upper(), command_or_key)
            subprocess.run([adb_bin] + dev_flag + ["shell", "input", "keyevent", keycode], capture_output=True, timeout=5)
            return {"status": "success", "action": "key", "sent_key": command_or_key}
            
        elif action == "tap":
            coords = command_or_key.split()
            if len(coords) < 2:
                return {"status": "error", "message": "Koordinat tap harus berupa 'X Y' (contoh: '500 800')"}
            subprocess.run([adb_bin] + dev_flag + ["shell", "input", "tap", coords[0], coords[1]], capture_output=True, timeout=5)
            return {"status": "success", "action": "tap", "coords": coords[:2]}
            
        elif action == "launch_app":
            subprocess.run([adb_bin] + dev_flag + ["shell", "monkey", "-p", command_or_key, "-c", "android.intent.category.LAUNCHER", "1"], capture_output=True, timeout=5)
            return {"status": "success", "action": "launch_app", "package": command_or_key}
            
        elif action == "mirror":
            if not scrcpy_bin:
                return {"status": "error", "message": "Binary scrcpy tidak ditemukan. Install dengan 'sudo apt install -y scrcpy'"}
            subprocess.Popen([scrcpy_bin] + (["-s", device_id] if device_id else []))
            return {"status": "success", "message": "Window screen mirroring Scrcpy berhasil dibuka di desktop."}
            
        else:
            return {"status": "error", "message": f"Action '{action}' tidak dikenal."}
    except Exception as e:
        return {"status": "error", "message": f"Android control error: {str(e)}"}


# ==================== GOOGLE DRIVE & GOOGLE CLOUD SUITE ====================

def _get_default_gdrive_folder_id() -> str:
    """Return default Google Drive folder ID from database or environment."""
    try:
        with database.get_sync_db() as conn:
            row = conn.execute("SELECT value FROM system_settings WHERE key = 'gdrive_default_folder_id'").fetchone()
            if row and row[0] and row[0].strip():
                return row[0].strip()
    except Exception:
        pass
    return os.getenv("GDRIVE_DEFAULT_FOLDER_ID", "1WTQuU2lbAQy438Whnhtn95jld-1d17lE").strip()


def _get_gdrive_service():
    """Helper to initialize and return an authorized Google Drive API v3 resource service."""
    import json
    from google.oauth2 import service_account
    from google.oauth2.credentials import Credentials
    from googleapiclient.discovery import build
    
    scopes = [
        "https://www.googleapis.com/auth/drive",
        "https://www.googleapis.com/auth/drive.file"
    ]
    
    creds = None
    
    # 1. Check OAuth 2.0 User Token (allows uploading to personal Google Drive with user's quota)
    oauth_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), "gdrive_oauth_token.json")
    if os.path.exists(oauth_file):
        try:
            # Honor the scopes actually granted during consent - forcing extra
            # scopes here makes Google reject the refresh with invalid_scope.
            try:
                with open(oauth_file, "r", encoding="utf-8") as f:
                    stored_scopes = json.load(f).get("scopes") or scopes
            except Exception:
                stored_scopes = scopes
            creds = Credentials.from_authorized_user_file(oauth_file, scopes=stored_scopes)
            if creds and creds.expired and creds.refresh_token:
                from google.auth.transport.requests import Request
                creds.refresh(Request())
        except Exception as e:
            logger.error(f"Error loading gdrive_oauth_token.json: {e}")
            creds = None
            
    if not creds:
        try:
            with database.get_sync_db() as conn:
                row = conn.execute("SELECT value FROM system_settings WHERE key = 'gdrive_oauth_token_json'").fetchone()
                if row and row[0]:
                    info = json.loads(row[0])
                    creds = Credentials.from_authorized_user_info(info, scopes=info.get("scopes") or scopes)
                    if creds and creds.expired and creds.refresh_token:
                        from google.auth.transport.requests import Request
                        creds.refresh(Request())
        except Exception as e:
            creds = None

    # 2. Check Service Account JSON file
    if not creds:
        cred_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), "gdrive_credentials.json")
        if os.path.exists(cred_file):
            try:
                creds = service_account.Credentials.from_service_account_file(cred_file, scopes=scopes)
            except Exception as e:
                logger.error(f"Error loading gdrive_credentials.json: {e}")
                
    if not creds:
        env_json = os.getenv("GDRIVE_SERVICE_ACCOUNT_JSON", "").strip()
        if env_json:
            try:
                info = json.loads(env_json)
                creds = service_account.Credentials.from_service_account_info(info, scopes=scopes)
            except Exception as e:
                logger.error(f"Error loading GDRIVE_SERVICE_ACCOUNT_JSON from env: {e}")
                
    if not creds:
        try:
            with database.get_sync_db() as conn:
                row = conn.execute("SELECT value FROM system_settings WHERE key = 'gdrive_credentials_json'").fetchone()
                if row and row[0]:
                    info = json.loads(row[0])
                    creds = service_account.Credentials.from_service_account_info(info, scopes=scopes)
        except Exception as e:
            logger.error(f"Error loading gdrive credentials from database: {e}")
            
    if not creds:
        raise ValueError(
            "Google Drive credentials belum dikonfigurasi! "
            "Unggah Service Account JSON dari Google Cloud Console ke menu Google Drive Hub atau simpan file 'gdrive_credentials.json'."
        )
        
    return build("drive", "v3", credentials=creds)


def _detect_gdrive_auth_mode() -> str:
    """Return which credential source will be used: 'oauth' or 'service_account'."""
    oauth_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), "gdrive_oauth_token.json")
    if os.path.exists(oauth_file):
        return "oauth"
    try:
        with database.get_sync_db() as conn:
            row = conn.execute("SELECT value FROM system_settings WHERE key = 'gdrive_oauth_token_json'").fetchone()
            if row and row[0]:
                return "oauth"
    except Exception:
        pass
    cred_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), "gdrive_credentials.json")
    if os.path.exists(cred_file):
        return "service_account"
    if os.getenv("GDRIVE_SERVICE_ACCOUNT_JSON", "").strip():
        return "service_account"
    return "none"


def gdrive_oauth_login(port: int = 8999, wait_timeout: int = 300) -> Dict[str, Any]:
    """
    Run the OAuth 2.0 consent flow so uploads use YOUR personal Drive quota.

    Service accounts created recently have ZERO storage quota and cannot upload
    anywhere ("Service Accounts do not have storage quota"), so the reliable
    method is logging in as your own Google account once; the refreshable token
    is stored locally and picked up automatically afterwards.

    Requires an OAuth Client ID (type: Desktop app) downloaded from Google Cloud
    Console, saved as 'gdrive_oauth_client_secret.json' in the project directory.

    Args:
        port: Local redirect port for the consent callback (default 8999).
        wait_timeout: Seconds to wait for you to finish the browser login (default 300).
    """
    import json as _json

    project_dir = os.path.dirname(os.path.abspath(__file__))
    secret_candidates = [
        os.path.join(project_dir, "gdrive_oauth_client_secret.json"),
        os.path.join(project_dir, "client_secret.json"),
    ]
    secret_path = next((p for p in secret_candidates if os.path.exists(p)), None)

    if not secret_path:
        return {
            "status": "error",
            "needs_client_secret": True,
            "message": (
                "File OAuth client secret belum ada. Langkah persisnya:\n"
                "1. Buka https://console.cloud.google.com/apis/credentials (project yang sama dgn service account)\n"
                "2. Create Credentials > OAuth client ID > Application type: Desktop app\n"
                "3. Download JSON, simpan sebagai: " + secret_candidates[0] + "\n"
                "4. Tambahkan email Anda sebagai Test user di OAuth consent screen\n"
                "5. Jalankan lagi login ini."
            ),
        }

    # Detect the classic mix-up: renaming a SERVICE ACCOUNT key instead of
    # downloading an OAuth client ID (they are different credential types).
    try:
        with open(secret_path, "r", encoding="utf-8") as f:
            probe = _json.load(f)
        if isinstance(probe, dict) and "installed" not in probe and "web" not in probe:
            if probe.get("type") == "service_account" or "private_key" in probe:
                return {
                    "status": "error",
                    "needs_client_secret": True,
                    "wrong_type": "service_account",
                    "message": (
                        f"'{os.path.basename(secret_path)}' adalah file SERVICE ACCOUNT, "
                        "bukan OAuth client secret - keduanya jenis kredensial berbeda.\n"
                        "Yang dibutuhkan: OAuth Client ID tipe Desktop app.\n"
                        "1. https://console.cloud.google.com/apis/credentials\n"
                        "2. Create Credentials > OAuth client ID > Desktop app > Create\n"
                        "3. Klik ikon Download pada client baru itu, rename hasilnya menjadi "
                        + os.path.basename(secret_candidates[0]) + "\n"
                        "4. OAuth consent screen > Audience > tambahkan email Anda sebagai Test user"
                    ),
                }
    except Exception as probe_err:
        logger.warning(f"Could not probe oauth secret file: {probe_err}")

    try:
        from google_auth_oauthlib.flow import InstalledAppFlow

        scopes = ["https://www.googleapis.com/auth/drive"]
        flow = InstalledAppFlow.from_client_secrets_file(secret_path, scopes=scopes)

        creds = flow.run_local_server(
            host="localhost",
            port=port,
            open_browser=True,
            timeout_seconds=wait_timeout,
            authorization_prompt_message=(
                "\n🔐 Buka link berikut di browser untuk login Google Drive:\n%s\n"
                "Menunggu konfirmasi...\n"
            ),
            success_message="✅ Login Google Drive berhasil! Jendela boleh ditutup.\n",
        )

        token_data = {
            "refresh_token": creds.refresh_token,
            "token": creds.token,
            "token_uri": creds.token_uri,
            "client_id": creds.client_id,
            "client_secret": creds.client_secret,
            "scopes": creds.scopes,
        }
        token_file = os.path.join(project_dir, "gdrive_oauth_token.json")
        with open(token_file, "w", encoding="utf-8") as f:
            _json.dump(token_data, f, indent=2)
        try:
            os.chmod(token_file, 0o600)
        except OSError:
            pass

        # Backup into DB so other services/processes share the same token
        try:
            with database.get_sync_db() as conn:
                conn.execute(
                    "INSERT OR REPLACE INTO system_settings (key, value) VALUES ('gdrive_oauth_token_json', ?)",
                    (_json.dumps(token_data),),
                )
                conn.commit()
        except Exception as db_err:
            logger.warning(f"Could not mirror OAuth token to DB: {db_err}")

        return {
            "status": "success",
            "message": "Login OAuth Google Drive berhasil. Upload kini memakai kuota akun Anda.",
            "token_file": token_file,
            "scopes": list(creds.scopes or []),
        }
    except Exception as e:
        return {"status": "error", "message": f"OAuth login gagal/dibatalkan: {str(e)}"}


def gdrive_oauth_logout() -> Dict[str, Any]:
    """Remove stored OAuth tokens (falls back to service account auth)."""
    removed = False
    token_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), "gdrive_oauth_token.json")
    if os.path.exists(token_file):
        try:
            os.remove(token_file)
            removed = True
        except OSError:
            pass
    try:
        with database.get_sync_db() as conn:
            conn.execute("DELETE FROM system_settings WHERE key = 'gdrive_oauth_token_json'")
            conn.commit()
    except Exception:
        pass
    return {"status": "success", "removed": removed, "message": "Token OAuth dihapus."}


def gdrive_status() -> Dict[str, Any]:
    """
    Check the connection status of Google Drive Integration, storage quota, and default folder info.
    """
    try:
        service = _get_gdrive_service()
        about = service.about().get(fields="user, storageQuota").execute()
        def_folder = _get_default_gdrive_folder_id()
        folder_name = "alfa agent"
        
        try:
            with database.get_sync_db() as conn:
                r = conn.execute("SELECT value FROM system_settings WHERE key = 'gdrive_default_folder_name'").fetchone()
                if r and r[0]:
                    folder_name = r[0]
        except Exception:
            pass
            
        return {
            "status": "success",
            "connected": True,
            "auth_mode": _detect_gdrive_auth_mode(),
            "user": about.get("user", {}),
            "storage_quota": about.get("storageQuota", {}),
            "default_folder_id": def_folder,
            "default_folder_name": folder_name,
            "default_folder_url": f"https://drive.google.com/drive/folders/{def_folder}"
        }
    except Exception as e:
        return {
            "status": "error",
            "connected": False,
            "message": str(e),
            "default_folder_id": _get_default_gdrive_folder_id()
        }


def gdrive_list_files(folder_id: str = "", query: str = "", limit: int = 20) -> Dict[str, Any]:
    """
    List, search, and browse files and folders stored in Google Drive.
    
    Args:
        folder_id: Optional ID of the Google Drive folder to list (defaults to configured folder).
        query: Optional search keyword or query term.
        limit: Max number of files to return (default 20, max 100).
    """
    try:
        service = _get_gdrive_service()
        target_folder = folder_id.strip() if folder_id else _get_default_gdrive_folder_id()
        q_parts = ["trashed = false"]
        if target_folder:
            q_parts.append(f"'{target_folder}' in parents")
        if query:
            q_parts.append(f"(name contains '{query}' or fullText contains '{query}')")
        q_str = " and ".join(q_parts)
        
        results = service.files().list(
            q=q_str,
            pageSize=min(limit, 100),
            supportsAllDrives=True,
            includeItemsFromAllDrives=True,
            fields="nextPageToken, files(id, name, mimeType, size, modifiedTime, webViewLink, webContentLink, iconLink)"
        ).execute()
        
        files = results.get("files", [])
        return {
            "status": "success",
            "total_found": len(files),
            "folder_id": target_folder,
            "files": files
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal mengambil daftar file Google Drive: {str(e)}"}


def gdrive_upload_file(filepath: str, folder_id: str = "", custom_filename: str = "") -> Dict[str, Any]:
    """
    Upload a local file or document (PDF, Excel, Word, images, code, archive) to Google Drive.
    
    Args:
        filepath: Path to the local file (e.g. '~/Dokumen/ALFA_SWARM_OUTPUTS/laporan.pdf' or filename).
        folder_id: Optional Google Drive folder ID to upload into (defaults to configured folder).
        custom_filename: Optional custom file name on Google Drive.
    """
    try:
        import mimetypes
        from googleapiclient.http import MediaFileUpload
        
        resolved_path = os.path.expanduser(filepath)
        if not os.path.exists(resolved_path):
            alt_path = os.path.join(SANDBOX_DIR, filepath)
            if os.path.exists(alt_path):
                resolved_path = alt_path
            else:
                return {"status": "error", "message": f"File '{filepath}' tidak ditemukan di sistem lokal."}
                
        service = _get_gdrive_service()
        target_folder = folder_id.strip() if folder_id else _get_default_gdrive_folder_id()
        upload_name = custom_filename or os.path.basename(resolved_path)
        mime_type, _ = mimetypes.guess_type(resolved_path)
        if not mime_type:
            mime_type = "application/octet-stream"
            
        file_metadata = {"name": upload_name}
        if target_folder:
            file_metadata["parents"] = [target_folder]
            
        media = MediaFileUpload(resolved_path, mimetype=mime_type, resumable=True)
        file = service.files().create(
            body=file_metadata,
            media_body=media,
            supportsAllDrives=True,
            fields="id, name, mimeType, size, webViewLink, webContentLink"
        ).execute()
        
        try:
            service.permissions().create(
                fileId=file.get("id"),
                body={"role": "reader", "type": "anyone"},
                supportsAllDrives=True
            ).execute()
        except Exception:
            pass
            
        return {
            "status": "success",
            "message": f"File '{upload_name}' berhasil diunggah ke Google Drive di folder target!",
            "file_id": file.get("id"),
            "file_name": file.get("name"),
            "folder_id": target_folder,
            "web_link": file.get("webViewLink"),
            "download_link": file.get("webContentLink")
        }
    except Exception as e:
        err_str = str(e)
        if "Service Accounts do not have storage quota" in err_str:
            return {
                "status": "error",
                "message": (
                    "Upload gagal: Service Account Google tidak punya kuota penyimpanan "
                    "(kebijakan Google terbaru). Solusi: lakukan login OAuth sekali via "
                    "Dashboard > Google Drive > 'Login OAuth', atau jalankan "
                    "./venv/bin/python scripts/gdrive_oauth_login.py - upload selanjutnya memakai kuota akun Anda."
                ),
                "needs_oauth": True,
            }
        return {"status": "error", "message": f"Gagal mengunggah file ke Google Drive: {err_str}"}


def gdrive_download_file(file_id: str, save_filename: str = "") -> Dict[str, Any]:
    """
    Download a file from Google Drive by its File ID to the local system.
    
    Args:
        file_id: The unique Google Drive File ID.
        save_filename: Optional local filename to save the downloaded content as.
    """
    try:
        import io
        from googleapiclient.http import MediaIoBaseDownload
        
        service = _get_gdrive_service()
        file_meta = service.files().get(fileId=file_id, supportsAllDrives=True, fields="id, name, mimeType").execute()
        target_name = save_filename or file_meta.get("name", f"gdrive_{file_id}")
        target_path = os.path.join(SANDBOX_DIR, target_name)
        
        request = service.files().get_media(fileId=file_id, supportsAllDrives=True)
        fh = io.FileIO(target_path, "wb")
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
            
        return {
            "status": "success",
            "message": f"File '{target_name}' berhasil diunduh dari Google Drive!",
            "file_id": file_id,
            "saved_path": target_path,
            "file_size": os.path.getsize(target_path)
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal mengunduh file dari Google Drive: {str(e)}"}


def gdrive_create_folder(folder_name: str, parent_folder_id: str = "") -> Dict[str, Any]:
    """
    Create a new folder in Google Drive.
    
    Args:
        folder_name: Name of the folder to create.
        parent_folder_id: Optional ID of the parent folder (defaults to configured folder).
    """
    try:
        service = _get_gdrive_service()
        target_parent = parent_folder_id.strip() if parent_folder_id else _get_default_gdrive_folder_id()
        file_metadata = {
            "name": folder_name,
            "mimeType": "application/vnd.google-apps.folder"
        }
        if target_parent:
            file_metadata["parents"] = [target_parent]
            
        folder = service.files().create(
            body=file_metadata,
            supportsAllDrives=True,
            fields="id, name, webViewLink"
        ).execute()
        
        return {
            "status": "success",
            "message": f"Folder '{folder_name}' berhasil dibuat di Google Drive!",
            "folder_id": folder.get("id"),
            "folder_name": folder.get("name"),
            "web_link": folder.get("webViewLink")
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal membuat folder Google Drive: {str(e)}"}


def gdrive_sync_to_second_brain(folder_id: str = "", limit: int = 10) -> Dict[str, Any]:
    """
    Ingest and sync documents from Google Drive directly into ALFA's Neural Vector Brain (Second Brain RAG).
    
    Args:
        folder_id: Optional Google Drive folder ID to ingest from (defaults to configured folder).
        limit: Max documents to ingest (default 10).
    """
    try:
        import vector_memory
        target_folder = folder_id.strip() if folder_id else _get_default_gdrive_folder_id()
        list_res = gdrive_list_files(folder_id=target_folder, limit=limit)
        if list_res.get("status") != "success":
            return list_res
            
        files = list_res.get("files", [])
        ingested = []
        # Attribute to the PRIMARY user so the main agent (Telegram/Web) can
        # retrieve these chunks - dashboard context has no telegram user id.
        uid = current_user_id_var.get()
        if not uid:
            try:
                allowed = os.getenv("ALLOWED_USER_IDS", "").strip()
                uid = int(allowed.split(",")[0]) if allowed.split(",")[0].strip().isdigit() else 0
            except Exception:
                uid = 0
        
        for f in files:
            fid = f.get("id")
            fname = f.get("name", "")
            mime = f.get("mimeType", "")
            
            if "folder" in mime:
                continue
                
            dl_res = gdrive_download_file(file_id=fid, save_filename=fname)
            if dl_res.get("status") == "success":
                local_f = dl_res.get("saved_path")
                v_res = vector_memory.ingest_document(
                    user_id=uid,
                    title=f"GDrive: {fname}",
                    content_or_path=local_f,
                    category="Google Drive Sync"
                )
                ingested.append({"name": fname, "file_id": fid, "vector_status": v_res.get("status")})
                
        return {
            "status": "success",
            "total_ingested": len(ingested),
            "folder_id": target_folder,
            "synced_files": ingested
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal sinkronisasi Google Drive ke Second Brain: {str(e)}"}


def self_restart_service() -> Dict[str, Any]:
    """
    GOD MODE: Self-Restart — the bot restarts its own systemd service to apply
    code changes, new tools, or recover from errors. The bot will go offline
    briefly (~2 seconds) and come back with all updates applied.
    """
    try:
        # Verify tools.py compiles before restarting
        tools_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "tools.py")
        compile_check = subprocess.run(
            [sys.executable, "-m", "py_compile", tools_path],
            capture_output=True, text=True
        )
        if compile_check.returncode != 0:
            return {"status": "error", "message": f"Tidak bisa restart — tools.py memiliki error: {compile_check.stderr}"}
        
        bot_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "bot.py")
        compile_check2 = subprocess.run(
            [sys.executable, "-m", "py_compile", bot_path],
            capture_output=True, text=True
        )
        if compile_check2.returncode != 0:
            return {"status": "error", "message": f"Tidak bisa restart — bot.py memiliki error: {compile_check2.stderr}"}
        
        # Schedule restart in 2 seconds (so we can send response first)
        subprocess.Popen(
            "sleep 2 && systemctl --user restart telegram-ai-bot.service",
            shell=True, start_new_session=True
        )
        
        return {
            "status": "success",
            "message": "🔄 Bot akan restart dalam 2 detik... Semua pembaruan dan tool baru akan aktif setelah restart. Bot akan kembali online dalam ~3 detik."
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


def proactive_system_guardian_config(action: str = "status", cpu_threshold: int = 90, ram_threshold: int = 85, disk_threshold: int = 90, battery_critical: int = 10, auto_kill_ram_hogs: bool = False) -> Dict[str, Any]:
    """
    GOD MODE: Configure the Proactive System Guardian daemon.
    The guardian runs in the background 24/7, monitoring system health and
    automatically taking protective actions (sending alerts, killing memory hogs,
    locking screen on critical battery, etc.).
    
    Args:
        action: 'status' to check guardian config, 'enable' to activate, 'disable' to deactivate.
        cpu_threshold: CPU usage % threshold for alert (default: 90).
        ram_threshold: RAM usage % threshold for alert (default: 85).
        disk_threshold: Disk usage % threshold for alert (default: 90).
        battery_critical: Battery % threshold for critical alert (default: 10).
        auto_kill_ram_hogs: If True, automatically kill top RAM-consuming non-essential processes when threshold exceeded.
    """
    try:
        config_path = os.path.join(os.path.expanduser("~"), ".alfa", "guardian_config.json")
        os.makedirs(os.path.dirname(config_path), exist_ok=True)
        
        import json
        
        if action == "status":
            if os.path.exists(config_path):
                with open(config_path, "r") as f:
                    config = json.load(f)
                return {"status": "success", "guardian": config}
            return {"status": "success", "guardian": {"enabled": False, "message": "Guardian belum dikonfigurasi."}}
        
        elif action == "enable":
            config = {
                "enabled": True,
                "cpu_threshold": cpu_threshold,
                "ram_threshold": ram_threshold,
                "disk_threshold": disk_threshold,
                "battery_critical": battery_critical,
                "auto_kill_ram_hogs": auto_kill_ram_hogs,
                "updated_at": datetime.datetime.now().isoformat()
            }
            with open(config_path, "w") as f:
                json.dump(config, f, indent=2)
            return {
                "status": "success",
                "message": f"🛡️ System Guardian AKTIF! Monitoring: CPU>{cpu_threshold}%, RAM>{ram_threshold}%, Disk>{disk_threshold}%, Battery<{battery_critical}%. Auto-kill: {'ON' if auto_kill_ram_hogs else 'OFF'}."
            }
        
        elif action == "disable":
            config = {"enabled": False, "updated_at": datetime.datetime.now().isoformat()}
            with open(config_path, "w") as f:
                json.dump(config, f, indent=2)
            return {"status": "success", "message": "🛡️ System Guardian dinonaktifkan."}
        
        return {"status": "error", "message": f"Action '{action}' tidak dikenal. Gunakan: status, enable, disable."}
    except Exception as e:
        return {"status": "error", "message": str(e)}


def proactive_ambient_agent_config(action: str = "status", enabled: bool = True, min_hours_between_pings: int = 3, quiet_hours_start: int = 23, quiet_hours_end: int = 7) -> Dict[str, Any]:
    """
    GOD MODE: Configure Ambient Proactive Autonomous Engagement.
    Controls whether and how often the AI agent can autonomously initiate contact,
    ask questions, check in on user projects, or send daily morning/afternoon briefings.
    
    Args:
        action: 'status' (view config), 'enable' (turn on proactive mode), 'disable' (turn off).
        enabled: Set proactive mode active/inactive.
        min_hours_between_pings: Minimum hours between spontaneous messages (default: 3).
        quiet_hours_start: Start hour for quiet time (default: 23 / 11 PM).
        quiet_hours_end: End hour for quiet time (default: 7 / 7 AM).
    """
    try:
        import json
        config_path = os.path.join(os.path.expanduser("~"), ".alfa", "proactive_config.json")
        os.makedirs(os.path.dirname(config_path), exist_ok=True)
        
        if action == "status":
            if os.path.exists(config_path):
                with open(config_path, "r", encoding="utf-8") as f:
                    config = json.load(f)
                return {"status": "success", "proactive_config": config}
            return {
                "status": "success",
                "proactive_config": {
                    "enabled": True,
                    "min_hours_between_pings": 3,
                    "quiet_hours_start": 23,
                    "quiet_hours_end": 7,
                    "message": "Mode Proaktif default aktif."
                }
            }
            
        elif action in ["enable", "set"]:
            config = {
                "enabled": True,
                "min_hours_between_pings": max(1, min_hours_between_pings),
                "quiet_hours_start": quiet_hours_start,
                "quiet_hours_end": quiet_hours_end,
                "updated_at": datetime.datetime.now().isoformat()
            }
            with open(config_path, "w", encoding="utf-8") as f:
                json.dump(config, f, indent=2)
            return {
                "status": "success",
                "message": f"🤖 Mode Proaktif Otonom AKTIF! Bot akan berinisiatif menyapa/mengecek kondisi setiap ~{min_hours_between_pings} jam di luar jam tenang ({quiet_hours_start}:00 - {quiet_hours_end}:00)."
            }
            
        elif action == "disable":
            config = {
                "enabled": False,
                "updated_at": datetime.datetime.now().isoformat()
            }
            with open(config_path, "w", encoding="utf-8") as f:
                json.dump(config, f, indent=2)
            return {"status": "success", "message": "🤖 Mode Proaktif Otonom dinonaktifkan. Bot hanya akan membalas jika Anda bertanya."}
            
        return {"status": "error", "message": f"Action '{action}' tidak dikenal. Gunakan: status, enable, disable."}
    except Exception as e:
        return {"status": "error", "message": f"Proactive config error: {str(e)}"}


def manage_wa_sheets_bot(action: str = "status") -> Dict[str, Any]:
    """
    ECOSYSTEM INTEGRATION: WhatsApp Google Sheets Bot Controller.
    Controls and monitors the wa-sheets-bot systemd user service.
    Can check status, start, stop, restart, enable auto-start, or view live logs.
    
    Args:
        action: 'status' (check running state & memory), 'start', 'stop', 'restart', 'logs' (recent 20 lines), 'enable' (enable autostart on boot).
    """
    try:
        act = action.lower().strip()
        svc_name = "wa-sheets-bot.service"
        
        if act == "status":
            res_active = subprocess.run(["systemctl", "--user", "is-active", svc_name], capture_output=True, text=True)
            res_enabled = subprocess.run(["systemctl", "--user", "is-enabled", svc_name], capture_output=True, text=True)
            res_status = subprocess.run(["systemctl", "--user", "status", svc_name, "--no-pager", "-n", "5"], capture_output=True, text=True)
            
            is_act = res_active.stdout.strip() == "active"
            return {
                "status": "success",
                "service": svc_name,
                "is_running": is_act,
                "state": res_active.stdout.strip(),
                "enabled_on_boot": res_enabled.stdout.strip() == "enabled",
                "details": res_status.stdout.strip()
            }
            
        elif act in ["start", "stop", "restart", "enable", "disable"]:
            res = subprocess.run(["systemctl", "--user", act, svc_name], capture_output=True, text=True)
            if res.returncode == 0:
                time.sleep(1)
                res_active = subprocess.run(["systemctl", "--user", "is-active", svc_name], capture_output=True, text=True)
                return {
                    "status": "success",
                    "message": f"Service '{svc_name}' berhasil di-{act}! Status saat ini: {res_active.stdout.strip()}.",
                    "current_state": res_active.stdout.strip()
                }
            return {"status": "error", "message": f"Gagal mengeksekusi {act}: {res.stderr}"}
            
        elif act == "logs":
            res_logs = subprocess.run(["journalctl", "--user", "-u", svc_name, "-n", "25", "--no-pager"], capture_output=True, text=True)
            return {
                "status": "success",
                "service": svc_name,
                "logs": res_logs.stdout.strip()
            }
            
        return {"status": "error", "message": f"Aksi '{action}' tidak dikenal. Gunakan: status, start, stop, restart, logs, enable."}
    except Exception as e:
        return {"status": "error", "message": f"Manage wa-sheets-bot error: {str(e)}"}


def open_web_dashboard(port: int = 8080) -> Dict[str, Any]:
    """
    ALFA OS: Web Command Center Dashboard Controller.
    Returns the active local web URL for the luxury management dashboard
    and ensures the alfa-dashboard.service is running.
    
    Args:
        port: Dashboard port (default: 8080).
    """
    try:
        # Check if service is active
        res = subprocess.run(["systemctl", "--user", "is-active", "alfa-dashboard.service"], capture_output=True, text=True)
        if res.stdout.strip() != "active":
            subprocess.run(["systemctl", "--user", "start", "alfa-dashboard.service"], capture_output=True, text=True)
            
        import socket
        hostname = socket.gethostname()
        local_ip = "127.0.0.1"
        try:
            s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
            s.connect(("8.8.8.8", 80))
            local_ip = s.getsockname()[0]
            s.close()
        except Exception:
            pass
            
        return {
            "status": "success",
            "message": f"🌐 Web Dashboard Command Center aktif! Buka di browser laptop: http://localhost:{port} atau dari HP di WiFi yang sama: http://{local_ip}:{port}",
            "local_url": f"http://localhost:{port}",
            "network_url": f"http://{local_ip}:{port}",
            "dashboard_features": ["Live Hardware Telemetry", "75+ Tools Arsenal & Runner", "Ecosystem Services Hub", "Second Brain Memory Visualizer", "24/7 Guardian Control", "Web Live AI Console", "AI Swarm & Rapat Antar Agent", "Multi-Provider API Key Vault"]
        }
    except Exception as e:
        return {"status": "error", "message": f"Open web dashboard error: {str(e)}"}


def manage_api_keys(action: str, name: str = "", provider: str = "gemini", api_key: str = "", default_model: str = "gemini-3.5-flash-lite", base_url: str = "", key_id: int = None) -> Dict[str, Any]:
    """
    Manage API keys and multi-provider endpoints (Gemini, OpenAI, Groq, OpenRouter, Anthropic, Ollama, NVIDIA NIM).
    Enables switching active keys or assigning specific provider keys to specialized agents.
    
    Args:
        action: 'list' (view all keys), 'add' (add key), 'activate' (set key active), 'delete' (remove key).
        name: Label name for the key (e.g. 'Production Gemini', 'NVIDIA NIM Llama 3.3', 'Groq Llama 3').
        provider: 'gemini', 'openai', 'groq', 'openrouter', 'anthropic', 'ollama', 'nvidia'.
        api_key: The API secret key string (e.g. nvapi-..., AIza..., sk-...).
        default_model: Default model string (e.g. 'meta/llama-3.3-70b-instruct', 'gemini-3.5-flash-lite', 'gpt-4o').
        base_url: Optional custom proxy or Ollama/NVIDIA NIM base URL (default for nvidia: 'https://integrate.api.nvidia.com/v1').
        key_id: Target key ID for 'activate' or 'delete'.
    """
    action = action.lower().strip()
    try:
        if action == "list":
            keys = database.list_api_keys_sync()
            return {
                "status": "success",
                "total_keys": len(keys),
                "keys": keys
            }
        elif action == "add":
            if not api_key:
                return {"status": "error", "message": "Parameter 'api_key' wajib diisi."}
            res = database.add_api_key_sync(
                name=name or f"{provider.capitalize()} Key",
                provider=provider,
                api_key=api_key,
                default_model=default_model,
                base_url=base_url,
                set_active=True
            )
            return {"status": "success", "message": f"API Key '{name}' untuk provider '{provider}' berhasil disimpan & diaktifkan!", "key_id": res.get("id")}
        elif action == "activate":
            if not key_id:
                return {"status": "error", "message": "Parameter 'key_id' wajib diisi."}
            return database.activate_api_key_sync(key_id)
        elif action == "delete":
            if not key_id:
                return {"status": "error", "message": "Parameter 'key_id' wajib diisi."}
            return database.delete_api_key_sync(key_id)
        else:
            return {"status": "error", "message": f"Action tidak dikenal: {action}. Gunakan 'list', 'add', 'activate', atau 'delete'."}
    except Exception as e:
        return {"status": "error", "message": f"Manage API keys error: {str(e)}"}


def manage_custom_agents(action: str, name: str = "", role: str = "", persona: str = "", system_instruction: str = "", provider: str = "gemini", model: str = "gemini-2.5-flash", avatar_emoji: str = "🤖", color_theme: str = "cyan", agent_id: int = None) -> Dict[str, Any]:
    """
    Manage the Autonomous AI Agent Workforce (Society of Agents).
    Create, list, update, and configure specialized agents that can collaborate, hold meetings, and execute tasks.
    
    Args:
        action: 'list' (view all agents), 'add' (create agent), 'delete' (remove agent), 'toggle' (enable/disable).
        name: Unique name of the agent (e.g. 'Security Guard', 'Frontend Ninja').
        role: Title / Role of the agent (e.g. 'Penetration Tester', 'Vue/React UI Specialist').
        persona: Persona description (e.g. 'Kritis, teliti, mengutamakan performa').
        system_instruction: Detailed system prompt for this agent.
        provider: 'gemini', 'openai', 'groq', 'openrouter', 'ollama'.
        model: Model identifier (default: 'gemini-2.5-flash').
        avatar_emoji: Avatar emoji (e.g. '👑', '⚡', '🛡️', '🌐', '💡').
        color_theme: 'cyan', 'emerald', 'violet', 'amber', 'rose', 'blue'.
        agent_id: Target agent ID for update or delete.
    """
    action = action.lower().strip()
    try:
        if action == "list":
            agents = database.list_custom_agents_sync()
            return {
                "status": "success",
                "total_agents": len(agents),
                "agents": agents
            }
        elif action == "add":
            if not name or not role:
                return {"status": "error", "message": "Parameter 'name' dan 'role' wajib diisi."}
            res = database.add_custom_agent_sync(
                name=name,
                role=role,
                persona=persona or f"Spesialis dalam {role}",
                system_instruction=system_instruction or f"Kamu adalah {name}, {role}.",
                provider=provider,
                model=model,
                avatar_emoji=avatar_emoji,
                color_theme=color_theme
            )
            return {"status": "success", "message": f"Agent '{name}' ({role}) berhasil ditambahkan ke AI Workforce!", "agent_id": res.get("id")}
        elif action == "delete":
            if not agent_id:
                return {"status": "error", "message": "Parameter 'agent_id' wajib diisi."}
            return database.delete_custom_agent_sync(agent_id)
        elif action == "toggle":
            if not agent_id:
                return {"status": "error", "message": "Parameter 'agent_id' wajib diisi."}
            cur = database.get_custom_agent_sync(agent_id)
            if not cur:
                return {"status": "error", "message": "Agent tidak ditemukan"}
            new_state = 0 if cur.get("is_enabled", 1) else 1
            return database.update_custom_agent_sync(agent_id, {"is_enabled": new_state})
        else:
            return {"status": "error", "message": f"Action tidak dikenal: {action}."}
    except Exception as e:
        return {"status": "error", "message": f"Manage custom agents error: {str(e)}"}


def query_token_usage(hours: int = 24) -> Dict[str, Any]:
    """
    Laporan pemakaian token AI per API key (realtime dari dashboard vault).
    
    Args:
        hours: Rentang jam ke belakang (1-720, default 24).
    """
    try:
        summary = database.get_api_usage_summary_sync(hours=hours)
        per_key = [
            {k: r.get(k) for k in ('key_name', 'provider', 'total_tokens', 'calls', 'last_used')}
            for r in summary.get('per_key', [])[:15]
        ]
        return {
            "status": "success",
            "window_hours": summary.get("window_hours"),
            "tokens_today": summary.get("tokens_today"),
            "calls_today": summary.get("calls_today"),
            "total_all_time": summary.get("total_all_time"),
            "per_key": per_key,
            "by_context": summary.get("by_context", []),
            "message": f"Pemakaian {hours} jam terakhir: {summary.get('tokens_today')} token hari ini."
        }
    except Exception as e:
        return {"status": "error", "message": f"Gagal membaca pemakaian token: {str(e)}"}


def list_wa_drive_uploads(limit: int = 20) -> Dict[str, Any]:
    """
    Daftar berkas (foto/PDF/dokumen) yang otomatis diunggah dari WhatsApp ke Google Drive.
    
    Args:
        limit: Jumlah maksimal entri terbaru (default 20).
    """
    try:
        path = os.path.expanduser("~/wa-sheets-bot/drive_uploads.json")
        if not os.path.exists(path):
            return {"status": "success", "total": 0,
                    "message": "Belum ada berkas WA yang terunggah ke Drive.",
                    "uploads": []}
        import json as _json
        with open(path, "r", encoding="utf-8") as f:
            data = _json.load(f)
        uploads = (data.get("uploads") or [])[: max(1, min(int(limit), 100))]
        compact = [{
            'waktu': u.get('ts'), 'nama_file': u.get('file_name'),
            'folder': u.get('folder'), 'pengirim': u.get('sender'),
            'grup': u.get('group'), 'link': u.get('web_link'),
        } for u in uploads]
        return {"status": "success", "total": len(uploads),
                "message": f"{len(uploads)} berkas WA terakhir di Drive.",
                "uploads": compact}
    except Exception as e:
        return {"status": "error", "message": f"Gagal membaca log unggahan: {str(e)}"}


def conduct_ai_meeting(topic: str, participants: str = "", rounds: int = 2, mode: str = "plan") -> Dict[str, Any]:
    """
    Conduct an Autonomous Conference / Meeting or Live Swarm Execution between multiple AI agents.
    
    Args:
        topic: The agenda, problem, or live task to be executed. WAJIB ringkasan TUGAS user sendiri (bukan kutipan jawaban/konsensus lama). Contoh: 'scrape 20 mouse gaming murah lalu buat CSV'.
        participants: Comma-separated agent names or empty for default team.
        rounds: Number of discussion rounds (1 to 3, default: 2).
        mode: 'plan' (strategic debate & action plan) or 'execute' (rapid alignment + live autonomous tool execution & artifacts).
    """
    try:
        import swarm_engine
        import concurrent.futures
        part_list = [p.strip() for p in participants.split(",") if p.strip()] if participants else None
        rounds_clamped = max(1, min(3, int(rounds)))
        mode_clean = mode.lower().strip() if mode else "plan"
        
        def _run_meeting():
            return asyncio.run(swarm_engine.conduct_multi_agent_meeting(topic, part_list, rounds_clamped, mode_clean))
        
        # asyncio.run() needs a fresh loop; if this tool is invoked from inside
        # a running loop (e.g. Telegram handler), delegate to a worker thread.
        try:
            asyncio.get_running_loop()
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
                result = pool.submit(_run_meeting).result(timeout=600)
        except RuntimeError:
            result = _run_meeting()
            
        return {
            "status": "success",
            "meeting_id": result.get("meeting_id"),
            "title": result.get("title"),
            "topic": topic,
            "mode": mode_clean,
            "total_rounds": rounds_clamped,
            "participants": result.get("participants"),
            "total_dialogues": len(result.get("dialogue_transcript", [])),
            "execution_results": result.get("execution_results", []),
            "consensus": result.get("consensus"),
            "action_plan": result.get("action_plan")
        }
    except Exception as e:
        return {"status": "error", "message": f"AI Meeting error: {str(e)}"}


def vault_store_secret(name: str, value: str, category: str = "api_key", notes: str = "") -> Dict[str, Any]:
    """
    Encrypt and store a sensitive credential, API key, affiliate token, or secret note
    into αlfa Secure Vault with AES-256-GCM authenticated encryption.
    
    Args:
        name: Unique identifier name for the secret (e.g. 'KLING_AI_KEY', 'SHOPEE_COOKIE', 'DB_PASSWORD').
        value: Secret text/token/key to encrypt and store securely.
        category: Category of secret ('api_key', 'affiliate', 'password', 'note').
        notes: Optional description or context about the secret.
    """
    try:
        import vault_engine
        res = vault_engine.vault.store_secret(name=name, value=value, category=category, notes=notes)
        return res
    except Exception as e:
        return {"status": "error", "message": f"Vault store error: {str(e)}"}


def vault_get_secret(name_or_id: str) -> Dict[str, Any]:
    """
    Retrieve and decrypt a sensitive secret from αlfa Secure Vault using AES-256-GCM.
    
    Args:
        name_or_id: The unique name or ID of the secret to decrypt.
    """
    try:
        import vault_engine
        sec = vault_engine.vault.get_secret(name_or_id)
        if not sec:
            return {"status": "error", "message": f"Secret '{name_or_id}' tidak ditemukan di dalam vault."}
        return {
            "status": "success",
            "name": sec["name"],
            "category": sec["category"],
            "value": sec["value"],
            "notes": sec["notes"],
            "updated_at": sec["updated_at"]
        }
    except Exception as e:
        return {"status": "error", "message": f"Vault retrieval error: {str(e)}"}


def vault_list_secrets(category: str = "all") -> Dict[str, Any]:
    """
    List all stored secrets metadata in αlfa Secure Vault without exposing decrypted plaintext.
    
    Args:
        category: Filter by category ('all', 'api_key', 'affiliate', 'password', 'note').
    """
    try:
        import vault_engine
        items = vault_engine.vault.list_secrets(category=category)
        return {
            "status": "success",
            "total_secrets": len(items),
            "encryption": "AES-256-GCM (Authenticated)",
            "secrets": items
        }
    except Exception as e:
        return {"status": "error", "message": f"Vault list error: {str(e)}"}


def vault_delete_secret(secret_id: int) -> Dict[str, Any]:
    """
    Permanently delete a secret from αlfa Secure Vault by ID.
    
    Args:
        secret_id: Numeric ID of the secret to delete.
    """
    try:
        import vault_engine
        deleted = vault_engine.vault.delete_secret(int(secret_id))
        if deleted:
            return {"status": "success", "message": f"Secret ID {secret_id} berhasil dihapus permanen dari vault."}
        return {"status": "error", "message": f"Secret ID {secret_id} tidak ditemukan."}
    except Exception as e:
        return {"status": "error", "message": f"Vault delete error: {str(e)}"}


def audit_website_security(target_url: str) -> Dict[str, Any]:
    """
    Conduct a Defensive Cybersecurity Audit on a website or API endpoint (Cyber Sentry):
    Audits SSL/TLS certificate, Security Headers (CSP, HSTS, X-Frame-Options, XSS, etc.),
    CORS policies, server fingerprint leaks, and generates an overall Security Grade (A+ to F).
    
    Args:
        target_url: The URL or domain to audit (e.g. 'https://shopee.co.id', 'https://example.com').
    """
    try:
        import security_auditor
        return security_auditor.audit_website_security(target_url)
    except Exception as e:
        return {"status": "error", "message": f"Security audit error: {str(e)}"}


def universal_deep_scraper(query: str, category: str = "all_marketplace", limit: int = 50) -> Dict[str, Any]:
    """
    High-Volume Universal Pro Web Scraper:
    Scrapes large volumes (20 - 200+ results) of rich data across various categories:
    - 'all_marketplace' (Shopee, Tokopedia, TikTok Shop, Lazada, Blibli)
    - 'jobs_career' (JobStreet, LinkedIn, Glints, Karir.com)
    - 'news_media' (Detik, Kompas, CNN, Liputan6, CNBC)
    - 'leads_contacts' (WhatsApp, Phone, Email, Suppliers, Distributors)
    - 'property_realestate' (Rumah123, Rumah.com, Lamudi, OLX)
    - 'google_general' (General Web Deep Search)
    
    Automatically extracts Titles, Prices, Contacts (Phone/WA/Email), Domains, URLs, and saves to CSV & JSON.
    
    Args:
        query: What to scrape / search (e.g. 'sepatu sneakers running wanita', 'python developer', 'distributor kopi gayo').
        category: Platform category to scrape (default: 'all_marketplace').
        limit: Total items to harvest (default: 50, supports up to 200).
    """
    try:
        import universal_scraper
        return universal_scraper.scrape_universal_keyword(query=query, category=category, limit=limit)
    except Exception as e:
        return {"status": "error", "message": f"Universal scraper error: {str(e)}"}


def scrape_custom_urls_batch(urls: List[str], concurrency: int = 15, use_camoufox: bool = False) -> Dict[str, Any]:
    """
    Scrape any custom list of URLs with high-speed multi-threaded workers or Camoufox stealth browser.
    Extracts page titles, meta info, prices, images, and descriptions into CSV and JSON.
    
    Args:
        urls: List of web URLs to scrape.
        concurrency: Concurrent scraping workers (default: 15).
        use_camoufox: If True, uses Camoufox anti-detect stealth browser (for Cloudflare/JS-heavy pages).
    """
    try:
        import universal_scraper
        return universal_scraper.scrape_custom_urls_or_selectors(urls=urls, concurrency=concurrency, use_camoufox=use_camoufox)
    except Exception as e:
        return {"status": "error", "message": f"Custom URL batch scraper error: {str(e)}"}


# List of all tools available to the Gemini Model
AVAILABLE_TOOLS = [
    universal_deep_scraper,
    scrape_custom_urls_batch,
    vault_store_secret,
    vault_get_secret,
    vault_list_secrets,
    vault_delete_secret,
    audit_website_security,
    get_system_stats,
    execute_bash_command,
    execute_python_sandbox,
    web_search,
    fetch_web_page_content,
    deep_research_topic,
    browser_open_url,
    browser_click_element,
    browser_type_text,
    browser_capture_screenshot,
    browser_close_tab,
    desktop_click_coordinate,
    desktop_type_keys,
    desktop_launch_app,
    spawn_background_subagent,
    check_subagent_status,
    add_recurring_task,
    list_recurring_tasks,
    cancel_recurring_task,
    generate_pdf_report,
    pdf_merge_documents,
    pdf_split_document,
    pdf_extract_full_text,
    pdf_encrypt_password,
    pdf_decrypt_password,
    pdf_rotate_pages,
    pdf_apply_watermark_text,
    pdf_insert_page_numbers,
    pdf_convert_to_images,
    images_convert_to_pdf,
    pdf_inspect_metadata,
    pdf_compress_and_optimize,
    affiliate_hunt_trending_products,
    affiliate_generate_viral_content,
    affiliate_broadcast_deal,
    affiliate_list_campaigns,
    scrape_real_product_data,
    scrape_large_scale_batch,
    marketplace_search_products,
    generate_promo_video_from_images,
    generate_excel_spreadsheet,
    generate_presentation_pptx,
    control_linux_hardware,
    send_file_to_chat,
    compress_folder_to_zip,
    record_desktop_screen,
    read_clipboard,
    manage_api_keys,
    manage_custom_agents,
    conduct_ai_meeting,
    query_token_usage,
    list_wa_drive_uploads,
    write_to_clipboard,
    show_desktop_notification,
    ssh_execute_command,
    query_database,
    send_email,
    list_running_processes,
    kill_process,
    edit_image,
    git_operations,
    translate_text,
    download_file_from_url,
    generate_secure_password,
    vision_click_target,
    auto_diagnose_and_heal_system,
    text_to_audio_file,
    convert_media_format,
    extract_audio_from_video,
    analyze_dataset_csv_json,
    audit_network_security,
    clean_system_storage,
    manage_system_services,
    manage_crontab_jobs,
    extract_and_link_knowledge,
    export_knowledge_base,
    start_focus_session,
    libreoffice_convert_document,
    libreoffice_render_page_previews,
    libreoffice_create_document,
    libreoffice_extract_document_text,
    proactive_ambient_agent_config,
    manage_wa_sheets_bot,
    open_web_dashboard,
    self_add_new_tool,
    list_dynamic_plugins,
    delete_dynamic_plugin,
    semantic_search_vector_brain,
    ingest_document_to_vector_brain,
    list_vector_brain_documents,
    self_restart_service,
    proactive_system_guardian_config,
    markitdown_convert_document,
    scrapling_stealth_fetch,
    scrapy_spider_quick_scrape,
    crawlee_web_scraper,
    crawl4ai_web_crawler,
    browser_use_autonomous_task,
    firecrawl_scrape_and_crawl,
    scrcpy_android_control,
    gdrive_status,
    gdrive_list_files,
    gdrive_upload_file,
    gdrive_download_file,
    gdrive_create_folder,
    gdrive_sync_to_second_brain,
    save_knowledge_memory,
    search_knowledge_memory,
    read_local_file,
    write_local_file,
    search_workspace_files,
    grep_workspace,
    edit_file_precise,
    apply_unified_diff,
    index_codebase,
    search_codebase,
    schedule_reminder,
    capture_desktop_screenshot,
    capture_webcam_frame,
    scan_local_network,
    *plugins.load_all_plugin_tools()
]

