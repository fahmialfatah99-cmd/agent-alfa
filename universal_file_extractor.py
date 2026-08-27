"""
UNIVERSAL FILE & MULTIMODAL ATTACHMENT EXTRACTOR.
Enterprise-grade parser for all document, spreadsheet, code, archive, audio, and visual formats:
- PDF: PyMuPDF (fitz) + pypdf page extractor
- Word: python-docx + mammoth paragraph & table parser
- PowerPoint: python-pptx per-slide extractor with speaker notes
- Excel: openpyxl + pandas multi-sheet markdown table converter
- CSV/TSV: csv.Sniffer with auto delimiter detection
- HTML/XML: BeautifulSoup clean text extractor
- JSON/YAML/TOML: structure parser
- Code: 30+ programming language syntax blocks with line numbers
- Archives: ZIP/TAR directory hierarchy and nested file extractor
- Audio/Video/Images: Gemini native multimodal Part generator + Real File Path Persistence
"""

import io
import os
import re
import csv
import time
import json
import tarfile
import zipfile
import logging
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger("UniversalFileExtractor")

UPLOAD_DIR = os.path.expanduser("~/Dokumen/ALFA_SWARM_OUTPUTS/uploads")
os.makedirs(UPLOAD_DIR, exist_ok=True)


def process_uploaded_attachment(
    filename: str,
    mime_type: str,
    raw_bytes: bytes,
    save_disk: bool = True
) -> Tuple[Optional[str], Optional[Any]]:
    """
    Process any uploaded file and return (text_context, multimodal_part).
    - Persists file to ~/Dokumen/ALFA_SWARM_OUTPUTS/uploads/ so real tools can operate on it.
    - If file is an image or audio or native PDF: returns multimodal_part for Gemini vision/audio/doc models.
    - Returns rich structured text context with exact file disk path for immediate tool execution.
    """
    ext = os.path.splitext(filename)[1].lower()
    fname_clean = os.path.basename(filename)
    safe_name = re.sub(r'[^a-zA-Z0-9_.-]', '_', fname_clean)
    
    file_disk_path = ""
    if save_disk:
        try:
            file_disk_path = os.path.join(UPLOAD_DIR, f"{int(time.time())}_{safe_name}")
            with open(file_disk_path, "wb") as f_disk:
                f_disk.write(raw_bytes)
        except Exception as e:
            logger.warning(f"Gagal simpan upload ke disk: {e}")

    path_header = f"[FILE TERSIMPAN DI DISK: {file_disk_path}]\n" if file_disk_path else ""
    tool_hint = "*(Gunakan path file nyata di atas jika ada instruksi konversi seperti 'jadikan pdf', 'merge', 'kompres', 'ubah format', dll.)*\n"

    # 1. IMAGES (Vision + Real Tool Execution)
    if mime_type.startswith("image/") or ext in (".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tiff", ".ico"):
        multimodal_part = None
        try:
            from google.genai import types
            img_mime = mime_type if mime_type.startswith("image/") else f"image/{ext.lstrip('.')}"
            if img_mime == "image/jpg":
                img_mime = "image/jpeg"
            multimodal_part = types.Part.from_bytes(data=raw_bytes, mime_type=img_mime)
        except Exception as e:
            logger.warning(f"Native image part error: {e}")

        text_ctx = f"[LAMPIRAN GAMBAR/FOTO: {fname_clean}]\n{path_header}{tool_hint}"
        return text_ctx, multimodal_part

    # 2. AUDIO & VOICE (Native Audio + Real Tool Execution)
    if mime_type.startswith("audio/") or ext in (".mp3", ".m4a", ".wav", ".ogg", ".aac", ".flac"):
        multimodal_part = None
        try:
            from google.genai import types
            aud_mime = mime_type if mime_type.startswith("audio/") else f"audio/{ext.lstrip('.')}"
            multimodal_part = types.Part.from_bytes(data=raw_bytes, mime_type=aud_mime)
        except Exception as e:
            logger.warning(f"Native audio part error: {e}")

        text_ctx = f"[LAMPIRAN AUDIO/SUARA: {fname_clean}]\n{path_header}{tool_hint}"
        return text_ctx, multimodal_part

    # 3. PDF DOCUMENTS (Dual mode: PyMuPDF text + Native Part)
    if ext == ".pdf" or mime_type == "application/pdf":
        pdf_text = _extract_pdf(raw_bytes, fname_clean)
        multimodal_part = None
        try:
            from google.genai import types
            multimodal_part = types.Part.from_bytes(data=raw_bytes, mime_type="application/pdf")
        except Exception:
            pass
        return f"[DOKUMEN PDF: {fname_clean}]\n{path_header}{tool_hint}\n{pdf_text}", multimodal_part

    # 4. WORD DOCUMENTS (.docx, .doc)
    if ext in (".docx", ".doc"):
        doc_text = _extract_word(raw_bytes, fname_clean)
        return f"[DOKUMEN WORD: {fname_clean}]\n{path_header}{tool_hint}\n{doc_text}", None

    # 5. POWERPOINT PRESENTATIONS (.pptx, .ppt)
    if ext in (".pptx", ".ppt"):
        pptx_text = _extract_powerpoint(raw_bytes, fname_clean)
        return f"[PRESENTASI PPTX: {fname_clean}]\n{path_header}{tool_hint}\n{pptx_text}", None

    # 6. SPREADSHEETS (.xlsx, .xls, .xlsm, .ods)
    if ext in (".xlsx", ".xls", ".xlsm", ".xltx", ".ods"):
        sheet_text = _extract_excel(raw_bytes, fname_clean, ext)
        return f"[SPREADSHEET EXCEL: {fname_clean}]\n{path_header}{tool_hint}\n{sheet_text}", None

    # 7. CSV / TSV / DELIMITED DATA
    if ext in (".csv", ".tsv", ".dsv", ".txt") and ("," in raw_bytes[:1000].decode("utf-8", errors="ignore") or "\t" in raw_bytes[:1000].decode("utf-8", errors="ignore")):
        csv_text = _extract_csv(raw_bytes, fname_clean)
        if csv_text:
            return f"[DATA TABEL CSV: {fname_clean}]\n{path_header}{tool_hint}\n{csv_text}", None

    # 8. ARCHIVES (.zip, .tar, .tar.gz, .tgz)
    if ext in (".zip", ".tar", ".gz", ".tgz", ".bz2"):
        archive_text = _extract_archive(raw_bytes, fname_clean, ext)
        return f"[ARSIP TERKOMPRESI: {fname_clean}]\n{path_header}{tool_hint}\n{archive_text}", None

    # 9. HTML / XML / SVG
    if ext in (".html", ".htm", ".xhtml", ".xml", ".svg"):
        html_text = _extract_html_xml(raw_bytes, fname_clean, ext)
        return f"[KODE WEB / STRUKTUR: {fname_clean}]\n{path_header}\n{html_text}", None

    # 10. JSON / YAML / TOML
    if ext in (".json", ".jsonl", ".ndjson", ".yaml", ".yml", ".toml"):
        data_text = _extract_data_formats(raw_bytes, fname_clean, ext)
        return f"[DATA KONFIGURASI: {fname_clean}]\n{path_header}\n{data_text}", None

    # 11. SOURCE CODE & SCRIPTS
    code_exts = {
        ".py": "python", ".js": "javascript", ".ts": "typescript", ".tsx": "tsx",
        ".jsx": "jsx", ".java": "java", ".c": "c", ".cpp": "cpp", ".h": "c",
        ".cs": "csharp", ".go": "go", ".rs": "rust", ".php": "php", ".rb": "ruby",
        ".sh": "bash", ".bash": "bash", ".zsh": "bash", ".sql": "sql", ".css": "css",
        ".scss": "scss", ".vue": "vue", ".svelte": "svelte", ".proto": "protobuf",
        ".graphql": "graphql", ".dockerfile": "dockerfile", ".md": "markdown", ".txt": "text",
        ".log": "log", ".env": "ini"
    }
    if ext in code_exts or mime_type.startswith("text/"):
        lang = code_exts.get(ext, "text")
        try:
            decoded = raw_bytes.decode("utf-8", errors="replace")
            lines = decoded.splitlines()
            numbered = "\n".join([f"{i+1:4d} | {line}" for i, line in enumerate(lines[:1000])])
            total_hint = f" (Menampilkan {min(len(lines), 1000)} dari total {len(lines)} baris)" if len(lines) > 1000 else ""
            return f"[KODE SUMBER: {fname_clean}{total_hint}]\n{path_header}\n```{lang}\n{numbered}\n```", None
        except Exception:
            pass

    # 12. Fallback UTF-8 text decoder
    try:
        decoded = raw_bytes.decode("utf-8", errors="replace")
        return f"[BERKAS TEKS: {fname_clean}]\n{path_header}\n{decoded}", None
    except Exception:
        return f"[BERKAS BINER: {fname_clean} ({len(raw_bytes)} bytes)]\n{path_header}", None


# ── EXTRACTOR IMPLEMENTATIONS ─────────────────────────────────────────────────

def _extract_pdf(raw_bytes: bytes, fname: str) -> str:
    try:
        import fitz
        doc = fitz.open(stream=raw_bytes, filetype="pdf")
        pages = []
        for i, page in enumerate(doc):
            text = page.get_text("text").strip()
            if text:
                pages.append(f"--- Halaman {i+1} ---\n{text}")
        if pages:
            return f"Total {len(doc)} halaman:\n\n" + "\n\n".join(pages[:50])
    except Exception as e:
        logger.debug(f"fitz pdf error: {e}")

    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(raw_bytes))
        pages = []
        for i, page in enumerate(reader.pages):
            txt = page.extract_text()
            if txt:
                pages.append(f"--- Halaman {i+1} ---\n{txt}")
        if pages:
            return f"Total {len(reader.pages)} halaman:\n\n" + "\n\n".join(pages[:50])
    except Exception as e:
        logger.debug(f"pypdf error: {e}")

    return f"Gagal membaca teks dari PDF {fname}."


def _extract_word(raw_bytes: bytes, fname: str) -> str:
    try:
        import docx
        doc = docx.Document(io.BytesIO(raw_bytes))
        out = []
        for p in doc.paragraphs:
            if p.text.strip():
                if p.style.name.startswith("Heading"):
                    out.append(f"### {p.text.strip()}")
                else:
                    out.append(p.text.strip())

        for t_idx, table in enumerate(doc.tables, 1):
            out.append(f"\n#### Tabel {t_idx}")
            rows = []
            for r in table.rows:
                rows.append([cell.text.strip().replace("\n", " ") for cell in r.cells])
            if rows:
                headers = rows[0]
                cols = len(headers)
                out.append("| " + " | ".join(headers) + " |")
                out.append("| " + " | ".join(["---"] * cols) + " |")
                for r in rows[1:]:
                    out.append("| " + " | ".join(r[:cols]) + " |")

        return "\n\n".join(out)
    except Exception as e:
        logger.warning(f"docx error: {e}")
        return f"Gagal membaca dokumen Word {fname}."


def _extract_powerpoint(raw_bytes: bytes, fname: str) -> str:
    try:
        import pptx
        prs = pptx.Presentation(io.BytesIO(raw_bytes))
        slides_out = []
        for i, slide in enumerate(prs.slides, 1):
            s_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    s_texts.append(shape.text.strip())
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()

            slide_md = f"### Slide {i}\n" + "\n".join(f"- {t}" for t in s_texts)
            if notes:
                slide_md += f"\n*Catatan Presenter:* {notes}"
            slides_out.append(slide_md)

        return f"Total {len(prs.slides)} Slide Presentasi:\n\n" + "\n\n".join(slides_out)
    except Exception as e:
        logger.warning(f"pptx error: {e}")
        return f"Gagal membaca presentasi PPTX {fname}."


def _extract_excel(raw_bytes: bytes, fname: str, ext: str) -> str:
    if ext in (".xlsx", ".xlsm", ".xltx"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(raw_bytes), data_only=True)
            sheets_out = []
            for sname in wb.sheetnames:
                ws = wb[sname]
                rows = list(ws.iter_rows(values_only=True))
                clean_rows = []
                for r in rows[:300]:
                    if any(c is not None for c in r):
                        clean_rows.append([str(c) if c is not None else "" for c in r])
                if clean_rows:
                    headers = clean_rows[0]
                    cols = len(headers)
                    table_lines = [f"### Sheet: {sname} ({len(rows)} baris)"]
                    table_lines.append("| " + " | ".join(headers) + " |")
                    table_lines.append("| " + " | ".join(["---"] * cols) + " |")
                    for r in clean_rows[1:]:
                        padded = r + [""] * (cols - len(r))
                        table_lines.append("| " + " | ".join(padded[:cols]) + " |")
                    sheets_out.append("\n".join(table_lines))
            if sheets_out:
                return "\n\n".join(sheets_out)
        except Exception as e:
            logger.warning(f"openpyxl error: {e}")

    try:
        import pandas as pd
        if ext == ".ods":
            df = pd.read_excel(io.BytesIO(raw_bytes), engine="odf")
        else:
            df = pd.read_excel(io.BytesIO(raw_bytes))
        return f"### Data Spreadsheet: {fname} ({df.shape[0]} baris x {df.shape[1]} kolom)\n\n" + df.head(150).to_markdown(index=False)
    except Exception as e:
        logger.warning(f"pandas excel error: {e}")

    return f"Gagal membaca spreadsheet {fname}."


def _extract_csv(raw_bytes: bytes, fname: str) -> Optional[str]:
    try:
        import pandas as pd
        text_preview = raw_bytes[:4096].decode("utf-8", errors="replace")
        dialect = csv.Sniffer().sniff(text_preview)
        sep = dialect.delimiter
        df = pd.read_csv(io.BytesIO(raw_bytes), sep=sep)
        return f"### Data Tabel ({df.shape[0]} baris x {df.shape[1]} kolom, delimiter '{sep}'):\n\n" + df.head(150).to_markdown(index=False)
    except Exception:
        try:
            decoded = raw_bytes.decode("utf-8", errors="replace")
            lines = decoded.splitlines()[:150]
            return "```csv\n" + "\n".join(lines) + "\n```"
        except Exception:
            return None


def _extract_archive(raw_bytes: bytes, fname: str, ext: str) -> str:
    out = [f"### 📦 Struktur Berkas Arsip: {fname}"]
    try:
        if ext == ".zip":
            with zipfile.ZipFile(io.BytesIO(raw_bytes)) as zf:
                info_list = zf.infolist()
                out.append(f"Total {len(info_list)} item di dalam zip:\n")
                for info in info_list[:60]:
                    size_kb = f"{info.file_size / 1024:.1f} KB" if not info.is_dir() else "DIR"
                    out.append(f"- `/{info.filename}` ({size_kb})")
                
                for info in info_list[:10]:
                    if not info.is_dir() and info.file_size < 30000 and any(info.filename.endswith(e) for e in (".txt", ".md", ".json", ".py", ".csv")):
                        try:
                            content = zf.read(info.filename).decode("utf-8", errors="replace")
                            out.append(f"\n#### 📄 Pratinjau Isi: `{info.filename}`\n```\n{content[:2000]}\n```")
                            break
                        except Exception:
                            pass

        elif ext in (".tar", ".gz", ".tgz", ".bz2"):
            with tarfile.open(fileobj=io.BytesIO(raw_bytes)) as tf:
                members = tf.getmembers()
                out.append(f"Total {len(members)} item di dalam tar:\n")
                for m in members[:60]:
                    size_kb = f"{m.size / 1024:.1f} KB" if not m.isdir() else "DIR"
                    out.append(f"- `/{m.name}` ({size_kb})")

        return "\n".join(out)
    except Exception as e:
        return f"Gagal membaca arsip {fname}: {e}"


def _extract_html_xml(raw_bytes: bytes, fname: str, ext: str) -> str:
    try:
        from bs4 import BeautifulSoup
        decoded = raw_bytes.decode("utf-8", errors="replace")
        soup = BeautifulSoup(decoded, "html.parser")
        for tag in soup(["script", "style", "nav", "footer"]):
            tag.decompose()
        text = soup.get_text(separator="\n").strip()
        lines = [l.strip() for l in text.splitlines() if l.strip()]
        clean_text = "\n".join(lines[:500])
        return f"### Konten Web/Struktur {fname}:\n\n{clean_text}"
    except Exception:
        return raw_bytes.decode("utf-8", errors="replace")[:4000]


def _extract_data_formats(raw_bytes: bytes, fname: str, ext: str) -> str:
    try:
        decoded = raw_bytes.decode("utf-8", errors="replace")
        if ext in (".json", ".jsonl", ".ndjson"):
            if ext == ".json":
                parsed = json.loads(decoded)
                return "```json\n" + json.dumps(parsed, indent=2, ensure_ascii=False)[:4000] + "\n```"
            else:
                lines = decoded.splitlines()
                return f"Total {len(lines)} baris JSONL:\n```json\n" + "\n".join(lines[:50]) + "\n```"
        elif ext in (".yaml", ".yml"):
            import yaml
            parsed = yaml.safe_load(decoded)
            return "```yaml\n" + yaml.dump(parsed, sort_keys=False)[:4000] + "\n```"
        return "```\n" + decoded[:4000] + "\n```"
    except Exception:
        return raw_bytes.decode("utf-8", errors="replace")[:4000]
